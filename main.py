import streamlit as st
import io
import google.generativeai as genai
from PIL import Image
import requests
import datetime
import os
from pymongo import MongoClient
from bson import ObjectId
import json
import hashlib
from google.genai import types
import uuid
from typing import List, Dict
import openai
import pandas as pd
import csv
from perplexity import Perplexity
import openpyxl
from openpyxl.styles import Font, Alignment, Border, Side
from openpyxl.utils import get_column_letter

# Configure a API key do Perplexity
perp_api_key = os.getenv("PERP_API_KEY")
if perp_api_key:
    perplexity_client = Perplexity(api_key=perp_api_key)
else:
    st.warning("PERP_API_KEY não encontrada. Busca web estará desativada.")
    perplexity_client = None

# Configurações das credenciais
OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
ASTRA_DB_API_ENDPOINT = os.getenv('ASTRA_DB_API_ENDPOINT')
ASTRA_DB_APPLICATION_TOKEN = os.getenv('ASTRA_DB_APPLICATION_TOKEN')
ASTRA_DB_NAMESPACE = os.getenv('ASTRA_DB_NAMESPACE')
ASTRA_DB_COLLECTION = os.getenv('ASTRA_DB_COLLECTION')

senha_admin = os.getenv('SENHA_ADMIN')

mongo_uri = os.getenv('MONGO_URI')


senha_syn = os.getenv('SENHA_SYN')
senha_sme = os.getenv('SENHA_SME')
senha_ent = os.getenv('SENHA_ENT')

class AstraDBClient:
    def __init__(self):
        self.base_url = f"{ASTRA_DB_API_ENDPOINT}/api/json/v1/{ASTRA_DB_NAMESPACE}"
        self.headers = {
            "Content-Type": "application/json",
            "x-cassandra-token": ASTRA_DB_APPLICATION_TOKEN,
            "Accept": "application/json"
        }
    
    def vector_search(self, collection: str, vector: List[float], limit: int = 6) -> List[Dict]:
        """Realiza busca por similaridade vetorial"""
        url = f"{self.base_url}/{collection}"
        payload = {
            "find": {
                "sort": {"$vector": vector},
                "options": {"limit": limit}
            }
        }
        try:
            response = requests.post(url, json=payload, headers=self.headers, timeout=30)
            response.raise_for_status()
            data = response.json()
            return data.get("data", {}).get("documents", [])
        except Exception as e:
            st.error(f"Erro na busca vetorial: {str(e)}")
            return []

# Inicializa o cliente AstraDB
astra_client = AstraDBClient()

def get_embedding(text: str) -> List[float]:
    """Obtém embedding do texto usando OpenAI"""
    try:
        client = openai.OpenAI(api_key=OPENAI_API_KEY)
        response = client.embeddings.create(
            input=text,
            model="text-embedding-3-small"
        )
        return response.data[0].embedding
    except Exception as e:
        st.warning(f"Embedding OpenAI não disponível: {str(e)}")
        # Fallback para embedding simples
        import hashlib
        import numpy as np
        text_hash = hashlib.md5(text.encode()).hexdigest()
        vector = [float(int(text_hash[i:i+2], 16) / 255.0) for i in range(0, 32, 2)]
        # Preenche com valores aleatórios para ter 1536 dimensões
        while len(vector) < 1536:
            vector.append(0.0)
        return vector[:1536]

def reescrever_com_rag_blog(content: str) -> str:
    """REESCREVE conteúdo de blog usando RAG - SAÍDA DIRETA DO CONTEÚDO REESCRITO"""
    try:
        # Gera embedding para busca
        embedding = get_embedding(content[:800])
        
        # Busca documentos relevantes
        relevant_docs = astra_client.vector_search(ASTRA_DB_COLLECTION, embedding, limit=10)
        
        # Constrói contexto dos documentos
        rag_context = ""
        if relevant_docs:
            rag_context = "INFORMAÇÕES TÉCNICAS RELEVANTES DA BASE:\n"
            for i, doc in enumerate(relevant_docs, 1):
                doc_content = str(doc)
                # Limpa e formata o documento
                doc_clean = doc_content.replace('{', '').replace('}', '').replace("'", "").replace('"', '')
                rag_context += f"--- Fonte {i} ---\n{doc_clean[:500]}...\n\n"
        else:
            rag_context = "Base de conhecimento não retornou resultados específicos."

        # Prompt de entendimento RAG
        rewrite_prompt = f"""

        Entenda o que no texto original de fato é enriquecido e corrigido pelo referencial teórico. Considere que você não pode tangenciar o assunto do texto original.
    
        ###BEGIN TEXTO ORIGINAL###
        {content}
        ###END TEXTO ORIGINAL###

        ###BEGIN REFERENCIAL TEÓRICO###
        {rag_context}
        ###END REFERENCIAL TEÓRICO###
        
        
        """

        # Gera conteúdo REEESCRITO
        pre_response = modelo_texto.generate_content(rewrite_prompt)

        # Saída final
        final_prompt = f"""
    
        ###BEGIN TEXTO ORIGINAL###
        {content}
        ###END TEXTO ORIGINAL###

        ###BEGIN REFERENCIAL TEÓRICO###
        {pre_response}
        ###END REFERENCIAL TEÓRICO###
        
        Aplique isso ao texto original:

        1. SUBSTITUA termos vagos por terminologia técnica precisa da área agrícola que são relevantes ao texto original.
        2. CORRIGIR automaticamente qualquer imprecisão técnica ou científica no texto original
        3. ENRIQUECER com dados concretos, números e informações específicas da base
        4. MANTER tom {tom_voz} mas com precisão técnica absoluta
        5. MANTENHA a estrutura do texto original. Não reescreva por inteiro. Apenas corrija
        7. O agente revisor precisaria entregar o texto exatamente como no original, mas apontando os ajustes técnicos necessários/feitos, sem reescrever tudo automaticamente OU reescrevendo e sinalizando o que foi alterado no texto, mostrando como estava > como ficou > fonte/referência utilizada.
        8. NÃO acrescente informações que tangem o tema do texto original
        9. Mantenha o tamanho do texto original (com um delta de no máximo 5%)
        
        ESTRUTURA OBRIGATÓRIA:
        - Mantenha a estrutura original. O seu papel é REVISAR TECNICAMENTE O CONTEÚDO DE ENTRADA ENRIQUECENDO-O E, QUANDO NECESSÁRIO, CORRIJINDO-O COM O REFERENCIAL TEÓRICO.


        RETORNE O CONTEÚDO REEESCRITO FINAL, apontando as mudanças em uma subseção ao final.
        """
        
        response = modelo_texto.generate_content(final_prompt)
 
        return response.text
        
    except Exception as e:
        st.error(f"Erro no RAG rewrite para blog: {str(e)}")
        return content

def reescrever_com_rag_revisao_SEO(content: str) -> str:
    """REESCREVE conteúdo técnico para revisão - SAÍDA DIRETA DO CONTEÚDO REESCRITO"""
    try:
        # Gera embedding para busca
        embedding = get_embedding(content[:800])
        
        # Busca documentos relevantes
        relevant_docs = astra_client.vector_search(ASTRA_DB_COLLECTION, embedding, limit=10)
        
        # Constrói contexto dos documentos
        rag_context = ""
        if relevant_docs:
            rag_context = "DOCUMENTAÇÃO TÉCNICA ESPECIALIZADA:\n"
            for i, doc in enumerate(relevant_docs, 1):
                doc_content = str(doc)
                doc_clean = doc_content.replace('{', '').replace('}', '').replace("'", "").replace('"', '')
                rag_context += f"--- Documento Técnico {i} ---\n{doc_clean[:400]}...\n\n"
        else:
            rag_context = "Consulta técnica não retornou documentos específicos."

        # Prompt de REWRITE TÉCNICO AVANÇADO
        rewrite_prompt = f"""
        CONTEÚDO TÉCNICO ORIGINAL PARA REESCRITA COMPLETA:
        {content}

        
        
        BASE DE CONHECIMENTO TÉCNICO:
        {rag_context}

        Aplique isso ao texto original:

        

        1. SUBSTITUA termos vagos por terminologia técnica precisa da área agrícola que são relevantes ao texto original.
        2. CORRIGIR automaticamente qualquer imprecisão técnica ou científica no texto original
        3. ENRIQUECER com dados concretos, números e informações específicas da base
        4. MANTER tom {tom_voz} mas com precisão técnica absoluta
        5. MANTENHA a estrutura do texto original. Não reescreva por inteiro. Apenas corrija
        7. O agente revisor precisaria entregar o texto exatamente como no original, mas apontando os ajustes técnicos necessários/feitos, sem reescrever tudo automaticamente OU reescrevendo e sinalizando o que foi alterado no texto, mostrando como estava > como ficou > fonte/referência utilizada.
        8. NÃO acrescente informações que tangem o tema do texto original
        9. Mantenha o tamanho do texto original (com um delta de no máximo 5%)
        
        ESTRUTURA OBRIGATÓRIA:
        - Mantenha a estrutura original. O seu papel é REVISAR TECNICAMENTE O CONTEÚDO DE ENTRADA ENRIQUECENDO-O COM O REFERENCIAL TEÓRICO.


        RETORNE O CONTEÚDO REEESCRITO FINAL, apontando as mudanças em uma subseção ao final.
        """

        # Gera conteúdo técnico REEESCRITO
        response = modelo_texto.generate_content(rewrite_prompt)
        return response.text
        
    except Exception as e:
        st.error(f"Erro no RAG rewrite técnico: {str(e)}")
        return content

def reescrever_com_rag_revisao_NORM(content: str) -> str:
    """REESCREVE conteúdo técnico para revisão - SAÍDA DIRETA DO CONTEÚDO REESCRITO"""
    try:
        # Gera embedding para busca
        embedding = get_embedding(content[:800])
        
        # Busca documentos relevantes
        relevant_docs = astra_client.vector_search(ASTRA_DB_COLLECTION, embedding, limit=10)
        
        # Constrói contexto dos documentos
        rag_context = ""
        if relevant_docs:
            rag_context = "DOCUMENTAÇÃO TÉCNICA ESPECIALIZADA:\n"
            for i, doc in enumerate(relevant_docs, 1):
                doc_content = str(doc)
                doc_clean = doc_content.replace('{', '').replace('}', '').replace("'", "").replace('"', '')
                rag_context += f"--- Documento Técnico {i} ---\n{doc_clean[:400]}...\n\n"
        else:
            rag_context = "Consulta técnica não retornou documentos específicos."

        # Prompt de REWRITE TÉCNICO AVANÇADO
        rewrite_prompt = f"""
        CONTEÚDO TÉCNICO ORIGINAL PARA REESCRITA COMPLETE:
        {content}

        
        
        BASE DE CONHECIMENTO TÉCNICO:
        {rag_context}

        Aplique isso ao texto original:

        1. SUBSTITUA termos vagos por terminologia técnica precisa da área agrícola que são relevantes ao texto original.
        2. CORRIGIR automaticamente qualquer imprecisão técnica ou científica no texto original
        3. ENRIQUECER com dados concretos, números e informações específicas da base
        4. MANTER tom {tom_voz} mas com precisão técnica absoluta
        5. MANTENHA a estrutura do texto original. Não reescreva por inteiro. Apenas corrija
        7. O agente revisor precisaria entregar o texto exatamente como no original, mas apontando os ajustes técnicos necessários/feitos, sem reescrever tudo automaticamente OU reescrevendo e sinalizando o que foi alterado no texto, mostrando como estava > como ficou > fonte/referência utilizada.
        8. NÃO acrescente informações que tangem o tema do texto original
        9. Mantenha o tamanho do texto original (com um delta de no máximo 5%)
        10. NÃO USE BULLETS EM TODO O CONTEÚDO. MANTENHA OS ORIGINAIS E O RESTO DEVE VIR EM FORMATO DE PARÁGRAFO
        
        ESTRUTURA OBRIGATÓRIA:
        - Mantenha a estrutura original. O seu papel é REVISAR TECNICAMENTE O CONTEÚDO DE ENTRADA ENRIQUECENDO-O COM O REFERENCIAL TEÓRICO.


        RETORNE O CONTEÚDO REEESCRITO FINAL, apontando as mudanças em uma subseção ao final.
        """

        # Gera conteúdo técnico REEESCRITO
        response = modelo_texto.generate_content(rewrite_prompt)
        return response.text
        
    except Exception as e:
        st.error(f"Erro no RAG rewrite técnico: {str(e)}")
        return content

# Configuração inicial
st.set_page_config(
    layout="wide",
    page_title="Conteúdo")

# --- Sistema de Autenticação ---
def make_hashes(password):
    return hashlib.sha256(str.encode(password)).hexdigest()

def check_hashes(password, hashed_text):
    return make_hashes(password) == hashed_text

# Dados de usuário (em produção, isso deve vir de um banco de dados seguro)
users = {
    "admin": make_hashes(senha_admin),  # admin/senha1234
    "SYN": make_hashes(senha_syn),  # user1/password1
    "SME": make_hashes(senha_sme),   # user2/password2
    "Enterprise": make_hashes(senha_ent)   # user2/password2
}

def get_current_user():
    """Retorna o usuário atual da sessão"""
    return st.session_state.get('user', 'unknown')

def login():
    """Formulário de login"""
    
    with st.form("login_form"):
        username = st.text_input("Usuário")
        password = st.text_input("Senha", type="password")
        submit_button = st.form_submit_button("Login")
        
        if submit_button:
            if username in users and check_hashes(password, users[username]):
                st.session_state.logged_in = True
                st.session_state.user = username
                st.success("Login realizado com sucesso!")
                st.rerun()
            else:
                st.error("Usuário ou senha incorretos")

# Verificar se o usuário está logado
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False

if not st.session_state.logged_in:
    login()
    st.stop()

# --- CONEXÃO MONGODB (após login) ---
client = MongoClient(mongo_uri)
db = client['agentes_personalizados']
collection_agentes = db['agentes']
collection_conversas = db['conversas']

# Configuração da API do Gemini
gemini_api_key = os.getenv("GEM_API_KEY")
if not gemini_api_key:
    st.error("GEMINI_API_KEY não encontrada nas variáveis de ambiente")
    st.stop()

genai.configure(api_key=gemini_api_key)
modelo_vision = genai.GenerativeModel("gemini-2.5-flash", generation_config={"temperature": 0.0})
modelo_texto = genai.GenerativeModel("gemini-2.5-flash")
modelo_texto2 = genai.GenerativeModel("gemini-2.5-pro")

# --- Funções CRUD para Agentes ---
def criar_agente(nome, system_prompt, base_conhecimento, comments, planejamento, categoria, agente_mae_id=None, herdar_elementos=None):
    """Cria um novo agente no MongoDB"""
    agente = {
        "nome": nome,
        "system_prompt": system_prompt,
        "base_conhecimento": base_conhecimento,
        "comments": comments,
        "planejamento": planejamento,
        "categoria": categoria,
        "agente_mae_id": agente_mae_id,
        "herdar_elementos": herdar_elementos or [],
        "ativo": True,
        "data_criacao": datetime.datetime.now(),
        "criado_por": get_current_user()  # NOVO CAMPO
    }
    result = collection_agentes.insert_one(agente)
    return result.inserted_id

def listar_agentes():
    """Retorna todos os agentes ativos do usuário atual ou todos se admin"""
    current_user = get_current_user()
    if current_user == "admin":
        return list(collection_agentes.find({"ativo": True}).sort("data_criacao", -1))
    else:
        return list(collection_agentes.find({
            "ativo": True, 
            "criado_por": current_user
        }).sort("data_criacao", -1))

def listar_agentes_para_heranca(agente_atual_id=None):
    """Retorna todos os agentes ativos que podem ser usados como mãe"""
    current_user = get_current_user()
    query = {"ativo": True}
    
    # Filtro por usuário (admin vê todos, outros só os seus)
    if current_user != "admin":
        query["criado_por"] = current_user
    
    if agente_atual_id:
        # Excluir o próprio agente da lista de opções para evitar auto-herança
        if isinstance(agente_atual_id, str):
            agente_atual_id = ObjectId(agente_atual_id)
        query["_id"] = {"$ne": agente_atual_id}
    
    return list(collection_agentes.find(query).sort("data_criacao", -1))

def obter_agente(agente_id):
    """Obtém um agente específico pelo ID com verificação de permissão"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    agente = collection_agentes.find_one({"_id": agente_id})
    
    # Verificar permissão
    if agente and agente.get('ativo', True):
        current_user = get_current_user()
        if current_user == "admin" or agente.get('criado_por') == current_user:
            return agente
    
    return None

def atualizar_agente(agente_id, nome, system_prompt, base_conhecimento, comments, planejamento, categoria, agente_mae_id=None, herdar_elementos=None):
    """Atualiza um agente existente com verificação de permissão"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    # Verificar se o usuário tem permissão para editar este agente
    agente_existente = obter_agente(agente_id)
    if not agente_existente:
        raise PermissionError("Agente não encontrado ou sem permissão de edição")
    
    return collection_agentes.update_one(
        {"_id": agente_id},
        {
            "$set": {
                "nome": nome,
                "system_prompt": system_prompt,
                "base_conhecimento": base_conhecimento,
                "comments": comments,
                "planejamento": planejamento,
                "categoria": categoria,
                "agente_mae_id": agente_mae_id,
                "herdar_elementos": herdar_elementos or [],
            }
        }
    )

def desativar_agente(agente_id):
    """Desativa um agente (soft delete) com verificação de permissão"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    
    # Verificar se o usuário tem permissão para desativar este agente
    agente_existente = obter_agente(agente_id)
    if not agente_existente:
        raise PermissionError("Agente não encontrado ou sem permissão para desativar")
    
    return collection_agentes.update_one(
        {"_id": agente_id},
        {"$set": {"ativo": False}}
    )

def obter_agente_com_heranca(agente_id):
    """Obtém um agente com os elementos herdados aplicados"""
    agente = obter_agente(agente_id)
    if not agente or not agente.get('agente_mae_id'):
        return agente
    
    agente_mae = obter_agente(agente['agente_mae_id'])
    if not agente_mae:
        return agente
    
    elementos_herdar = agente.get('herdar_elementos', [])
    agente_completo = agente.copy()
    
    for elemento in elementos_herdar:
        if elemento == 'system_prompt' and not agente_completo.get('system_prompt'):
            agente_completo['system_prompt'] = agente_mae.get('system_prompt', '')
        elif elemento == 'base_conhecimento' and not agente_completo.get('base_conhecimento'):
            agente_completo['base_conhecimento'] = agente_mae.get('base_conhecimento', '')
        elif elemento == 'comments' and not agente_completo.get('comments'):
            agente_completo['comments'] = agente_mae.get('comments', '')
        elif elemento == 'planejamento' and not agente_completo.get('planejamento'):
            agente_completo['planejamento'] = agente_mae.get('planejamento', '')
    
    return agente_completo

def salvar_conversa(agente_id, mensagens, segmentos_utilizados=None):
    """Salva uma conversa no histórico"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    conversa = {
        "agente_id": agente_id,
        "mensagens": mensagens,
        "segmentos_utilizados": segmentos_utilizados,
        "data_criacao": datetime.datetime.now()
    }
    return collection_conversas.insert_one(conversa)

def obter_conversas(agente_id, limite=10):
    """Obtém o histórico de conversas de um agente"""
    if isinstance(agente_id, str):
        agente_id = ObjectId(agente_id)
    return list(collection_conversas.find(
        {"agente_id": agente_id}
    ).sort("data_criacao", -1).limit(limite))

# --- Função para construir contexto com segmentos selecionados ---
def construir_contexto(agente, segmentos_selecionados, historico_mensagens=None):
    """Constrói o contexto com base nos segmentos selecionados"""
    contexto = ""
    
    if "system_prompt" in segmentos_selecionados and agente.get('system_prompt'):
        contexto += f"### INSTRUÇÕES DO SISTEMA ###\n{agente['system_prompt']}\n\n"
    
    if "base_conhecimento" in segmentos_selecionados and agente.get('base_conhecimento'):
        contexto += f"### BASE DE CONHECIMENTO ###\n{agente['base_conhecimento']}\n\n"
    
    if "comments" in segmentos_selecionados and agente.get('comments'):
        contexto += f"### COMENTÁRIOS DO CLIENTE ###\n{agente['comments']}\n\n"
    
    if "planejamento" in segmentos_selecionados and agente.get('planejamento'):
        contexto += f"### PLANEJAMENTO ###\n{agente['planejamento']}\n\n"
    
    # Adicionar histórico se fornecido
    if historico_mensagens:
        contexto += "### HISTÓRICO DA CONVERSA ###\n"
        for msg in historico_mensagens:
            contexto += f"{msg['role']}: {msg['content']}\n"
        contexto += "\n"
    
    contexto += "### RESPOSTA ATUAL ###\nassistant:"
    
    return contexto

# --- Funções para Transcrição de Áudio/Video ---
def transcrever_audio_video(arquivo, tipo_arquivo):
    """Transcreve áudio ou vídeo usando a API do Gemini"""
    try:
        client = genai.Client(api_key=gemini_api_key)
        
        if tipo_arquivo == "audio":
            mime_type = f"audio/{arquivo.name.split('.')[-1]}"
        else:  # video
            mime_type = f"video/{arquivo.name.split('.')[-1]}"
        
        # Lê os bytes do arquivo
        arquivo_bytes = arquivo.read()
        
        # Para arquivos maiores, usa upload
        if len(arquivo_bytes) > 20 * 1024 * 1024:  # 20MB
            uploaded_file = client.files.upload(file=arquivo_bytes, mime_type=mime_type)
            response = client.models.generate_content(
                model="gemini-2.0-flash", 
                contents=["Transcreva este arquivo em detalhes:", uploaded_file]
            )
        else:
            # Para arquivos menores, usa inline
            response = client.models.generate_content(
                model="gemini-2.0-flash",
                contents=[
                    "Transcreva este arquivo em detalhes:",
                    types.Part.from_bytes(data=arquivo_bytes, mime_type=mime_type)
                ]
            )
        
        return response.text
    except Exception as e:
        return f"Erro na transcrição: {str(e)}"

# --- Configuração de Autenticação de Administrador ---
def check_admin_password():
    """Retorna True se o usuário fornecer a senha de admin correta."""
    
    def admin_password_entered():
        """Verifica se a senha de admin está correta."""
        if st.session_state["admin_password"] == senha_admin:
            st.session_state["admin_password_correct"] = True
            st.session_state["admin_user"] = "admin"
            del st.session_state["admin_password"]
        else:
            st.session_state["admin_password_correct"] = False

    if "admin_password_correct" not in st.session_state:
        # Mostra o input para senha de admin
        st.text_input(
            "Senha de Administrador", 
            type="password", 
            on_change=admin_password_entered, 
            key="admin_password"
        )
        return False
    elif not st.session_state["admin_password_correct"]:
        # Senha incorreta, mostra input + erro
        st.text_input(
            "Senha de Administrador", 
            type="password", 
            on_change=admin_password_entered, 
            key="admin_password"
        )
        st.error("😕 Senha de administrador incorreta")
        return False
    else:
        # Senha correta
        return True

# ========== SELEÇÃO EXTERNA DE AGENTE ==========
st.image('macLogo.png', width=300)
st.title("Conteúdo")

# Botão de logout na sidebar
if st.button("🚪 Sair", key="logout_btn"):
    for key in ["logged_in", "user", "admin_password_correct", "admin_user"]:
        if key in st.session_state:
            del st.session_state[key]
    st.rerun()

# --- SELEÇÃO DE AGENTE EXTERNA ---
st.header("🤖 Selecione a base de conhecimento")

# Inicializar estado da sessão para agente selecionado
if "agente_selecionado" not in st.session_state:
    st.session_state.agente_selecionado = None
if "segmentos_selecionados" not in st.session_state:
    st.session_state.segmentos_selecionados = ["system_prompt", "base_conhecimento", "comments", "planejamento"]

# Carregar agentes (agora filtrados por usuário)
agentes = listar_agentes()

# Container para seleção de agente
with st.container():
    col1, col2, col3 = st.columns([3, 1, 1])
    
    with col1:
        if agentes:
            # Agrupar agentes por categoria
            agentes_por_categoria = {}
            for agente in agentes:
                categoria = agente.get('categoria', 'Social')
                if categoria not in agentes_por_categoria:
                    agentes_por_categoria[categoria] = []
                agentes_por_categoria[categoria].append(agente)
            
            # Criar opções de seleção com agrupamento
            agente_options = {}
            for categoria, agentes_cat in agentes_por_categoria.items():
                for agente in agentes_cat:
                    agente_completo = obter_agente_com_heranca(agente['_id'])
                    if agente_completo:  # Só adiciona se tiver permissão
                        display_name = f"{agente['nome']} ({categoria})"
                        if agente.get('agente_mae_id'):
                            display_name += " 🔗"
                        # Adicionar indicador de proprietário se não for admin
                        if get_current_user() != "admin" and agente.get('criado_por'):
                            display_name += f" 👤"
                        agente_options[display_name] = agente_completo
            
            if agente_options:
                # Seletor de agente
                agente_selecionado_display = st.selectbox(
                    "Selecione um agente para trabalhar:", 
                    list(agente_options.keys()),
                    key="seletor_agente_global"
                )
                
                # Botão para aplicar agente
                if st.button("🔄 Aplicar Agente", key="aplicar_agente"):
                    st.session_state.agente_selecionado = agente_options[agente_selecionado_display]
                    st.success(f"Agente '{agente_selecionado_display}' selecionado!")
                    st.rerun()
            else:
                st.info("Nenhum agente disponível com as permissões atuais.")
        
        else:
            st.info("Nenhum agente disponível. Crie um agente primeiro na aba de Gerenciamento.")
    
    with col2:
        # Botão para limpar agente selecionado
        if st.session_state.agente_selecionado:
            if st.button("🗑️ Limpar Agente", key="limpar_agente"):
                st.session_state.agente_selecionado = None
                st.session_state.messages = []
                st.success("Agente removido!")
                st.rerun()
    
    with col3:
        # Botão para recarregar lista
        if st.button("🔄 Recarregar", key="recarregar_agentes"):
            st.rerun()

# Mostrar agente atual selecionado
if st.session_state.agente_selecionado:
    agente_atual = st.session_state.agente_selecionado
    
    # Container para informações do agente
    with st.container():
        st.success(f"**✅ Agente Ativo:** {agente_atual['nome']} ({agente_atual.get('categoria', 'Social')})")
        
        # Mostrar informações de herança se aplicável
        if 'agente_mae_id' in agente_atual and agente_atual['agente_mae_id']:
            agente_original = obter_agente(agente_atual['_id'])
            if agente_original and agente_original.get('herdar_elementos'):
                st.info(f"🔗 Este agente herda {len(agente_original['herdar_elementos'])} elementos do agente mãe")
        
        # Mostrar segmentos ativos
        st.info(f"📋 Segmentos ativos: {', '.join(st.session_state.segmentos_selecionados)}")
        
        # Botão para alterar segmentos
        if st.button("⚙️ Alterar Segmentos", key="alterar_segmentos"):
            # Toggle para mostrar/ocultar configuração de segmentos
            if "mostrar_segmentos" not in st.session_state:
                st.session_state.mostrar_segmentos = True
            else:
                st.session_state.mostrar_segmentos = not st.session_state.mostrar_segmentos
        
        # Mostrar configuração de segmentos se solicitado
        if st.session_state.get('mostrar_segmentos', False):
            with st.expander("🔧 Configurar Segmentos do Agente", expanded=True):
                st.write("Selecione quais elementos do agente serão utilizados:")
                
                col_seg1, col_seg2, col_seg3, col_seg4 = st.columns(4)
                
                with col_seg1:
                    system_prompt_ativado = st.checkbox("System Prompt", 
                                                      value="system_prompt" in st.session_state.segmentos_selecionados,
                                                      key="seg_system")
                with col_seg2:
                    base_conhecimento_ativado = st.checkbox("Brand Guidelines", 
                                                          value="base_conhecimento" in st.session_state.segmentos_selecionados,
                                                          key="seg_base")
                with col_seg3:
                    comments_ativado = st.checkbox("Comentários", 
                                                 value="comments" in st.session_state.segmentos_selecionados,
                                                 key="seg_comments")
                with col_seg4:
                    planejamento_ativado = st.checkbox("Planejamento", 
                                                     value="planejamento" in st.session_state.segmentos_selecionados,
                                                     key="seg_planejamento")
                
                if st.button("✅ Aplicar Segmentos", key="aplicar_segmentos"):
                    novos_segmentos = []
                    if system_prompt_ativado:
                        novos_segmentos.append("system_prompt")
                    if base_conhecimento_ativado:
                        novos_segmentos.append("base_conhecimento")
                    if comments_ativado:
                        novos_segmentos.append("comments")
                    if planejamento_ativado:
                        novos_segmentos.append("planejamento")
                    
                    st.session_state.segmentos_selecionados = novos_segmentos
                    st.success(f"Segmentos atualizados: {', '.join(novos_segmentos)}")
                    st.session_state.mostrar_segmentos = False
                    st.rerun()

else:
    st.warning("⚠️ Nenhum agente selecionado. Selecione um agente acima para começar.")

st.markdown("---")

# Menu de abas - AGORA COM A NOVA ABA DE CALENDÁRIO
tab_chat, tab_gerenciamento, tab_conteudo, tab_blog, tab_revisao_ortografica, tab_revisao_tecnica, tab_otimizacao, tab_calendario, tab_briefings, tab_revisao_tecnica2 = st.tabs([
    "💬 Chat", 
    "⚙️ Gerenciar Agentes",
    "✨ Geração de Conteúdo", 
    "🌱 Geração de Conteúdo Blog",
    "📝 Revisão Ortográfica",
    "🔧 Revisão Técnica",
    "🚀 Otimização de Conteúdo",
    "📅 Criadora de Calendário",
    "📋 Gerador de Briefings",
    "Revisão Técnica Sem RAG" # NOVA ABA
])

# ========== ABA: CHAT ==========
with tab_chat:
    st.header("💬 Chat com Agente")
    
    # Inicializar estado da sessão
    if "messages" not in st.session_state:
        st.session_state.messages = []
    
    # Verificar se há agente selecionado
    if not st.session_state.agente_selecionado:
        st.info("Selecione um agente na parte superior do app para iniciar o chat.")
    else:
        agente = st.session_state.agente_selecionado
        st.subheader(f"Conversando com: {agente['nome']}")
        
        # Exibir histórico de mensagens
        for message in st.session_state.messages:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])
        
        # Input do usuário
        if prompt := st.chat_input("Digite sua mensagem..."):
            # Adicionar mensagem do usuário ao histórico
            st.session_state.messages.append({"role": "user", "content": prompt})
            with st.chat_message("user"):
                st.markdown(prompt)
            
            # Construir contexto com segmentos selecionados
            contexto = construir_contexto(
                agente, 
                st.session_state.segmentos_selecionados, 
                st.session_state.messages
            )
            
            # Gerar resposta
            with st.chat_message("assistant"):
                with st.spinner('Pensando...'):
                    try:
                        resposta = modelo_texto.generate_content(contexto)
                        st.markdown(resposta.text)
                        
                        # Adicionar ao histórico
                        st.session_state.messages.append({"role": "assistant", "content": resposta.text})
                        
                        # Salvar conversa com segmentos utilizados
                        salvar_conversa(
                            agente['_id'], 
                            st.session_state.messages,
                            st.session_state.segmentos_selecionados
                        )
                        
                    except Exception as e:
                        st.error(f"Erro ao gerar resposta: {str(e)}")

# ========== ABA: GERENCIAMENTO DE AGENTES ==========
with tab_gerenciamento:
    st.header("⚙️ Gerenciamento de Agentes")
    
    # Verificar autenticação apenas para gerenciamento
    current_user = get_current_user()
    
    if current_user not in ["admin", "SYN", "SME", "Enterprise"]:
        st.warning("Acesso restrito a usuários autorizados")
    else:
        # Para admin, verificar senha adicional
        if current_user == "admin":
            if not check_admin_password():
                st.warning("Digite a senha de administrador")
            else:
                st.write(f'Bem-vindo administrador!')
        else:
            st.write(f'Bem-vindo {current_user}!')
            
        # Subabas para gerenciamento
        sub_tab1, sub_tab2, sub_tab3 = st.tabs(["Criar Agente", "Editar Agente", "Gerenciar Agentes"])
        
        with sub_tab1:
            st.subheader("Criar Novo Agente")
            
            with st.form("form_criar_agente"):
                nome_agente = st.text_input("Nome do Agente:")
                
                # Seleção de categoria
                categoria = st.selectbox(
                    "Categoria:",
                    ["Social", "SEO", "Conteúdo"],
                    help="Organize o agente por área de atuação"
                )
                
                # Opção para criar como agente filho
                criar_como_filho = st.checkbox("Criar como agente filho (herdar elementos)")
                
                agente_mae_id = None
                herdar_elementos = []
                
                if criar_como_filho:
                    # Listar TODOS os agentes disponíveis para herança
                    agentes_mae = listar_agentes_para_heranca()
                    if agentes_mae:
                        agente_mae_options = {f"{agente['nome']} ({agente.get('categoria', 'Social')})": agente['_id'] for agente in agentes_mae}
                        agente_mae_selecionado = st.selectbox(
                            "Agente Mãe:",
                            list(agente_mae_options.keys()),
                            help="Selecione o agente do qual este agente irá herdar elementos"
                        )
                        agente_mae_id = agente_mae_options[agente_mae_selecionado]
                        
                        st.subheader("Elementos para Herdar")
                        herdar_elementos = st.multiselect(
                            "Selecione os elementos a herdar do agente mãe:",
                            ["system_prompt", "base_conhecimento", "comments", "planejamento"],
                            help="Estes elementos serão herdados do agente mãe se não preenchidos abaixo"
                        )
                    else:
                        st.info("Nenhum agente disponível para herança. Crie primeiro um agente mãe.")
                
                system_prompt = st.text_area("Prompt de Sistema:", height=150, 
                                            placeholder="Ex: Você é um assistente especializado em...",
                                            help="Deixe vazio se for herdar do agente mãe")
                base_conhecimento = st.text_area("Brand Guidelines:", height=200,
                                               placeholder="Cole aqui informações, diretrizes, dados...",
                                               help="Deixe vazio se for herdar do agente mãe")
                comments = st.text_area("Comentários do cliente:", height=200,
                                               placeholder="Cole aqui os comentários de ajuste do cliente (Se houver)",
                                               help="Deixe vazio se for herdar do agente mãe")
                planejamento = st.text_area("Planejamento:", height=200,
                                           placeholder="Estratégias, planejamentos, cronogramas...",
                                           help="Deixe vazio se for herdar do agente mãe")
                
                submitted = st.form_submit_button("Criar Agente")
                if submitted:
                    if nome_agente:
                        agente_id = criar_agente(
                            nome_agente, 
                            system_prompt, 
                            base_conhecimento, 
                            comments, 
                            planejamento,
                            categoria,
                            agente_mae_id if criar_como_filho else None,
                            herdar_elementos if criar_como_filho else []
                        )
                        st.success(f"Agente '{nome_agente}' criado com sucesso na categoria {categoria}!")
                    else:
                        st.error("Nome é obrigatório!")
        
        with sub_tab2:
            st.subheader("Editar Agente Existente")
            
            agentes = listar_agentes()
            if agentes:
                agente_options = {agente['nome']: agente for agente in agentes}
                agente_selecionado_nome = st.selectbox("Selecione o agente para editar:", 
                                                     list(agente_options.keys()))
                
                if agente_selecionado_nome:
                    agente = agente_options[agente_selecionado_nome]
                    
                    with st.form("form_editar_agente"):
                        novo_nome = st.text_input("Nome do Agente:", value=agente['nome'])
                        
                        # Categoria
                        
                        
                        # Informações de herança
                        if agente.get('agente_mae_id'):
                            agente_mae = obter_agente(agente['agente_mae_id'])
                            if agente_mae:
                                st.info(f"🔗 Este agente é filho de: {agente_mae['nome']}")
                                st.write(f"Elementos herdados: {', '.join(agente.get('herdar_elementos', []))}")
                        
                        # Opção para tornar independente
                        if agente.get('agente_mae_id'):
                            tornar_independente = st.checkbox("Tornar agente independente (remover herança)")
                            if tornar_independente:
                                agente_mae_id = None
                                herdar_elementos = []
                            else:
                                agente_mae_id = agente.get('agente_mae_id')
                                herdar_elementos = agente.get('herdar_elementos', [])
                        else:
                            agente_mae_id = None
                            herdar_elementos = []
                            # Opção para adicionar herança
                            adicionar_heranca = st.checkbox("Adicionar herança de agente mãe")
                            if adicionar_heranca:
                                # Listar TODOS os agentes disponíveis para herança (excluindo o próprio)
                                agentes_mae = listar_agentes_para_heranca(agente['_id'])
                                if agentes_mae:
                                    agente_mae_options = {f"{agente_mae['nome']} ({agente_mae.get('categoria', 'Social')})": agente_mae['_id'] for agente_mae in agentes_mae}
                                    if agente_mae_options:
                                        agente_mae_selecionado = st.selectbox(
                                            "Agente Mãe:",
                                            list(agente_mae_options.keys()),
                                            help="Selecione o agente do qual este agente irá herdar elementos"
                                        )
                                        agente_mae_id = agente_mae_options[agente_mae_selecionado]
                                        herdar_elementos = st.multiselect(
                                            "Elementos para herdar:",
                                            ["system_prompt", "base_conhecimento", "comments", "planejamento"],
                                            default=herdar_elementos
                                        )
                                    else:
                                        st.info("Nenhum agente disponível para herança.")
                                else:
                                    st.info("Nenhum agente disponível para herança.")
                        
                        novo_prompt = st.text_area("Prompt de Sistema:", value=agente['system_prompt'], height=150)
                        nova_base = st.text_area("Brand Guidelines:", value=agente.get('base_conhecimento', ''), height=200)
                        nova_comment = st.text_area("Comentários:", value=agente.get('comments', ''), height=200)
                        novo_planejamento = st.text_area("Planejamento:", value=agente.get('planejamento', ''), height=200)
                        
                        submitted = st.form_submit_button("Atualizar Agente")
                        if submitted:
                            if novo_nome:
                                atualizar_agente(
                                    agente['_id'], 
                                    novo_nome, 
                                    novo_prompt, 
                                    nova_base, 
                                    nova_comment, 
                                    novo_planejamento,
                                    agente_mae_id,
                                    herdar_elementos
                                )
                                st.success(f"Agente '{novo_nome}' atualizado com sucesso!")
                                st.rerun()
                            else:
                                st.error("Nome é obrigatório!")
            else:
                st.info("Nenhum agente criado ainda.")
        
        with sub_tab3:
            st.subheader("Gerenciar Agentes")
            
            # Mostrar informações do usuário atual
            if current_user == "admin":
                st.info("👑 Modo Administrador: Visualizando todos os agentes do sistema")
            else:
                st.info(f"👤 Visualizando apenas seus agentes ({current_user})")
            
            # Filtros por categoria
            categorias = ["Todos", "Social", "SEO", "Conteúdo"]
            categoria_filtro = st.selectbox("Filtrar por categoria:", categorias)
            
            agentes = listar_agentes()
            
            # Aplicar filtro
            if categoria_filtro != "Todos":
                agentes = [agente for agente in agentes if agente.get('categoria') == categoria_filtro]
            
            if agentes:
                for i, agente in enumerate(agentes):
                    with st.container():
                        # Mostrar proprietário se for admin
                        owner_info = ""
                        if current_user == "admin" and agente.get('criado_por'):
                            owner_info = f" | 👤 {agente['criado_por']}"
                        
                        st.write(f"**{agente['nome']} - {agente.get('categoria', 'Social')}{owner_info} - Criado em {agente['data_criacao'].strftime('%d/%m/%Y')}**")
                        
                        # Mostrar informações de herança
                        if agente.get('agente_mae_id'):
                            agente_mae = obter_agente(agente['agente_mae_id'])
                            if agente_mae:
                                st.write(f"**🔗 Herda de:** {agente_mae['nome']}")
                                st.write(f"**Elementos herdados:** {', '.join(agente.get('herdar_elementos', []))}")
                        
                        st.write(f"**Prompt de Sistema:** {agente['system_prompt'][:100]}..." if agente['system_prompt'] else "**Prompt de Sistema:** (herdado ou vazio)")
                        if agente.get('base_conhecimento'):
                            st.write(f"**Brand Guidelines:** {agente['base_conhecimento'][:200]}...")
                        if agente.get('comments'):
                            st.write(f"**Comentários do cliente:** {agente['comments'][:200]}...")
                        if agente.get('planejamento'):
                            st.write(f"**Planejamento:** {agente['planejamento'][:200]}...")
                        
                        col1, col2 = st.columns(2)
                        with col1:
                            if st.button("Selecionar para Chat", key=f"select_{i}"):
                                st.session_state.agente_selecionado = obter_agente_com_heranca(agente['_id'])
                                st.session_state.messages = []
                                st.success(f"Agente '{agente['nome']}' selecionado!")
                        with col2:
                            if st.button("Desativar", key=f"delete_{i}"):
                                desativar_agente(agente['_id'])
                                st.success(f"Agente '{agente['nome']}' desativado!")
                                st.rerun()
                        st.divider()
            else:
                st.info("Nenhum agente encontrado para esta categoria.")


# ========== ABA: GERAÇÃO DE CONTEÚDO ==========
with tab_conteudo:
    st.header("✨ Geração de Conteúdo com Múltiplos Insumos")
    
    # Conexão com MongoDB para briefings
    try:
        client2 = MongoClient(mongo_uri)
        db_briefings = client2['briefings_Broto_Tecnologia']
        collection_briefings = db_briefings['briefings']
        mongo_connected_conteudo = True
    except Exception as e:
        st.error(f"Erro na conexão com MongoDB: {str(e)}")
        mongo_connected_conteudo = False

    # Função para extrair texto de diferentes tipos de arquivo
    def extrair_texto_arquivo(arquivo):
        """Extrai texto de diferentes formatos de arquivo"""
        try:
            extensao = arquivo.name.split('.')[-1].lower()
            
            if extensao == 'pdf':
                return extrair_texto_pdf(arquivo)
            elif extensao == 'txt':
                return extrair_texto_txt(arquivo)
            elif extensao in ['pptx', 'ppt']:
                return extrair_texto_pptx(arquivo)
            elif extensao in ['docx', 'doc']:
                return extrair_texto_docx(arquivo)
            else:
                return f"Formato {extensao} não suportado para extração de texto."
                
        except Exception as e:
            return f"Erro ao extrair texto do arquivo {arquivo.name}: {str(e)}"

    def extrair_texto_pdf(arquivo):
        """Extrai texto de arquivos PDF"""
        try:
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(arquivo)
            texto = ""
            for pagina in pdf_reader.pages:
                texto += pagina.extract_text() + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do PDF: {str(e)}"

    def extrair_texto_txt(arquivo):
        """Extrai texto de arquivos TXT"""
        try:
            return arquivo.read().decode('utf-8')
        except:
            try:
                return arquivo.read().decode('latin-1')
            except Exception as e:
                return f"Erro na leitura do TXT: {str(e)}"

    def extrair_texto_pptx(arquivo):
        """Extrai texto de arquivos PowerPoint"""
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(arquivo.read()))
            texto = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto += shape.text + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do PowerPoint: {str(e)}"

    def extrair_texto_docx(arquivo):
        """Extrai texto de arquivos Word"""
        try:
            import docx
            import io
            doc = docx.Document(io.BytesIO(arquivo.read()))
            texto = ""
            for para in doc.paragraphs:
                texto += para.text + "\n"
            return texto
        except Exception as e:
            return f"Erro na leitura do Word: {str(e)}"

    # Layout principal
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.subheader("📝 Fontes de Conteúdo")
        
        # Opção 1: Upload de múltiplos arquivos
        st.write("📎 Upload de Arquivos (PDF, TXT, PPTX, DOCX):")
        arquivos_upload = st.file_uploader(
            "Selecione um ou mais arquivos:",
            type=['pdf', 'txt', 'pptx', 'ppt', 'docx', 'doc'],
            accept_multiple_files=True,
            help="Arquivos serão convertidos para texto e usados como base para geração de conteúdo"
        )
        
        # Processar arquivos uploadados
        textos_arquivos = ""
        if arquivos_upload:
            st.success(f"✅ {len(arquivos_upload)} arquivo(s) carregado(s)")
            
            with st.expander("📋 Visualizar Conteúdo dos Arquivos", expanded=False):
                for i, arquivo in enumerate(arquivos_upload):
                    st.write(f"**{arquivo.name}** ({arquivo.size} bytes)")
                    with st.spinner(f"Processando {arquivo.name}..."):
                        texto_extraido = extrair_texto_arquivo(arquivo)
                        textos_arquivos += f"\n\n--- CONTEÚDO DE {arquivo.name.upper()} ---\n{texto_extraido}"
                        
                        # Mostrar preview
                        if len(texto_extraido) > 500:
                            st.text_area(f"Preview - {arquivo.name}", 
                                       value=texto_extraido[:500] + "...", 
                                       height=100,
                                       key=f"preview_{i}")
                        else:
                            st.text_area(f"Preview - {arquivo.name}", 
                                       value=texto_extraido, 
                                       height=100,
                                       key=f"preview_{i}")
        
        # Opção 2: Selecionar briefing do banco de dados
        st.write("🗃️ Briefing do Banco de Dados:")
        if mongo_connected_conteudo:
            briefings_disponiveis = list(collection_briefings.find().sort("data_criacao", -1).limit(20))
            if briefings_disponiveis:
                briefing_options = {f"{briefing['nome_projeto']} ({briefing['tipo']}) - {briefing['data_criacao'].strftime('%d/%m/%Y')}": briefing for briefing in briefings_disponiveis}
                briefing_selecionado = st.selectbox("Escolha um briefing:", list(briefing_options.keys()))
                
                if briefing_selecionado:
                    briefing_data = briefing_options[briefing_selecionado]
                    st.info(f"Briefing selecionado: {briefing_data['nome_projeto']}")
            else:
                st.info("Nenhum briefing encontrado no banco de dados.")
        else:
            st.warning("Conexão com MongoDB não disponível")
        
        # Opção 3: Inserir briefing manualmente
        st.write("✍️ Briefing Manual:")
        briefing_manual = st.text_area("Ou cole o briefing completo aqui:", height=150,
                                      placeholder="""Exemplo:
Título: Campanha de Lançamento
Objetivo: Divulgar novo produto
Público-alvo: Empresários...
Pontos-chave: [lista os principais pontos]""")
        
        # Transcrição de áudio/vídeo
        st.write("🎤 Transcrição de Áudio/Video:")
        arquivos_midia = st.file_uploader(
            "Áudios/Vídeos para transcrição:",
            type=['mp3', 'wav', 'mp4', 'mov', 'avi'],
            accept_multiple_files=True,
            help="Arquivos de mídia serão transcritos automaticamente"
        )
        
        transcricoes_texto = ""
        if arquivos_midia:
            st.info(f"🎬 {len(arquivos_midia)} arquivo(s) de mídia carregado(s)")
            if st.button("🔄 Transcrever Todos os Arquivos de Mídia"):
                with st.spinner("Transcrevendo arquivos de mídia..."):
                    for arquivo in arquivos_midia:
                        tipo = "audio" if arquivo.type.startswith('audio') else "video"
                        transcricao = transcrever_audio_video(arquivo, tipo)
                        transcricoes_texto += f"\n\n--- TRANSCRIÇÃO DE {arquivo.name.upper()} ---\n{transcricao}"
                        st.success(f"✅ {arquivo.name} transcrito!")
    
    with col2:
        st.subheader("⚙️ Configurações")
        
        tipo_conteudo = st.selectbox("Tipo de Conteúdo:", 
                                   ["Post Social", "Artigo Blog", "Email Marketing", 
                                    "Landing Page", "Script Vídeo", "Relatório Técnico",
                                    "Press Release", "Newsletter", "Case Study"])
        
        tom_voz = st.selectbox("Tom de Voz:", 
                              ["Formal", "Informal", "Persuasivo", "Educativo", 
                               "Inspirador", "Técnico", "Jornalístico"], key = 'qq')
        
        palavras_chave = st.text_input("Palavras-chave (opcional):",
                                      placeholder="separadas por vírgula")
        
        numero_palavras = st.slider("Número de Palavras:", 100, 3000, 800)
        
        # Configurações avançadas
        with st.expander("🔧 Configurações Avançadas"):
            usar_contexto_agente = st.checkbox("Usar contexto do agente selecionado", 
                                             value=bool(st.session_state.agente_selecionado))
            
            nivel_detalhe = st.select_slider("Nível de Detalhe:", 
                                           ["Resumido", "Balanceado", "Detalhado", "Completo"])
            
            incluir_cta = st.checkbox("Incluir Call-to-Action", value=True)
            
            formato_saida = st.selectbox("Formato de Saída:", 
                                       ["Texto Simples", "Markdown", "HTML Básico"])

    # Área de instruções específicas
    st.subheader("🎯 Instruções Específicas")
    instrucoes_especificas = st.text_area(
        "Diretrizes adicionais para geração:",
        placeholder="""Exemplos:
- Focar nos benefícios para o usuário final
- Incluir estatísticas quando possível
- Manter linguagem acessível
- Evitar jargões técnicos excessivos
- Seguir estrutura: problema → solução → benefícios""",
        height=100
    )

    # Botão para gerar conteúdo
    if st.button("🚀 Gerar Conteúdo com Todos os Insumos", type="primary", use_container_width=True):
        # Verificar se há pelo menos uma fonte de conteúdo
        tem_conteudo = (arquivos_upload or 
                       briefing_manual or 
                       ('briefing_data' in locals() and briefing_data) or
                       arquivos_midia)
        
        if not tem_conteudo:
            st.error("❌ Por favor, forneça pelo menos uma fonte de conteúdo (arquivos, briefing ou mídia)")
        else:
            with st.spinner("Processando todos os insumos e gerando conteúdo..."):
                try:
                    # Construir o contexto combinado de todas as fontes
                    contexto_completo = "## FONTES DE CONTEÚDO COMBINADAS:\n\n"
                    
                    # Adicionar conteúdo dos arquivos uploadados
                    if textos_arquivos:
                        contexto_completo += "### CONTEÚDO DOS ARQUIVOS:\n" + textos_arquivos + "\n\n"
                    
                    # Adicionar briefing do banco ou manual
                    if briefing_manual:
                        contexto_completo += "### BRIEFING MANUAL:\n" + briefing_manual + "\n\n"
                    elif 'briefing_data' in locals() and briefing_data:
                        contexto_completo += "### BRIEFING DO BANCO:\n" + briefing_data['conteudo'] + "\n\n"
                    
                    # Adicionar transcrições
                    if transcricoes_texto:
                        contexto_completo += "### TRANSCRIÇÕES DE MÍDIA:\n" + transcricoes_texto + "\n\n"
                    
                    # Adicionar contexto do agente se selecionado
                    contexto_agente = ""
                    if usar_contexto_agente and st.session_state.agente_selecionado:
                        agente = st.session_state.agente_selecionado
                        contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                    
                    # Construir prompt final
                    prompt_final = f"""
                    {contexto_agente}
                    
                    ## INSTRUÇÕES PARA GERAÇÃO DE CONTEÚDO:
                    
                    **TIPO DE CONTEÚDO:** {tipo_conteudo}
                    **TOM DE VOZ:** {tom_voz}
                    **PALAVRAS-CHAVE:** {palavras_chave if palavras_chave else 'Não especificadas'}
                    **NÚMERO DE PALAVRAS:** {numero_palavras} (±10%)
                    **NÍVEL DE DETALHE:** {nivel_detalhe}
                    **INCLUIR CALL-TO-ACTION:** {incluir_cta}
                    
                    **INSTRUÇÕES ESPECÍFICAS:**
                    {instrucoes_especificas if instrucoes_especificas else 'Nenhuma instrução específica fornecida.'}
                    
                    ## FONTES E REFERÊNCIAS:
                    {contexto_completo}
                    
                    ## TAREFA:
                    Com base em TODAS as fontes fornecidas acima, gere um conteúdo do tipo {tipo_conteudo} que:
                    
                    1. **Síntese Eficiente:** Combine e sintetize informações de todas as fontes
                    2. **Coerência:** Mantenha consistência com as informações originais
                    3. **Valor Agregado:** Vá além da simples cópia, agregando insights
                    4. **Engajamento:** Crie conteúdo que engaje o público-alvo
                    5. **Clareza:** Comunique ideias complexas de forma acessível
                    
                    **FORMATO DE SAÍDA:** {formato_saida}
                    
                    Gere um conteúdo completo e profissional.
                    """
                    
                    resposta = modelo_texto.generate_content(prompt_final)
                    
                    # Processar saída baseada no formato selecionado
                    conteudo_gerado = resposta.text
                    
                    if formato_saida == "HTML Básico":
                        # Converter markdown para HTML básico
                        import re
                        conteudo_gerado = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'\*(.*?)\*', r'<em>\1</em>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'### (.*?)\n', r'<h3>\1</h3>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'## (.*?)\n', r'<h2>\1</h2>', conteudo_gerado)
                        conteudo_gerado = re.sub(r'# (.*?)\n', r'<h1>\1</h1>', conteudo_gerado)
                        conteudo_gerado = conteudo_gerado.replace('\n', '<br>')
                    
                    st.subheader("📄 Conteúdo Gerado")
                    
                    if formato_saida == "HTML Básico":
                        st.components.v1.html(conteudo_gerado, height=400, scrolling=True)
                    else:
                        st.markdown(conteudo_gerado)
                    
                    # Estatísticas
                    palavras_count = len(conteudo_gerado.split())
                    col_stat1, col_stat2, col_stat3 = st.columns(3)
                    with col_stat1:
                        st.metric("Palavras Geradas", palavras_count)
                    with col_stat2:
                        st.metric("Arquivos Processados", len(arquivos_upload) if arquivos_upload else 0)
                    with col_stat3:
                        st.metric("Fontes Utilizadas", 
                                 (1 if arquivos_upload else 0) + 
                                 (1 if briefing_manual or 'briefing_data' in locals() else 0) +
                                 (1 if transcricoes_texto else 0))
                    
                    # Botões de download
                    extensao = ".html" if formato_saida == "HTML Básico" else ".md" if formato_saida == "Markdown" else ".txt"
                    
                    st.download_button(
                        f"💾 Baixar Conteúdo ({formato_saida})",
                        data=conteudo_gerado,
                        file_name=f"conteudo_gerado_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}{extensao}",
                        mime="text/html" if formato_saida == "HTML Básico" else "text/plain"
                    )
                    
                    # Salvar no histórico se MongoDB disponível
                    if mongo_connected_conteudo:
                        try:
                            from bson import ObjectId
                            historico_data = {
                                "tipo_conteudo": tipo_conteudo,
                                "tom_voz": tom_voz,
                                "palavras_chave": palavras_chave,
                                "numero_palavras": numero_palavras,
                                "conteudo_gerado": conteudo_gerado,
                                "fontes_utilizadas": {
                                    "arquivos_upload": [arquivo.name for arquivo in arquivos_upload] if arquivos_upload else [],
                                    "briefing_manual": bool(briefing_manual),
                                    "transcricoes": len(arquivos_midia) if arquivos_midia else 0
                                },
                                "data_criacao": datetime.datetime.now()
                            }
                            db_briefings['historico_geracao'].insert_one(historico_data)
                            st.success("✅ Conteúdo salvo no histórico!")
                        except Exception as e:
                            st.warning(f"Conteúdo gerado, mas não salvo no histórico: {str(e)}")
                    
                except Exception as e:
                    st.error(f"❌ Erro ao gerar conteúdo: {str(e)}")
                    st.info("💡 Dica: Verifique se os arquivos não estão corrompidos e tente novamente.")

    # Seção de histórico rápido
    if mongo_connected_conteudo:
        with st.expander("📚 Histórico de Gerações Recentes"):
            try:
                historico = list(db_briefings['historico_geracao'].find().sort("data_criacao", -1).limit(5))
                if historico:
                    for item in historico:
                        st.write(f"**{item['tipo_conteudo']}** - {item['data_criacao'].strftime('%d/%m/%Y %H:%M')}")
                        st.caption(f"Palavras-chave: {item.get('palavras_chave', 'Nenhuma')} | Tom: {item['tom_voz']}")
                        with st.expander("Ver conteúdo"):
                            st.write(item['conteudo_gerado'][:500] + "..." if len(item['conteudo_gerado']) > 500 else item['conteudo_gerado'])
                else:
                    st.info("Nenhuma geração no histórico")
            except Exception as e:
                st.warning(f"Erro ao carregar histórico: {str(e)}")

# ========== ABA: BLOG INTELIGENTE COM RAG TÉCNICO + PERPLEXITY ==========
with tab_blog:
    st.header("🌱 Blog Inteligente - Geração Avançada")
    st.markdown("**Cole tudo o que você quer abordar em uma única caixa de texto. O sistema fará o resto.**")
    
    # ============================================
    # 1. INICIALIZAÇÃO E CONFIGURAÇÕES
    # ============================================
    
    # Conexão com MongoDB
    try:
        client_blog_rag = MongoClient(mongo_uri)
        db_blog_rag = client_blog_rag['blog_rag_tecnico']
        collection_posts_rag = db_blog_rag['posts_rag']
        collection_versoes_rag = db_blog_rag['versoes_ajustes']
        mongo_connected_blog_rag = True
    except Exception as e:
        st.error(f"❌ Erro na conexão com MongoDB: {str(e)}")
        mongo_connected_blog_rag = False
    
    # Estado da sessão - INICIALIZAR TODAS AS VARIÁVEIS
    if 'conteudo_gerado_blog' not in st.session_state:
        st.session_state.conteudo_gerado_blog = None
    if 'versoes_blog' not in st.session_state:
        st.session_state.versoes_blog = []
    if 'relatorio_fontes_blog' not in st.session_state:
        st.session_state.relatorio_fontes_blog = None
    if 'briefing_original_blog' not in st.session_state:
        st.session_state.briefing_original_blog = None
    if 'fontes_perplexity_blog' not in st.session_state:
        st.session_state.fontes_perplexity_blog = []
    if 'usou_perplexity_blog' not in st.session_state:
        st.session_state.usou_perplexity_blog = False
    
    # ============================================
    # 2. INTERFACE SIMPLIFICADA - ÚNICA CAIXA DE TEXTO
    # ============================================
    
    st.markdown("---")
    
    # ÁREA PRINCIPAL - ÚNICA CAIXA DE TEXTO
    texto_briefing = st.text_area(
        "📋 **DESCREVA O CONTEÚDO QUE VOCÊ QUER GERAR**",
        height=250,
        placeholder="""Exemplo de briefing completo:

Título: Manejo de nematoides na cultura da soja com produtos biológicos

Cultura: Soja
Problema: Aumento da população de nematoides (Meloidogyne e Heterodera) em solos com palhada de milho
Produtos: NemaControl (bionematicida) e Victrato (bioativador)

Objetivo: Educar o produtor sobre a importância do manejo biológico de nematoides, mostrando resultados práticos e posicionando os produtos como solução eficaz.

Público-alvo: Produtores de soja do Centro-Oeste, nível técnico médio a alto.

Palavras-chave principais: manejo de nematoides, bionematicida, soja
Palavras-chave secundárias: Meloidogyne, Heterodera, tratamento de sementes, produtividade

Observações importantes:
- Tom técnico mas acessível
- Incluir dados de eficácia dos produtos
- Citar resultados de pesquisas da Embrapa
- Evitar linguagem muito comercial
- Referenciar fontes confiáveis
- Texto com ~1500 palavras
- Incluir CTA para falar com consultor técnico

Contexto do mês: Fevereiro - período de desenvolvimento vegetativo da soja no Centro-Oeste, momento crítico para manejo de nematoides.
""",
        key="briefing_unico"
    )
    
    st.markdown("---")
    
    # ============================================
    # 3. CONFIGURAÇÕES AVANÇADAS (EXPANDÍVEIS)
    # ============================================
    
    with st.expander("⚙️ Configurações Avançadas (opcional)", expanded=False):
        col_adv1, col_adv2 = st.columns(2)
        
        with col_adv1:
            # Palavras-chave e densidade
            palavras_chave_input = st.text_input(
                "Palavras-chave (separadas por vírgula):",
                placeholder="ex: manejo de nematoides, bionematicida, soja",
                help="Se não preencher, serão extraídas do briefing"
            )
            
            densidade_palavras = st.slider(
                "Densidade desejada para palavras-chave (%):",
                min_value=1,
                max_value=10,
                value=3,
                step=1,
                help="Percentual aproximado de vezes que as palavras-chave devem aparecer no texto"
            )
            
            palavras_primeira_linha = st.text_input(
                "Palavras que devem aparecer na primeira linha:",
                placeholder="ex: nematoides, soja, manejo",
                help="Estas palavras devem estar obrigatoriamente nas primeiras 100 caracteres"
            )
        
        with col_adv2:
            # Configurações de busca web
            usar_perplexity_blog = st.checkbox(
                "🌐 Buscar informações atualizadas na web", 
                value=True,
                help="Ativa busca no Perplexity para enriquecer com dados atualizados"
            )
            
            if usar_perplexity_blog:
                profundidade_busca = st.select_slider(
                    "Profundidade da busca:",
                    options=["Básica", "Moderada", "Avançada"],
                    value="Avançada"
                )
            
            # Configurações de formato
            tom_voz = st.selectbox(
                "Tom de voz:",
                ["Técnico-científico", "Jornalístico", "Educativo", "Consultivo"],
                index=0
            )
            
            numero_palavras = st.number_input(
                "Número aproximado de palavras:",
                min_value=500,
                max_value=5000,
                value=1500,
                step=100
            )
    
    # ============================================
    # 4. BOTÃO PRINCIPAL - GERAÇÃO
    # ============================================
    
    col_btn1, col_btn2, col_btn3 = st.columns([1, 2, 1])
    with col_btn2:
        if st.button("🚀 GERAR CONTEÚDO DO BLOG", type="primary", use_container_width=True):
            if not texto_briefing.strip():
                st.error("❌ Por favor, descreva o conteúdo que deseja gerar.")
            else:
                with st.spinner("🔄 Processando briefing e gerando conteúdo..."):
                    try:
                        # Salvar briefing original
                        st.session_state.briefing_original_blog = texto_briefing
                        
                        # Processar palavras-chave
                        palavras_chave_lista = []
                        if palavras_chave_input:
                            palavras_chave_lista = [p.strip() for p in palavras_chave_input.split(',') if p.strip()]
                        
                        palavras_primeira_linha_lista = []
                        if palavras_primeira_linha:
                            palavras_primeira_linha_lista = [p.strip() for p in palavras_primeira_linha.split(',') if p.strip()]
                        
                        # ============================================
                        # 5. FUNÇÃO DE BUSCA PERPLEXITY
                        # ============================================
                        
                        def buscar_perplexity_blog(briefing: str, profundidade: str) -> Dict:
                            """Busca informações atualizadas na web"""
                            try:
                                from perplexity import Perplexity
                                
                                perp_api_key = os.getenv("PERP_API_KEY")
                                if not perp_api_key:
                                    return {"erro": "PERP_API_KEY não encontrada", "resultado": None, "fontes": []}
                                
                                client = Perplexity(api_key=perp_api_key)
                                
                                prompt_busca = f"""
                                Você é um pesquisador agrícola. Busque informações técnicas atualizadas e confiáveis sobre:
                                
                                {briefing[:800]}
                                
                                REQUISITOS:
                                1. Fontes: Embrapa, universidades, artigos científicos, boletins técnicos
                                2. Dados concretos: números, estatísticas, resultados de pesquisa
                                3. Informações dos últimos 2-3 anos sempre que possível
                                4. Para CADA informação, forneça a fonte completa
                                
                                FORMATO:
                                ## INFORMAÇÕES ENCONTRADAS
                                
                                ### [Tópico 1]
                                - Informação: [dado técnico]
                                - Fonte: [instituição, ano]
                                - Relevância: [por que é relevante para o tema]
                                - URL/Link: [se disponível]
                                
                                ### [Tópico 2]
                                ...
                                
                                ## LISTA DE FONTES
                                [Lista numerada com todas as fontes utilizadas]
                                """
                                
                                response = client.chat.completions.create(
                                    model="sonar",
                                    messages=[{"role": "user", "content": prompt_busca}],
                                    temperature=0.0,
                                    max_tokens=20000
                                )
                                
                                if response and response.choices:
                                    resultado = response.choices[0].message.content
                                    
                                    # Extrair fontes
                                    fontes = []
                                    linhas = resultado.split('\n')
                                    for linha in linhas:
                                        if 'http://' in linha or 'https://' in linha:
                                            import re
                                            urls = re.findall(r'(https?://[^\s\)]+)', linha)
                                            fontes.extend(urls)
                                        elif 'Fonte:' in linha and '[' not in linha:
                                            fontes.append(linha.strip())
                                    
                                    return {
                                        "erro": None,
                                        "resultado": resultado,
                                        "fontes": list(set(fontes))[:15]
                                    }
                                else:
                                    return {"erro": "Sem resposta", "resultado": None, "fontes": []}
                                    
                            except Exception as e:
                                return {"erro": str(e), "resultado": None, "fontes": []}
                        
                        # ============================================
                        # 6. EXECUTAR BUSCA PERPLEXITY
                        # ============================================
                        
                        resultados_perplexity = {"resultado": None, "fontes": [], "erro": None}
                        
                        if usar_perplexity_blog:
                            with st.spinner("🌐 Buscando informações atualizadas na web..."):
                                resultados_perplexity = buscar_perplexity_blog(texto_briefing, profundidade_busca if 'profundidade_busca' in locals() else "Avançada")
                                
                                if resultados_perplexity.get('erro'):
                                    st.warning(f"⚠️ Busca web: {resultados_perplexity['erro']}")
                                else:
                                    fontes_count = len(resultados_perplexity.get('fontes', []))
                                    st.success(f"✅ {fontes_count} fontes encontradas na web")
                                    
                                    # Salvar no session state
                                    st.session_state.fontes_perplexity_blog = resultados_perplexity.get('fontes', [])
                                    st.session_state.usou_perplexity_blog = True
                        else:
                            st.session_state.usou_perplexity_blog = False
                        
                        # ============================================
                        # 7. CONTEXTO DO AGENTE
                        # ============================================
                        
                        contexto_agente = ""
                        if st.session_state.agente_selecionado:
                            agente = st.session_state.agente_selecionado
                            contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                        
                        # ============================================
                        # 8. PROMPT DE GERAÇÃO AVANÇADA
                        # ============================================
                        
                        prompt_geracao_blog = f"""
                        {contexto_agente}
                        
                        ## INSTRUÇÕES PARA GERAÇÃO DE CONTEÚDO TÉCNICO AGRÍCOLA
                        
                        ### BRIEFING DO USUÁRIO:
                        {texto_briefing}
                        
                        ### INFORMAÇÕES DA WEB (PERPLEXITY) COM FONTES:
                        {resultados_perplexity.get('resultado', 'Nenhuma informação da web disponível.')}
                        
                        ### CONFIGURAÇÕES ESPECÍFICAS:
                        - Tom de voz: {tom_voz}
                        - Número aproximado de palavras: {numero_palavras} (±10%)
                        - Palavras-chave para densidade: {', '.join(palavras_chave_lista) if palavras_chave_lista else 'extraídas do briefing'}
                        - Densidade desejada: {densidade_palavras}%
                        - Palavras obrigatórias na primeira linha: {', '.join(palavras_primeira_linha_lista) if palavras_primeira_linha_lista else 'não especificadas'}
                        
                        ---
                        
                        ## DIRETRIZES CRÍTICAS DE GERAÇÃO (COM BASE EM FEEDBACK DE USUÁRIOS):
                        
                        **1. EVITE REPETIÇÕES E CONTEÚDO RASO:**
                           - NÃO fique repetindo as mesmas informações
                           - APROFUNDE o conteúdo técnico em cada seção
                           - CONSTRUA uma narrativa coesa que evolui do problema para a solução
                           - POSICIONE os produtos de forma estratégica, explicando seu papel no manejo
                           - DISCORRA sobre os benefícios com dados concretos, não apenas afirmações genéricas
                        
                        **2. QUEBRE PARÁGRAFOS LONGOS:**
                           - Parágrafos com NO MÁXIMO 4-5 linhas
                           - Use subtítulos (H2, H3) para organizar o conteúdo
                           - Use listas com bullets para informações concatenadas (máx 5 itens)
                           - Facilite a escaneabilidade do texto
                        
                        **3. LINKS ANCORADOS:**
                           - Quando usar informações da busca web, ANCRE links relevantes
                           - Formato: "informação [fonte](url)"
                           - Ao final, liste todas as referências completas
                        
                        **4. REFERÊNCIAS E JUSTIFICATIVAS:**
                           - Para CADA dado técnico apresentado, CITE a fonte
                           - JUSTIFIQUE por que aquela informação é relevante no contexto
                           - CONSTRUA a credibilidade do texto com fontes confiáveis
                        
                        **5. NARRATIVA E POSICIONAMENTO:**
                           - Comece contextualizando o problema agrícola
                           - Apresente dados que mostram a gravidade/importância
                           - Introduza as soluções (produtos) de forma natural
                           - EXPLIQUE o modo de ação, não apenas o nome
                           - MOSTRE resultados com dados de eficácia
                           - CONCLUA com recomendações práticas e CTA
                        
                        ---
                        
                        ## ESTRUTURA SUGERIDA (ADAPTE CONFORME O BRIEFING):
                        
                        # [TÍTULO PRINCIPAL COM PALAVRA-CHAVE]
                        
                        [Introdução contextualizando o problema - 2-3 parágrafos curtos]
                        
                        ## [PROBLEMA/DESAFIO TÉCNICO]
                        
                        [Parágrafo explicando o problema]
                        [Parágrafo com dados sobre impacto econômico/perdas]
                        
                        ### [Subtópico específico do problema - ex: Principais espécies]
                        [Conteúdo com dados e fontes]
                        
                        ## [SOLUÇÕES/MANEJO]
                        
                        [Parágrafo introdutório sobre manejo integrado]
                        
                        ### [Produto/Solução 1]
                        [O que é, modo de ação, benefícios, dados de eficácia, como aplicar]
                        
                        ### [Produto/Solução 2]
                        [O que é, modo de ação, benefícios, dados de eficácia, como aplicar]
                        
                        ## [RESULTADOS E BENEFÍCIOS]
                        
                        [Parágrafos com dados de campo, resultados de pesquisa, depoimentos técnicos]
                        
                        ## [RECOMENDAÇÕES PRÁTICAS]
                        
                        [Orientações para implementação, época de aplicação, doses, cuidados]
                        
                        ## [CONCLUSÃO]
                        
                        [Resumo dos pontos principais e CTA]
                        
                        ---
                        
                        ## LISTA DE REFERÊNCIAS
                        
                        [Listar todas as fontes utilizadas ao longo do texto]
                        
                        ---
                        
                        **GERAR O CONTEÚDO COMPLETO SEGUINDO TODAS AS DIRETRIZES ACIMA.**
                        **O CONTEÚDO DEVE SER TÉCNICO, PROFUNDO, BEM ESTRUTURADO E COM TODAS AS FONTES CITADAS.**
                        """
                        
                        # Gerar conteúdo
                        resposta = modelo_texto.generate_content(prompt_geracao_blog)
                        conteudo_gerado = resposta.text
                        
                        # ============================================
                        # 9. PÓS-PROCESSAMENTO
                        # ============================================
                        
                        # Construir relatório de fontes
                        relatorio_fontes = "## 📚 REFERÊNCIAS E FONTES UTILIZADAS\n\n"
                        
                        if resultados_perplexity.get('fontes'):
                            relatorio_fontes += "### Fontes da Web:\n"
                            for i, fonte in enumerate(resultados_perplexity['fontes'], 1):
                                relatorio_fontes += f"{i}. {fonte}\n"
                        else:
                            relatorio_fontes += "*Nenhuma fonte web específica foi capturada.*\n"
                        
                        # Salvar no session state
                        st.session_state.conteudo_gerado_blog = conteudo_gerado
                        st.session_state.relatorio_fontes_blog = relatorio_fontes
                        
                        # Salvar primeira versão
                        st.session_state.versoes_blog = [{
                            "versao": 1,
                            "conteudo": conteudo_gerado,
                            "data": datetime.datetime.now(),
                            "descricao": "Geração inicial"
                        }]
                        
                        # Salvar no MongoDB
                        if mongo_connected_blog_rag:
                            try:
                                documento = {
                                    "briefing": texto_briefing,
                                    "conteudo": conteudo_gerado,
                                    "fontes": resultados_perplexity.get('fontes', []),
                                    "configuracoes": {
                                        "tom_voz": tom_voz,
                                        "palavras_chave": palavras_chave_lista,
                                        "usou_perplexity": usar_perplexity_blog
                                    },
                                    "data_criacao": datetime.datetime.now()
                                }
                                collection_posts_rag.insert_one(documento)
                            except Exception as e:
                                st.warning(f"⚠️ Conteúdo gerado mas não salvo no banco: {str(e)}")
                        
                        st.success("✅ Conteúdo gerado com sucesso!")
                        st.rerun()
                        
                    except Exception as e:
                        st.error(f"❌ Erro na geração: {str(e)}")
    
    # ============================================
    # 10. EXIBIÇÃO DO RESULTADO
    # ============================================
    
    if st.session_state.conteudo_gerado_blog:
        st.markdown("---")
        
        # Métricas - USANDO SESSION STATE EM VEZ DE VARIÁVEIS LOCAIS
        palavras_count = len(st.session_state.conteudo_gerado_blog.split())
        
        col_m1, col_m2, col_m3, col_m4 = st.columns(4)
        with col_m1:
            st.metric("📊 Palavras", palavras_count)
        with col_m2:
            versoes = len(st.session_state.versoes_blog)
            st.metric("📋 Versões", versoes)
        with col_m3:
            # Pegar tom_voz do session state ou usar padrão
            tom_exibicao = st.session_state.get('tom_voz_blog', 'Técnico-científico')
            st.metric("🎯 Tom", tom_exibicao)
        with col_m4:
            # Usar session state para verificar se usou perplexity
            usou_perplexity = st.session_state.get('usou_perplexity_blog', False)
            tem_fontes = len(st.session_state.get('fontes_perplexity_blog', [])) > 0
            st.metric("🌐 Fontes", "✅" if usou_perplexity and tem_fontes else "❌")
        
        # Abas para visualização
        tab_conteudo, tab_ref, tab_versoes, tab_export = st.tabs([
            "📝 Conteúdo Gerado", "📚 Referências", "📋 Histórico", "💾 Exportar"
        ])
        
        with tab_conteudo:
            st.markdown(st.session_state.conteudo_gerado_blog)
        
        with tab_ref:
            if st.session_state.relatorio_fontes_blog:
                st.markdown(st.session_state.relatorio_fontes_blog)
            else:
                st.info("Nenhuma referência disponível")
        
        with tab_versoes:
            if st.session_state.versoes_blog:
                for versao in reversed(st.session_state.versoes_blog[-5:]):
                    with st.expander(f"Versão {versao['versao']} - {versao['data'].strftime('%d/%m/%Y %H:%M') if isinstance(versao['data'], datetime.datetime) else 'Data desconhecida'} - {versao['descricao']}"):
                        st.text_area(
                            f"Conteúdo da versão {versao['versao']}",
                            value=versao['conteudo'][:500] + "..." if len(versao['conteudo']) > 500 else versao['conteudo'],
                            height=200,
                            key=f"versao_{versao['versao']}"
                        )
                        
                        if st.button(f"Restaurar versão {versao['versao']}", key=f"restore_{versao['versao']}"):
                            st.session_state.conteudo_gerado_blog = versao['conteudo']
                            st.success(f"✅ Versão {versao['versao']} restaurada!")
                            st.rerun()
            else:
                st.info("Nenhuma versão disponível")
        
        with tab_export:
            col_exp1, col_exp2 = st.columns(2)
            
            with col_exp1:
                # TXT
                st.download_button(
                    "📥 Baixar como TXT",
                    data=st.session_state.conteudo_gerado_blog,
                    file_name=f"blog_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                    mime="text/plain",
                    use_container_width=True
                )
                
                # MD
                st.download_button(
                    "📥 Baixar como MD",
                    data=st.session_state.conteudo_gerado_blog,
                    file_name=f"blog_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.md",
                    mime="text/markdown",
                    use_container_width=True
                )
            
            with col_exp2:
                # Referências
                if st.session_state.relatorio_fontes_blog:
                    st.download_button(
                        "📥 Baixar Referências",
                        data=st.session_state.relatorio_fontes_blog,
                        file_name=f"referencias_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.md",
                        mime="text/markdown",
                        use_container_width=True
                    )
                
                # Pacote completo
                pacote = f"""# BLOG POST - {datetime.datetime.now().strftime('%d/%m/%Y')}

## BRIEFING ORIGINAL
{st.session_state.briefing_original_blog if st.session_state.briefing_original_blog else 'N/A'}

## CONTEÚDO GERADO
{st.session_state.conteudo_gerado_blog}

## REFERÊNCIAS
{st.session_state.relatorio_fontes_blog if st.session_state.relatorio_fontes_blog else 'N/A'}
"""
                st.download_button(
                    "📦 Pacote Completo",
                    data=pacote,
                    file_name=f"pacote_completo_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                    mime="text/plain",
                    use_container_width=True
                )
        
        # ============================================
        # 11. SEÇÃO DE AJUSTES
        # ============================================
        
        st.markdown("---")
        st.subheader("🔄 Ajustar Conteúdo")
        
        col_ajuste1, col_ajuste2 = st.columns([3, 1])
        
        with col_ajuste1:
            solicitacao_ajuste = st.text_area(
                "Descreva os ajustes desejados:",
                placeholder="Exemplos:\n- Aprofunde mais na seção sobre modo de ação dos produtos\n- Adicione mais dados de eficácia com fontes\n- Melhore a narrativa, conectando melhor problema e solução\n- Quebre parágrafos longos no início\n- Inclua mais informações sobre a cultura alvo\n- Aumente a densidade da palavra-chave 'manejo de nematoides'\n- Adicione um CTA mais forte no final",
                height=100,
                key="campo_ajuste_blog"
            )
        
        with col_ajuste2:
            st.markdown("#####")
            if st.button("✅ APLICAR AJUSTES", type="secondary", use_container_width=True):
                if solicitacao_ajuste.strip():
                    with st.spinner("🔄 Aplicando ajustes..."):
                        try:
                            # Preparar prompt de ajuste
                            prompt_ajuste = f"""
                            ## CONTEÚDO ATUAL:
                            {st.session_state.conteudo_gerado_blog}
                            
                            ## BRIEFING ORIGINAL:
                            {st.session_state.briefing_original_blog if st.session_state.briefing_original_blog else 'N/A'}
                            
                            ## AJUSTES SOLICITADOS:
                            {solicitacao_ajuste}
                            
                            ## INSTRUÇÕES:
                            1. APLIQUE os ajustes solicitados mantendo a estrutura geral
                            2. APROFUNDE o conteúdo técnico quando necessário
                            3. QUEBRE parágrafos longos
                            4. MANTENHA as fontes e referências
                            5. MELHORE a narrativa se solicitado
                            6. POSICIONE os produtos de forma estratégica
                            
                            RETORNE APENAS O CONTEÚDO AJUSTADO.
                            """
                            
                            resposta_ajuste = modelo_texto.generate_content(prompt_ajuste)
                            conteudo_ajustado = resposta_ajuste.text
                            
                            # Salvar versão anterior
                            nova_versao = {
                                "versao": len(st.session_state.versoes_blog) + 1,
                                "conteudo": st.session_state.conteudo_gerado_blog,
                                "data": datetime.datetime.now(),
                                "descricao": f"Ajuste: {solicitacao_ajuste[:50]}..."
                            }
                            st.session_state.versoes_blog.append(nova_versao)
                            
                            # Atualizar conteúdo atual
                            st.session_state.conteudo_gerado_blog = conteudo_ajustado
                            
                            st.success("✅ Ajustes aplicados com sucesso!")
                            st.rerun()
                            
                        except Exception as e:
                            st.error(f"❌ Erro ao aplicar ajustes: {str(e)}")
                else:
                    st.warning("⚠️ Descreva os ajustes desejados.")
    
    # ============================================
    # 12. HISTÓRICO DE GERAÇÕES
    # ============================================
    
    if mongo_connected_blog_rag:
        st.markdown("---")
        st.subheader("📚 Histórico de Gerações")
        
        try:
            historico = list(collection_posts_rag.find().sort("data_criacao", -1).limit(5))
            
            if historico:
                for post in historico:
                    with st.expander(f"📄 {post.get('data_criacao', '').strftime('%d/%m/%Y %H:%M') if post.get('data_criacao') else 'Data desconhecida'} - Briefing: {post.get('briefing', '')[:100]}..."):
                        st.write(f"**Palavras:** {len(post.get('conteudo', '').split())}")
                        st.write(f"**Fontes:** {len(post.get('fontes', []))}")
                        
                        if st.button(f"Carregar este post", key=f"load_{post.get('_id')}"):
                            st.session_state.conteudo_gerado_blog = post.get('conteudo', '')
                            st.session_state.briefing_original_blog = post.get('briefing', '')
                            st.success("✅ Post carregado!")
                            st.rerun()
            else:
                st.info("Nenhum post no histórico")
                
        except Exception as e:
            st.warning(f"Erro ao carregar histórico: {str(e)}")

# ========== ABA: REVISÃO ORTOGRÁFICA ==========
with tab_revisao_ortografica:
    st.header("📝 Revisão Ortográfica")
    
    texto_para_revisao = st.text_area("Cole o texto que deseja revisar:", height=300)
    
    if st.button("🔍 Realizar Revisão Ortográfica", type="primary"):
        if texto_para_revisao:
            with st.spinner("Revisando texto..."):
                try:
                    # Usar contexto do agente selecionado se disponível
                    if st.session_state.agente_selecionado:
                        agente = st.session_state.agente_selecionado
                        contexto = construir_contexto(agente, st.session_state.segmentos_selecionados)
                        prompt = f"""
                        
                        Faça uma revisão ortográfica e gramatical completa do seguinte texto:
                        
                        ###BEGIN TEXTO A SER REVISADO###
                        {texto_para_revisao}
                        ###END TEXTO A SER REVISADO###
                        
                        MANTENHA A ESTRUTURA DO TEXTO ORIGINAL. APENAS CORRIJA ERROS ORTOGRÁFICOS (SE PRESENTES) E APONTE QUAIS FORAM OS ERROS CORRIGIDOS
                        """
                    else:
                        prompt = f"""
                        Faça uma revisão ortográfica e gramatical completa do seguinte texto:
                        
                        ###BEGIN TEXTO A SER REVISADO###
                        {texto_para_revisao}
                        ###END TEXTO A SER REVISADO###
                        
                        MANTENHA A ESTRUTURA DO TEXTO ORIGINAL. APENAS CORRIJA ERROS ORTOGRÁFICOS (SE PRESENTES) E APONTE QUAIS FORAM OS ERROS CORRIGIDOS
                        """
                    
                    resposta = modelo_texto.generate_content(prompt)
                    st.subheader("📋 Resultado da Revisão")
                    st.markdown(resposta.text)
                    
                except Exception as e:
                    st.error(f"Erro na revisão: {str(e)}")
        else:
            st.warning("Por favor, cole um texto para revisão.")

# ========== ABA: REVISÃO TÉCNICA (VERSÃO COMPLETA COM RELATÓRIO DE MUDANÇAS) ==========
with tab_revisao_tecnica:
    st.header("🔧 Revisão Técnica com RAGs Especializados")
    st.markdown("**Análise em camadas: taxonomia, epidemiologia, produtos + reescrita final com relatório detalhado**")
    
    # Layout com duas colunas principais
    col_original_rag, col_revisado_rag = st.columns(2)
    
    with col_original_rag:
        st.subheader("📄 Conteúdo Original")
        texto_tecnico = st.text_area(
            "Cole o conteúdo técnico para revisão:", 
            height=300,
            placeholder="Cole aqui o conteúdo técnico agrícola que precisa ser revisado...",
            key="texto_tecnico_rag",
            label_visibility="collapsed"
        )
    
    with col_revisado_rag:
        st.subheader("✨ Conteúdo Revisado com RAG")
        # Placeholder para o conteúdo revisado com RAG
        revisado_rag_placeholder = st.empty()
        revisado_rag_placeholder.info("📝 Aguardando revisão com RAG... O conteúdo revisado aparecerá aqui.")
    
    # Configurações da revisão (abaixo das colunas)
    st.markdown("---")
    st.subheader("⚙️ Configurações da Revisão")
    
    col_config1, col_config2, col_config3 = st.columns([2, 1, 1])
    
    with col_config1:
        # Tipo de conteúdo específico
        tipo_conteudo = st.selectbox(
            "Tipo de Conteúdo:",
            ["Artigo Técnico", "Material Comercial", "Blog Post", "Manual Técnico", "Comunicado Técnico"],
            help="Define o rigor da revisão"
        )
    
    with col_config2:
        st.subheader("🔍 RAGs Especializados")
        
        rag_taxonomia = st.checkbox("RAG Taxonomia", value=True, 
                                  help="Busca específica por classificação de patógenos")
        rag_epidemiologia = st.checkbox("RAG Epidemiologia", value=True,
                                      help="Busca específica por condições ambientais")
        rag_produtos = st.checkbox("RAG Produtos", value=True,
                                 help="Busca específica por informações de produtos")
        rag_geral = st.checkbox("RAG Geral", value=True,
                              help="Busca geral por similaridade semântica")
    
    with col_config3:
        st.subheader("⚙️ Configurações")
        
        nivel_rigor = st.select_slider(
            "Nível de Rigor:",
            ["Leve", "Moderado", "Rigoroso", "Especialista"]
        )
        
        limite_documentos = st.number_input("Docs por RAG", min_value=3, max_value=20, value=12,
                                          help="Número de documentos resgatados por RAG especializado")
        
        usar_contexto_agente = st.checkbox("Usar contexto do agente", 
                                         value=bool(st.session_state.agente_selecionado))
        
        # NOVA OPÇÃO: Incluir relatório detalhado
        incluir_relatorio = st.checkbox("📋 Incluir relatório de mudanças", value=True,
                                      help="Gera um relatório detalhado mostrando todas as alterações")

    # Funções para RAGs especializados (mantidas iguais)
    def realizar_rag_taxonomia(texto: str, limite: int = 12) -> List[Dict]:
        """RAG especializado em taxonomia e classificação de patógenos"""
        perguntas_especificas = [
            "classificação taxonômica",
            "fungo ou oomiceto",
            "nome científico patógeno", 
            "reino filo classe ordem",
            "agente causal doença",
            "Peronospora Phakopsora Corynespora",
            "oomiceto vs fungo diferença",
            "taxonomia fitopatologia"
        ]
        
        documentos_combinados = []
        for pergunta in perguntas_especificas:
            query = f"{texto[:200]} {pergunta}"
            embedding = get_embedding(query)
            documentos = astra_client.vector_search(ASTRA_DB_COLLECTION, embedding, limit=limite//len(perguntas_especificas))
            documentos_combinados.extend(documentos)
        
        # Remover duplicados
        documentos_unicos = []
        ids_vistos = set()
        for doc in documentos_combinados:
            doc_id = str(doc.get('_id', ''))
            if doc_id not in ids_vistos:
                documentos_unicos.append(doc)
                ids_vistos.add(doc_id)
        
        return documentos_unicos[:limite]

    def realizar_rag_epidemiologia(texto: str, limite: int = 12) -> List[Dict]:
        """RAG especializado em condições epidemiológicas"""
        perguntas_especificas = [
            "condições ambientais doença",
            "temperatura umidade molhamento foliar",
            "condições ideais infecção",
            "epidemiologia doença plantas",
            "período molhamento temperatura ótima",
            "umidade relativa infecção",
            "condições climáticas favoráveis",
            "fatores epidemiológicos"
        ]
        
        documentos_combinados = []
        for pergunta in perguntas_especificas:
            query = f"{texto[:200]} {pergunta}"
            embedding = get_embedding(query)
            documentos = astra_client.vector_search(ASTRA_DB_COLLECTION, embedding, limit=limite//len(perguntas_especificas))
            documentos_combinados.extend(documentos)
        
        # Remover duplicados
        documentos_unicos = []
        ids_vistos = set()
        for doc in documentos_combinados:
            doc_id = str(doc.get('_id', ''))
            if doc_id not in ids_vistos:
                documentos_unicos.append(doc)
                ids_vistos.add(doc_id)
        
        return documentos_unicos[:limite]

    def realizar_rag_produtos(texto: str, limite: int = 12) -> List[Dict]:
        """RAG especializado em informações de produtos"""
        perguntas_especificas = [
            "modo de ação produto",
            "aplicação dose recomendada",
            "eficácia controle doença",
            "características técnicas produto",
            "benefícios produto agrícola",
            "tecnologia aplicação",
            "resultados eficácia",
            "recomendações uso produto"
        ]
        
        documentos_combinados = []
        for pergunta in perguntas_especificas:
            query = f"{texto[:200]} {pergunta}"
            embedding = get_embedding(query)
            documentos = astra_client.vector_search(ASTRA_DB_COLLECTION, embedding, limit=limite//len(perguntas_especificas))
            documentos_combinados.extend(documentos)
        
        # Remover duplicados
        documentos_unicos = []
        ids_vistos = set()
        for doc in documentos_combinados:
            doc_id = str(doc.get('_id', ''))
            if doc_id not in ids_vistos:
                documentos_unicos.append(doc)
                ids_vistos.add(doc_id)
        
        return documentos_unicos[:limite]

    def realizar_rag_geral(texto: str, limite: int = 12) -> List[Dict]:
        """RAG geral por similaridade semântica"""
        embedding = get_embedding(texto[:800])
        documentos = astra_client.vector_search(ASTRA_DB_COLLECTION, embedding, limit=limite)
        return documentos

    def processar_rags_especializados(texto: str, rags_ativos: dict, limite: int = 12) -> dict:
        """Executa todos os RAGs especializados e retorna resultados consolidados"""
        resultados = {}
        
        if rags_ativos.get('taxonomia'):
            with st.spinner("🔬 Buscando informações de taxonomia..."):
                resultados['taxonomia'] = realizar_rag_taxonomia(texto, limite)
        
        if rags_ativos.get('epidemiologia'):
            with st.spinner("🌡️ Buscando informações epidemiológicas..."):
                resultados['epidemiologia'] = realizar_rag_epidemiologia(texto, limite)
        
        if rags_ativos.get('produtos'):
            with st.spinner("🧪 Buscando informações de produtos..."):
                resultados['produtos'] = realizar_rag_produtos(texto, limite)
        
        if rags_ativos.get('geral'):
            with st.spinner("📚 Buscando informações gerais..."):
                resultados['geral'] = realizar_rag_geral(texto, limite)
        
        return resultados

    # NOVA FUNÇÃO: Reescrita com relatório detalhado de mudanças
    def reescrever_com_relatorio_mudancas(texto_original: str, resultados_rags: dict, contexto_agente: str = "") -> tuple:
        """Reescreve o conteúdo e gera um relatório detalhado das mudanças"""
        
        # Construir contexto consolidado dos RAGs
        contexto_rags = "## DOCUMENTOS TÉCNICOS DE REFERÊNCIA:\n\n"
        
        for categoria, documentos in resultados_rags.items():
            if documentos:
                contexto_rags += f"### {categoria.upper()} ({len(documentos)} documentos):\n"
                for i, doc in enumerate(documentos, 1):
                    doc_content = str(doc)
                    doc_limpo = doc_content.replace('{', '').replace('}', '').replace("'", "").replace('"', '')
                    if len(doc_limpo) > 300:
                        doc_limpo = doc_limpo[:300] + "..."
                    contexto_rags += f"- {doc_limpo}\n"
                contexto_rags += "\n"

        # Prompt para reescrita COM relatório
        prompt_reescrita = f"""
        {contexto_agente}

        ## TEXTO ORIGINAL PARA REESCRITA:
        {texto_original}

        ## BASE TÉCNICA DE REFERÊNCIA:
        {contexto_rags}

        ## INSTRUÇÕES CRÍTICAS:

        **SUA TAREFA:** 
        1. Reescrever o texto original aplicando correções técnicas baseadas nos documentos de referência
        2. Gerar um relatório DETALHADO de TODAS as mudanças realizadas
        3. Você deve manter a estrutura original do texto. Você deve realizar apenas mudanças e enriquecimentos conforme o contexto novo vindo da base técnica de referência. O texto original deve sempre ser o molde a ser seguido.

        **FORMATO DE SAÍDA EXIGIDO (use exatamente esta estrutura):**

        ### 📝 TEXTO REESCRITO
        [AQUI VOCÊ COLA O TEXTO COMPLETO REESCRITO E CORRIGIDO]

        ### 🔍 RELATÓRIO DETALHADO DE MUDANÇAS

        #### 📊 RESUMO EXECUTIVO
        - Total de correções aplicadas: [N]
        - Principais categorias de ajustes: [lista categorias]
        - Impacto na precisão técnica: [Alto/Médio/Baixo]

        #### 📋 MUDANÇAS DETALHADAS

        **1. CORREÇÕES TAXONÔMICAS:**
        [Lista cada correção taxonômica no formato:
        - **Original:** "texto original"
        - **Corrigido:** "texto corrigido" 
        - **Justificativa:** explicação técnica baseada nos documentos]

        **2. PRECISÃO EPIDEMIOLÓGICA:**
        [Lista cada correção epidemiológica no formato:
        - **Original:** "texto original"
        - **Corrigido:** "texto corrigido"
        - **Justificativa:** explicação com base científica]

        **3. INFORMAÇÕES DE PRODUTOS:**
        [Lista cada correção de produtos no formato:
        - **Original:** "texto original" 
        - **Corrigido:** "texto corrigido"
        - **Justificativa:** ajuste técnico necessário]

        **4. TERMINOLOGIA TÉCNICA:**
        [Lista cada ajuste de terminologia no formato:
        - **Original:** "termo vago/impreciso"
        - **Corrigido:** "termo técnico preciso"
        - **Justificativa:** padronização técnica]

        **5. DADOS E ESTATÍSTICAS:**
        [Lista cada correção de dados no formato:
        - **Original:** "dado impreciso"
        - **Corrigido:** "dado corrigido"
        - **Justificativa:** fonte/documento de referência]

        #### 🎯 IMPACTO DAS CORREÇÕES
        - Melhorias na precisão científica: [lista específica]
        - Ajustes na comunicação técnica: [lista específica]
        - Correções de segurança da informação: [lista específica]

        **CORREÇÕES TÉCNICAS OBRIGATÓRIAS:**

        1. **PRECISÃO TAXONÔMICA:**
           - Corrigir "fungo" para "oomiceto" quando aplicável
           - Validar nomes científicos e classificação
           - Ajustar descrições de ciclo de vida

        2. **ESPECIFICIDADE EPIDEMIOLÓGICA:**
           - Substituir termos vagos por faixas específicas
           - Especificar temperaturas exatas
           - Definir períodos de molhamento foliar
           - Vincular condições ao fechamento do dossel

        3. **DESCRIÇÃO PRECISA DE SINTOMAS:**
           - Corrigir descrições imprecisas
           - Especificar localização nas plantas
           - Detalhar evolução dos sintomas
           - Ajustar terminologia técnica

        4. **MANEJO E TIMING:**
           - Alinhar mensagens sobre timing de aplicação
           - Esclarecer momentos diferentes
           - Especificar rotação de MoA

        5. **INFORMAÇÕES DE PRODUTOS:**
           - Corrigir claims imprecisos
           - Especificar "conforme bula" quando necessário
           - Validar números de eficácia
           - Ajustar claims técnicos com precisão

        **REGRAS ADICIONAIS:**
        - Mantenha a estrutura e formatação do original
        - Preserve títulos, subtítulos e marcações
        - Apenas corrija o conteúdo técnico, não reinvente a estrutura
        - Se não houver informações nos RAGs para corrigir algo específico, mantenha o original
        - Para CADA mudança, forneça justificativa técnica específica

        **RETORNE EXATAMENTE no formato especificado acima.**
        """

        try:
            resposta = modelo_texto.generate_content(prompt_reescrita)
            texto_completo = resposta.text
            
            # Separar o texto reescrito do relatório
            if "### 📝 TEXTO REESCRITO" in texto_completo and "### 🔍 RELATÓRIO DETALHADO DE MUDANÇAS" in texto_completo:
                partes = texto_completo.split("### 🔍 RELATÓRIO DETALHADO DE MUDANÇAS")
                texto_reescrito = partes[0].replace("### 📝 TEXTO REESCRITO", "").strip()
                relatorio_mudancas = "### 🔍 RELATÓRIO DETALHADO DE MUDANÇAS" + partes[1]
            else:
                # Fallback se o formato não for seguido
                texto_reescrito = texto_completo
                relatorio_mudancas = "### ❌ Relatório não gerado automaticamente\nO modelo não seguiu o formato solicitado para o relatório."
            
            return texto_reescrito, relatorio_mudancas
            
        except Exception as e:
            st.error(f"Erro na reescrita: {str(e)}")
            return texto_original, f"### ❌ Erro na geração do relatório\n{str(e)}"

    def reescrever_sem_relatorio(texto_original: str, resultados_rags: dict, contexto_agente: str = "") -> str:
        """Reescreve o conteúdo sem gerar relatório (para opção rápida)"""
        
        contexto_rags = "## DOCUMENTOS TÉCNICOS DE REFERÊNCIA:\n\n"
        
        for categoria, documentos in resultados_rags.items():
            if documentos:
                contexto_rags += f"### {categoria.upper()} ({len(documentos)} documentos):\n"
                for i, doc in enumerate(documentos, 1):
                    doc_content = str(doc)
                    doc_limpo = doc_content.replace('{', '').replace('}', '').replace("'", "").replace('"', '')
                    if len(doc_limpo) > 300:
                        doc_limpo = doc_limpo[:300] + "..."
                    contexto_rags += f"- {doc_limpo}\n"
                contexto_rags += "\n"

        prompt_rapido = f"""
        {contexto_agente}

        ## TEXTO ORIGINAL PARA REESCRITA:
        {texto_original}

        ## BASE TÉCNICA DE REFERÊNCIA:
        {contexto_rags}

        **REESCREVA o texto aplicando correções técnicas baseadas nos documentos.**
        **RETORNE APENAS o texto reescrito, sem comentários ou relatórios.**

        Correções obrigatórias:
        - Precisão taxonômica (fungo vs oomiceto)
        - Especificidade epidemiológica (temperaturas, umidades)
        - Informações precisas de produtos
        - Terminologia técnica adequada

        Mantenha a estrutura original.
        """

        resposta = modelo_texto.generate_content(prompt_rapido)
        return resposta.text.strip()

    # Botão de revisão técnica com RAGs especializados - AGORA CENTRALIZADO
    st.markdown("---")
    col_btn1, col_btn2, col_btn3 = st.columns([1, 2, 1])
    
    with col_btn2:
        if st.button("🔬 Realizar Revisão com RAGs Especializados", type="primary", use_container_width=True):
            if texto_tecnico:
                # Configurar RAGs ativos
                rags_ativos = {
                    'taxonomia': rag_taxonomia,
                    'epidemiologia': rag_epidemiologia, 
                    'produtos': rag_produtos,
                    'geral': rag_geral
                }
                
                # Construir contexto do agente se solicitado
                contexto_agente = ""
                if usar_contexto_agente and st.session_state.agente_selecionado:
                    agente = st.session_state.agente_selecionado
                    contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                
                with st.spinner("🚀 Executando pipeline de RAGs especializados..."):
                    try:
                        # FASE 1: Executar RAGs especializados
                        st.subheader("📡 Fase 1: Busca com RAGs Especializados")
                        
                        resultados_rags = processar_rags_especializados(texto_tecnico, rags_ativos, limite_documentos)
                        
                        # Mostrar estatísticas dos RAGs
                        col_rag1, col_rag2, col_rag3, col_rag4 = st.columns(4)
                        with col_rag1:
                            st.metric("RAG Taxonomia", 
                                     len(resultados_rags.get('taxonomia', [])),
                                     help="Documentos sobre classificação de patógenos")
                        with col_rag2:
                            st.metric("RAG Epidemiologia", 
                                     len(resultados_rags.get('epidemiologia', [])),
                                     help="Documentos sobre condições ambientais")
                        with col_rag3:
                            st.metric("RAG Produtos", 
                                     len(resultados_rags.get('produtos', [])),
                                     help="Documentos sobre produtos e eficácia")
                        with col_rag4:
                            st.metric("RAG Geral", 
                                     len(resultados_rags.get('geral', [])),
                                     help="Documentos por similaridade semântica")
                        
                        # FASE 2: Reescrita com LLM
                        st.subheader("✍️ Fase 2: Reescrita com Base nos RAGs")
                        
                        with st.spinner("Reescrevendo conteúdo e gerando relatório de mudanças..."):
                            # Escolher qual função de reescrita usar baseado na configuração
                            if incluir_relatorio:
                                texto_reescrito, relatorio_mudancas = reescrever_com_relatorio_mudancas(
                                    texto_tecnico, resultados_rags, contexto_agente
                                )
                            else:
                                texto_reescrito = reescrever_sem_relatorio(texto_tecnico, resultados_rags, contexto_agente)
                                relatorio_mudancas = None
                        
                        # FASE 3: Atualizar visualização lado a lado
                        st.subheader("📋 Fase 3: Resultados da Revisão")
                        
                        # Atualizar a coluna direita com o conteúdo revisado
                        with col_revisado_rag:
                            revisado_rag_placeholder.empty()
                            st.success("✅ Conteúdo revisado com RAGs!")
                            
                            # Criar abas para organizar o conteúdo revisado
                            if incluir_relatorio and relatorio_mudancas:
                                tab_texto_reescrito, tab_relatorio_mudancas, tab_analise = st.tabs([
                                    "📝 Texto Reescrito", "📋 Relatório de Mudanças", "📊 Análise RAGs"
                                ])
                                
                                with tab_texto_reescrito:
                                    st.text_area(
                                        "Texto reescrito com base nos RAGs:",
                                        texto_reescrito,
                                        height=300,
                                        label_visibility="collapsed"
                                    )
                                
                                with tab_relatorio_mudancas:
                                    st.markdown(relatorio_mudancas)
                                
                                with tab_analise:
                                    # Estatísticas de comparação
                                    palavras_orig = len(texto_tecnico.split())
                                    palavras_reesc = len(texto_reescrito.split())
                                    diff_palavras = palavras_reesc - palavras_orig
                                    
                                    col_stat1, col_stat2, col_stat3 = st.columns(3)
                                    with col_stat1:
                                        st.metric("Palavras Original", palavras_orig)
                                    with col_stat2:
                                        st.metric("Palavras Reescrito", palavras_reesc)
                                    with col_stat3:
                                        st.metric("Diferença", 
                                                 f"{'+' if diff_palavras > 0 else ''}{diff_palavras}",
                                                 delta=f"{diff_palavras/palavras_orig*100:.1f}%" if palavras_orig > 0 else "0%")
                                    
                                    # Estatísticas dos RAGs
                                    st.markdown("### 📊 Estatísticas dos RAGs")
                                    for categoria, documentos in resultados_rags.items():
                                        if documentos:
                                            st.write(f"**{categoria.capitalize()}:** {len(documentos)} documentos encontrados")
                            else:
                                # Sem relatório - apenas mostrar texto reescrito
                                st.text_area(
                                    "Texto reescrito com base nos RAGs:",
                                    texto_reescrito,
                                    height=300,
                                    label_visibility="collapsed"
                                )
                        
                        # Botões de download
                        st.markdown("---")
                        col_dl1, col_dl2, col_dl3 = st.columns(3)
                        
                        with col_dl1:
                            st.download_button(
                                "💾 Baixar Texto Reescrito",
                                data=texto_reescrito,
                                file_name=f"texto_reescrito_rags_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                                mime="text/plain",
                                use_container_width=True
                            )
                        
                        with col_dl2:
                            if incluir_relatorio and relatorio_mudancas:
                                st.download_button(
                                    "💾 Baixar Relatório",
                                    data=relatorio_mudancas,
                                    file_name=f"relatorio_mudancas_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.md",
                                    mime="text/markdown",
                                    use_container_width=True
                                )
                        
                        with col_dl3:
                            # Pacote completo
                            pacote_completo = f"TEXTO ORIGINAL:\n{texto_tecnico}\n\n"
                            pacote_completo += "="*60 + "\n\n"
                            pacote_completo += f"TEXTO REESCRITO COM RAGs:\n{texto_reescrito}\n\n"
                            if incluir_relatorio and relatorio_mudancas:
                                pacote_completo += "="*60 + "\n\n"
                                pacote_completo += f"RELATÓRIO DE MUDANÇAS:\n{relatorio_mudancas}"
                            
                            st.download_button(
                                "📦 Baixar Pacote Completo",
                                data=pacote_completo,
                                file_name=f"revisao_completa_rags_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                                mime="text/plain",
                                use_container_width=True
                            )
                        
                        # Salvar no histórico se MongoDB disponível
                        if mongo_connected_blog:
                            try:
                                revisao_data = {
                                    "texto_original": texto_tecnico,
                                    "texto_reescrito": texto_reescrito,
                                    "relatorio_mudancas": relatorio_mudancas if incluir_relatorio else "Não gerado",
                                    "rags_utilizados": rags_ativos,
                                    "documentos_encontrados": {k: len(v) for k, v in resultados_rags.items()},
                                    "nivel_rigor": nivel_rigor,
                                    "incluiu_relatorio": incluir_relatorio,
                                    "data_criacao": datetime.datetime.now()
                                }
                                if 'revisoes_rags' not in db.list_collection_names():
                                    db.create_collection('revisoes_rags')
                                db['revisoes_rags'].insert_one(revisao_data)
                                st.success("✅ Revisão salva no histórico!")
                            except Exception as e:
                                st.warning(f"Revisão concluída, mas não salva: {str(e)}")
                    
                    except Exception as e:
                        st.error(f"❌ Erro no pipeline de RAGs: {str(e)}")
                        with col_revisado_rag:
                            revisado_rag_placeholder.error(f"❌ Erro: {str(e)}")
            else:
                st.warning("Por favor, cole um conteúdo técnico para revisão.")

    # Ferramentas avançadas para análise (mantidas iguais)
    if 'ultima_revisao' in st.session_state and 'ultima_revisao' in locals():
        st.markdown("---")
        st.subheader("🔄 Ajustes Incrementais para RAGs")
        
        st.info("Use o campo abaixo para solicitar ajustes específicos na última revisão com RAGs.")
        
        # Caixa de texto para comandos de ajuste específico para RAGs
        comando_ajuste_rag = st.text_area(
            "Comandos para ajustar a revisão RAG:",
            height=150,
            placeholder="Exemplos:\n- Aumente o foco na taxonomia dos patógenos\n- Inclua mais informações epidemiológicas\n- Corrija dados específicos de produtos\n- Adicione referências da base técnica",
            key="comando_ajuste_rag"
        )
        
        # Botão para ajustar a revisão RAG
        if st.button("🔄 Ajustar Revisão RAG", type="secondary", use_container_width=True):
            if comando_ajuste_rag and 'texto_reescrito' in locals():
                with st.spinner("🔄 Aplicando ajustes na revisão RAG..."):
                    try:
                        # Prompt para ajuste da revisão RAG
                        prompt_ajuste_rag = f"""
                        VOCÊ É: Um especialista técnico agrícola.

                        SUA TAREFA: Ajustar a revisão técnica anterior com base nas solicitações específicas.

                        TEXTO ORIGINAL:
                        {texto_tecnico}

                        TEXTO REESCRITO COM RAGs:
                        {texto_reescrito}

                        RELATÓRIO DE MUDANÇAS:
                        {relatorio_mudancas if 'relatorio_mudancas' in locals() and relatorio_mudancas else "Nenhum relatório disponível"}

                        SOLICITAÇÕES DE AJUSTE:
                        {comando_ajuste_rag}

                        INSTRUÇÕES:
                        1. Aplique TODOS os ajustes solicitados
                        2. Mantenha a precisão técnica
                        3. Considere as informações dos RAGs utilizados
                        4. Retorne o texto reescrito ajustado
                        5. Se solicitado, atualize também o relatório de mudanças

                        Retorne o texto reescrito ajustado.
                        """

                        resposta_ajuste_rag = modelo_texto2.generate_content(prompt_ajuste_rag)
                        texto_reescrito_ajustado = resposta_ajuste_rag.text
                        
                        # Atualizar a visualização
                        with col_revisado_rag:
                            revisado_rag_placeholder.empty()
                            st.success("✅ Revisão RAG ajustada!")
                            st.text_area(
                                "Texto reescrito ajustado:",
                                texto_reescrito_ajustado,
                                height=300,
                                label_visibility="collapsed"
                            )
                        
                        # Botão para baixar versão ajustada
                        st.download_button(
                            "💾 Baixar Versão Ajustada",
                            data=texto_reescrito_ajustado,
                            file_name=f"revisao_rag_ajustada_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                            mime="text/plain",
                            use_container_width=True
                        )
                    
                    except Exception as e:
                        st.error(f"❌ Erro ao ajustar revisão RAG: {str(e)}")

# O resto do código permanece igual...


# --- FUNÇÃO ATUALIZADA PARA BUSCA WEB COM PERPLEXITY ---
def buscar_perplexity(prompt: str) -> str:
    """Realiza busca na web usando a biblioteca Perplexity"""
    try:
        
        # Enviar prompt para o Perplexity
        response = perplexity_client.chat.completions.create(
            model="sonar",
            messages=[
                {"role": "user", "content": prompt}
            ],
            temperature=0.0  # Baixa temperatura para respostas mais precisas
        )
        
        # Pegar a resposta
        resposta = response.choices[0].message.content
        
        # Adicionar informações da resposta
        resposta_completa = f"""{resposta}"""
        
        return resposta_completa
        
    except Exception as e:
        return f"❌ Erro na busca Perplexity: {str(e)}"

# --- FUNÇÃO ESPECÍFICA PARA OTIMIZAÇÃO DE CONTEÚDO ---
def buscar_fontes_para_otimizacao(conteudo: str, tipo: str, tom: str) -> str:
    """Busca fontes específicas para otimização de conteúdo agrícola"""
    
    prompt = f"""
    
   
    DADOS TÉCNICOS ATUALIZADOS para este conteúdo:
    {conteudo[:800]}
    
    
    """
    
    return buscar_perplexity(prompt)
        

# ========== ABA: OTIMIZAÇÃO DE CONTEÚDO ==========
with tab_otimizacao:
    st.header("🚀 Otimização de Conteúdo")
    
    # Inicializar session state
    if 'conteudo_otimizado' not in st.session_state:
        st.session_state.conteudo_otimizado = None
    if 'ultima_otimizacao' not in st.session_state:
        st.session_state.ultima_otimizacao = None
    if 'ajustes_realizados' not in st.session_state:
        st.session_state.ajustes_realizados = []
    if 'fontes_busca_web' not in st.session_state:
        st.session_state.fontes_busca_web = ""
    
    # Área para entrada do conteúdo
    texto_para_otimizar = st.text_area("Cole o conteúdo para otimização:", height=300)
    
    # Configurações
    col_config1, col_config2 = st.columns([2, 1])
    
    with col_config1:
        tipo_otimizacao = st.selectbox("Tipo de Otimização:", 
                                      ["SEO", "Engajamento", "Conversão", "Clareza"])
        
    with col_config2:
        tom_voz = st.text_input("Tom de Voz (ex: Técnico, Persuasivo):", 
                               value="Técnico",
                               key="tom_voz_otimizacao")
        
        nivel_heading = st.selectbox("Nível de Heading Solicitado:", 
                                   ["H1", "H2", "H3", "H4"],
                                   help="Nível de heading que foi solicitado no briefing. CORRIJA se o texto usar nível diferente")

    # CONFIGURAÇÕES DE BUSCA WEB
    st.subheader("🔍 Busca Web e Links")
    
    usar_busca_web = st.checkbox("Usar busca web para enriquecer conteúdo", 
                               value=True,
                               help="Ativa a busca no Perplexity para encontrar informações atualizadas")
    
    incluir_links_internos = st.checkbox("Incluir links internos", 
                                       value=True,
                                       help="Sugere e ancora links relevantes no texto")

    # Área para briefing
    instrucoes_briefing = st.text_area(
        "Instruções do briefing (opcional):",
        height=80
    )

    # --- FUNÇÃO DE BUSCA WEB SEPARADA ---
    def realizar_busca_web_perplexity(texto, tipo_otimizacao, tom_voz):
        """Função separada para realizar busca web"""
        try:
            # Importar dentro da função para evitar erros de importação
            from perplexity import Perplexity
            
            # Obter API key
            perp_api_key = os.getenv("PERP_API_KEY")
            if not perp_api_key:
                return "❌ ERRO: PERP_API_KEY não encontrada nas variáveis de ambiente"
            
            # Inicializar cliente
            client = Perplexity(api_key=perp_api_key)
            
            # Construir prompt para busca
            prompt = f"""
            Você é um assistente especializado em pesquisa. Busque informações atualizadas e confiáveis sobre:
            
            TÓPICO PRINCIPAL: {texto}

        
            
            CRITÉRIOS DE PESQUISA:
            1. Fontes confiáveis: Embrapa, universidades, órgãos governamentais, institutos de pesquisa
            2. Informações técnicas atualizadas (últimos 2-3 anos)
            3. Dados concretos: números, estatísticas, resultados de pesquisa
            4. Melhores práticas agrícolas
            5. Soluções tecnológicas inovadoras
            
            FORMATO DE RESPOSTA:
            Para CADA fonte encontrada, forneça:
            - TÍTULO: Título do artigo/referência
            - CONTEÚDO: Resumo das informações relevantes (máx 200 palavras)
            - URL: Link completo para a fonte
            - RELEVÂNCIA: Por que esta fonte é relevante para o tópico
            
            Retorne no máximo 20 fontes mais relevantes.
            """
            
            # Fazer busca
            response = client.chat.completions.create(
                model="sonar",
                messages=[
                    {"role": "user", "content": prompt}
                ],
                temperature=0.0,
                max_tokens=20000
            )
            
            if response and response.choices:
                resultado = response.choices[0].message.content
                return resultado
            else:
                return "❌ ERRO: Nenhuma resposta recebida do Perplexity"
                
        except ImportError as e:
            return f"❌ ERRO: Biblioteca perplexity-api não instalada. Execute: pip install perplexity-api\nDetalhes: {str(e)}"
        except Exception as e:
            return f"❌ ERRO na busca web: {str(e)}"

    # Botão de otimização
    if st.button("🚀 Otimizar Conteúdo", type="primary", use_container_width=True):
        if texto_para_otimizar:
            with st.spinner("Processando otimização..."):
                try:
                    # FASE 1: BUSCA WEB (se ativada) - AGORA COM TRATAMENTO SEPARADO
                    fontes_encontradas = ""
                    if usar_busca_web:
                        # Container separado para busca web
                        with st.container():
                            st.info("🔍 Iniciando busca web no Perplexity...")
                            
                            # Criar um placeholder para os resultados
                            busca_placeholder = st.empty()
                            
                            # Executar busca web em um bloco try separado
                            try:
                                resultado_busca = realizar_busca_web_perplexity(
                                    texto_para_otimizar, 
                                    tipo_otimizacao, 
                                    tom_voz
                                )
                                
                                # Verificar resultado
                                if resultado_busca and not resultado_busca.startswith("❌"):
                                    fontes_encontradas = resultado_busca
                                    st.session_state.fontes_busca_web = resultado_busca
                                    busca_placeholder.success(f"✅ Busca web concluída: {len(resultado_busca.split())} palavras encontradas")
                                    
                                    # Mostrar preview
                                    with st.expander("📋 Prévia das fontes encontradas", expanded=False):
                                        st.markdown(resultado_busca[:1000] + "..." if len(resultado_busca) > 1000 else resultado_busca)
                                else:
                                    busca_placeholder.warning("⚠️ Busca web não retornou resultados válidos")
                                    st.info("⚠️ Continuando sem fontes externas da busca web")
                                    
                            except Exception as busca_error:
                                busca_placeholder.error(f"❌ Erro na busca web: {str(busca_error)}")
                                st.info("⚠️ Continuando sem fontes externas da busca web")
                    
                    # FASE 2: OTIMIZAÇÃO COM GEMINI
                    st.info("🤖 Iniciando otimização com Gemini...")
                    
                    # Contexto do agente
                    contexto_agente = ""
                    if st.session_state.agente_selecionado:
                        agente = st.session_state.agente_selecionado
                        contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                    
                    # Prompt de otimização
                    prompt = f"""
                    ###BEGIN contexto agente###
                    {contexto_agente}
                    ###END contexto agente###

                    Instruções: Você é um especialista redator técnico. Com base nas informações fornecidas no formato abaixo, gere um artigo completo e bem estruturado sobre o ciclo de desenvolvimento de uma cultura agrícola, seguindo rigorosamente a estrutura, diretrizes e marcação solicitadas.

                    ############BEGIN Formato de Entrada################
                    TÍTULO/H1 desejado: [Título do artigo]
                    Objetivo do conteúdo: [Objetivo descritivo do conteúdo]
                    Público-alvo (persona, nível técnico): [Descrição do público]
                    Palavra-chave principal (KW1): [Palavra-chave primária]
                    Palavras-chave secundárias: [Lista de palavras-chave secundárias, uma por linha]
                    Estrutura (H2/H3 em ordem):
                    [Estrutura completa do artigo com títulos H2 e H3]
                    Região/bioma/safra alvo: [Cultura e contexto]
                    CTA FINAL OBRIGATÓRIA:
                    [Texto do call-to-action]
                    link da CTA: [URL]
                    Interlinks prioritários (URLs internas existentes): [Lista ou "não aplicável"]
                    Links externos obrigatórios (se houver): [Lista ou "não aplicável"]
                    Diretrizes de tom/estilo (brand voice): [Ex.: técnico e leve]
                    Observações/restrições: [Informações adicionais]
                    Frases e parágrafos devem ser mais curtos
                    ############END Formato de Entrada################

                    
                    Sua tarefa: Ao receber uma entrada no formato acima, você deve gerar um documento de artigo completo que inclua:
                    
                        Metadados SEO:
                    
                            Meta title: Crie um com até 60 caracteres, incluindo a KW1.
                    
                            Meta description: Crie uma descrição persuasiva com até 160 caracteres, incluindo a KW1 e uma chamada para ação.
                    
                            URL: Sugira uma URL amigável para SEO baseada no título.
                    
                            Categoria: Sugira uma categoria temática.
                    
                            Imagem de capa: Sugira um tema genérico para imagem (ex.: "Lavouras de [cultura] em campo aberto") e um Alt text descritivo.
                    
                        Corpo do Artigo:
                    
                            Inicie com o TÍTULO/H1 fornecido.
                    
                            Escreva uma introdução envolvente que contextualize a importância da cultura e do manejo correto do seu ciclo.
                    
                            Desenvolva o conteúdo seguindo exatamente a ordem e a hierarquia (H2, H3) fornecidas na "Estrutura".
                    
                            Para cada H3 (que representa um estágio fenológico), estruture o texto com os seguintes subtópicos, sem usar marcadores na explicação:
                    
                                O que é: Definição clara do estágio.
                    
                                Características: Descrições morfológicas e fisiológicas principais.
                    
                                Práticas de Manejo: Recomendações técnicas específicas para essa fase (nutrição, irrigação, controle fitossanitário).
                    
                                Pontos Críticos e Cuidados: Principais riscos (estresses, pragas, doenças) e como mitigá-los.
                    
                            Incorpore naturalmente a KW principal e as palavras-chave secundárias ao longo do texto.
                    
                            Use um tom que equilibre precisão técnica e clareza, conforme as diretrizes de "brand voice".
                    
                            Onde a estrutura sugerir (ex.: após seções longas), insira uma caixa "Leia mais:" ou "Leia também:" com 2-3 sugestões de artigos relacionados baseadas no tema geral. Invente títulos plausíveis para estes interlinks.
                    
                            Finalize com uma conclusão que resuma a importância do manejo faseado.
                    
                            Inclua obrigatoriamente o CTA FINAL com o texto e link fornecidos.
                    
                        Elementos Adicionais (se aplicável na estrutura):
                    
                            Se a estrutura incluir "Tabela", crie uma tabela em markdown resumindo os estágios, características, práticas e pontos críticos.
                    
                            Se a estrutura incluir uma seção sobre "Quanto tempo dura o ciclo...", explique a variação de duração com base em cultivares, clima e região.
                    
                    Regras Gerais:
                    
                        Fidelidade: Siga a estrutura fornecida à risca. Não altere a ordem dos H2/H3.
                    
                        Objetividade: Forneça informações práticas e acionáveis. Evite linguagem excessivamente promocional no corpo do texto.
                    
                        Completude: Certifique-se de que todos os elementos da entrada foram atendidos (KWs, estrutura, CTA).

                        Frases e parágrafos devem ser mais curtos
                    
                        Formatação: Use negrito para termos técnicos importantes ou frases de impacto ocasionais. Use marcadores apenas em listas de itens muito concisos (ex.: características de um estágio). Prefira parágrafos fluidos.
                    
                    Exemplo de Saída (Estrutura Visual):
                    text
                    
                    Meta title: [Texto]
                    Meta description: [Texto]
                    URL: /url-sugerida
                    Categoria: [Categoria Sugerida]
                    Imagem de capa: [Tema sugerido]
                    Alt text: [Descrição da imagem]
                    
                    # TÍTULO/H1 FORNECIDO
                    
                    [Parágrafo de introdução]
                    
                    ## H2 FORNECIDO
                    [Texto explicativo da seção]
                    
                    ### H3 FORNECIDO
                    **O que é:** [Definição].
                    **Características:** [Descrição].
                    **Práticas de Manejo:** [Recomendações].
                    **Pontos Críticos e Cuidados:** [Riscos e soluções].
                    
                    [Continue para todos os H3s e H2s...]
                    
                    **Leia mais:**
                    *   Título de artigo relacionado 1
                    *   Título de artigo relacionado 2
                    
                    ## H2 FINAL (ex.: Conclusão)
                    [Texto de conclusão]
                    
                    [CTA FINAL OBRIGATÓRIO com link]

                    [Links que foram ancorados por extenso]



                    **TEXTO ORIGINAL:**
                    {texto_para_otimizar}

                    **FONTES DA BUSCA WEB (para serem usadas de forma ancorada ao longo do texto quando relevantes)**
                    {fontes_encontradas if fontes_encontradas else "Nenhuma fonte externa disponível."}

                    **INSTRUÇÕES DO BRIEFING:**
                    {instrucoes_briefing if instrucoes_briefing else 'Sem briefing específico'}

                    **CONFIGURAÇÕES:**
                    - Tipo: {tipo_otimizacao}
                    - Tom: {tom_voz}
                    - Heading level: {nivel_heading}
                    - Links internos: {"Sim" if incluir_links_internos else "Não"}
                    - Busca web usada: {"Sim" if fontes_encontradas else "Não"}
                    - Frases e parágrafos devem ser mais curtos

                    ## REQUISITOS OBRIGATÓRIOS:

                    VOCÊ É: Um especialista em SEO técnico agrícola.

                    TAREFA: Otimizar o conteúdo para SEO, focando especialmente na estrutura de headings.
                    
                    HEADING SOLICITADO NO BRIEFING: {nivel_heading}
                    
                    SIGA ESTE PROCESSO EM 4 ETAPAS:
                    
                    ETAPA 1 - ANÁLISE ESTRUTURAL:
                    1. Identificar TODOS os headings no texto
                    2. Verificar se estão no nível {nivel_heading}
                    3. Se não estiverem, corrigir para {nivel_heading}
                    4. Manter a hierarquia lógica
                    
                    ETAPA 2 - OTIMIZAÇÃO SEO:
                    1. Garantir que o primeiro heading seja {nivel_heading} com palavra-chave
                    2. Incluir palavra-chave em 50% dos headings
                    3. Headings devem ser descritivos e incluir benefícios
                    4. Máximo de 300 palavras entre headings
                    
                    ETAPA 3 - CONTEÚDO:
                    1. Bullets para listas (máx 5 itens)
                    2. Parágrafos curtos (3-4 frases)
                    3. CTAs claros
                    4. Dados concretos quando possível
                    
                    ETAPA 4 - META TAGS:
                    Gerar 3 opções de title/description otimizadas.
                    
                    FORMATO DE RESPOSTA:
                    [Conteúdo otimizado com headings corrigidos]
                    
                    Heading corrections applied:
                    - "Original heading" → "Corrected {nivel_heading}

                    1. **TITLES E DESCRIPTIONS (OBRIGATÓRIO):**
                       Gere 3 opções de meta title (≤60 chars) e description (≤155 chars)
                       Exemplo:
                       Title: Guia Prático de Adubação Nitrogenada no Milho - Aumente sua Produtividade
                       Description: Descubra como a adubação nitrogenada adequada pode aumentar em até 30% a produtividade do milho. Técnicas comprovadas!

                    2. **BULLETS QUANDO APLICÁVEL:**
                       - Use bullets para listas de benefícios
                       - Use bullets para características técnicas
                       - Use bullets para etapas de processo
                       - Máximo 5 itens por lista

                    3. **HEADING LEVEL {nivel_heading}:**
                       - Todos os headings principais devem ser {nivel_heading}
                       - Corrigir se estiver usando nível diferente
                       - Manter hierarquia consistente

                    4. **CORREÇÕES AUTOMÁTICAS:**
                       - Remova introduções genéricas - Você é um profissional experiente
                       - Quebre parágrafos longos (3-4 frases máx)
                       - Remova repetições
                       - Melhore escaneabilidade
                       - Divida frases complexas
                       - Incorpore dados das fontes quando relevante
                       - Frases e parágrafos devem ser mais curtos

                    5. **LINKS INTERNOS:**
                       Sugira 3-5 links relevantes no formato: [texto âncora](url)
                       Escreva os links que foram ancorados por extenso ao final
                    """

                    # Gerar otimização
                    resposta = modelo_texto.generate_content(prompt)
                    resultado = resposta.text
                    
                    # Processar resultado
                    partes_do_resultado = {
                        "📝 CONTEÚDO OTIMIZADO": resultado  # Default
                    }
                    
                    # Tentar extrair seções
                    secoes = ["📊 SUGESTÕES DE META TAGS", "✅ CORREÇÕES APLICADAS", "🔗 LINKS INTERNOS SUGERIDOS", "📝 CONTEÚDO OTIMIZADO"]
                    
                    for i in range(len(secoes)):
                        if secoes[i] in resultado:
                            inicio = resultado.find(secoes[i])
                            if i < len(secoes) - 1 and secoes[i+1] in resultado:
                                fim = resultado.find(secoes[i+1])
                                conteudo = resultado[inicio + len(secoes[i]):fim].strip()
                            else:
                                conteudo = resultado[inicio + len(secoes[i]):].strip()
                            
                            # Limpar formatação extra
                            conteudo = conteudo.strip(":#*-\n ")
                            partes_do_resultado[secoes[i]] = conteudo
                    
                    # Salvar no session state
                    st.session_state.conteudo_otimizado = partes_do_resultado.get("📝 CONTEÚDO OTIMIZADO", resultado)
                    st.session_state.ultima_otimizacao = resultado
                    st.session_state.texto_original = texto_para_otimizar
                    st.session_state.fontes_busca_web = fontes_encontradas
                    st.session_state.partes_resultado = partes_do_resultado
                    
                    # Exibir resultados
                    st.success("✅ Conteúdo otimizado com sucesso!")
                    
                    # 1. Meta Tags
                    st.subheader("📊 Meta Tags Geradas")
                    if "📊 SUGESTÕES DE META TAGS" in partes_do_resultado:
                        st.markdown(partes_do_resultado["📊 SUGESTÕES DE META TAGS"])
                    else:
                        # Procurar meta tags no texto
                        lines = resultado.split('\n')
                        meta_candidates = []
                        for line in lines:
                            line_lower = line.lower()
                            if ('title:' in line_lower or 'description:' in line_lower or 
                                'meta ' in line_lower or 'tag' in line_lower):
                                meta_candidates.append(line)
                        
                        if meta_candidates:
                            st.info("Meta tags encontradas:")
                            for line in meta_candidates[:6]:
                                st.write(line)
                        else:
                            st.warning("Meta tags não foram detectadas automaticamente")
                    
                    # 2. Correções
                    if "✅ CORREÇÕES APLICADAS" in partes_do_resultado:
                        with st.expander("✅ Correções Aplicadas", expanded=True):
                            st.markdown(partes_do_resultado["✅ CORREÇÕES APLICADAS"])
                    
                    # 3. Links Internos
                    if "🔗 LINKS INTERNOS SUGERIDOS" in partes_do_resultado and incluir_links_internos:
                        with st.expander("🔗 Links Sugeridos"):
                            st.markdown(partes_do_resultado["🔗 LINKS INTERNOS SUGERIDOS"])
                    
                    # 4. Conteúdo Otimizado
                    st.subheader("📝 Conteúdo Otimizado")
                    conteudo_final = partes_do_resultado.get("📝 CONTEÚDO OTIMIZADO", resultado)
                    st.markdown(conteudo_final)
                    
                    # Verificações
                    st.subheader("🔍 Verificação")
                    
                    col1, col2, col3 = st.columns(3)
                    with col1:
                        bullets = conteudo_final.count("- ") + conteudo_final.count("* ")
                        st.metric("Bullet Points", bullets)
                    with col2:
                        has_heading = nivel_heading.lower() in conteudo_final.lower()
                        st.metric(f"Heading {nivel_heading}", "✅" if has_heading else "❌")
                    with col3:
                        has_meta = 'title' in conteudo_final[:500].lower() or 'description' in conteudo_final[:500].lower()
                        st.metric("Meta Tags", "✅" if has_meta else "❌")
                    
                    # Download
                    st.download_button(
                        "💾 Baixar Conteúdo Otimizado",
                        data=conteudo_final,
                        file_name=f"conteudo_otimizado_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                        mime="text/plain"
                    )
                    
                except Exception as e:
                    st.error(f"❌ Erro na otimização: {str(e)}")
                    st.info("Dica: Verifique sua conexão com a API do Gemini")
        else:
            st.warning("Por favor, cole um conteúdo para otimizar")

    # Ajustes incrementais
    if st.session_state.conteudo_otimizado:
        st.divider()
        st.subheader("🔄 Ajustes Incrementais")
        
        comando_ajuste = st.text_area(
            "Ajustes desejados:",
            height=80,
            placeholder="Ex: Adicione mais bullets, corrija headings, melhore meta tags...",
            key="ajuste_text"
        )
        
        if st.button("🔄 Aplicar Ajustes", key="btn_ajuste"):
            if comando_ajuste:
                with st.spinner("Aplicando ajustes..."):
                    try:
                        prompt_ajuste = f"""
                        **CONTEÚDO ATUAL:** {st.session_state.conteudo_otimizado[:1000]}
                        
                        **AJUSTES SOLICITADOS:** {comando_ajuste}
                        
                        **MANTENHA:** 
                        - Meta tags existentes
                        - Heading level {nivel_heading}
                        - Bullets onde aplicável
                        - Frases e parágrafos devem ser mais curtos

                        VOCÊ É: Um especialista em SEO técnico agrícola.

                        TAREFA: Otimizar o conteúdo para SEO, focando especialmente na estrutura de headings.
                        
                        HEADING SOLICITADO NO BRIEFING: {nivel_heading}
                        
                        SIGA ESTE PROCESSO EM 4 ETAPAS:
                        
                        ETAPA 1 - ANÁLISE ESTRUTURAL:
                        1. Identificar TODOS os headings no texto
                        2. Verificar se estão no nível {nivel_heading}
                        3. Se não estiverem, corrigir para {nivel_heading}
                        4. Manter a hierarquia lógica
                        
                        ETAPA 2 - OTIMIZAÇÃO SEO:
                        1. Garantir que o primeiro heading seja {nivel_heading} com palavra-chave
                        2. Incluir palavra-chave em 50% dos headings
                        3. Headings devem ser descritivos e incluir benefícios
                        4. Máximo de 300 palavras entre headings
                        
                        ETAPA 3 - CONTEÚDO:
                        1. Bullets para listas (máx 5 itens)
                        2. Parágrafos curtos (3-4 frases)
                        3. CTAs claros
                        4. Dados concretos quando possível
                        
                        ETAPA 4 - META TAGS:
                        Gerar 3 opções de title/description otimizadas.
                        
                        FORMATO DE RESPOSTA:
                        [Conteúdo otimizado com headings corrigidos]
                        
                        Heading corrections applied:
                        - "Original heading" → "Corrected {nivel_heading}
                        
                        Aplique os ajustes e retorne APENAS o conteúdo atualizado.
                        """
                        
                        resposta = modelo_texto.generate_content(prompt_ajuste)
                        st.session_state.conteudo_otimizado = resposta.text
                        st.session_state.ajustes_realizados.append(comando_ajuste)
                        
                        st.success("✅ Ajustes aplicados!")
                        st.markdown(resposta.text)
                        
                    except Exception as e:
                        st.error(f"Erro: {str(e)}")
            else:
                st.warning("Digite os ajustes desejados")
        
        # Limpar histórico
        if st.button("🗑️ Limpar Histórico de Ajustes"):
            st.session_state.ajustes_realizados = []
            st.success("Histórico limpo")
            
# ========== ABA: CRIADORA DE CALENDÁRIO ==========
with tab_calendario:
    st.header("📅 Criadora de Calendário")
    
    if not st.session_state.agente_selecionado:
        st.warning("Nenhum agente selecionado.")
    else:
        agente = st.session_state.agente_selecionado
        st.success(f"Agente: {agente['nome']}")
        
        col1, col2 = st.columns([2, 1])
        
        with col1:
            mes_ano = st.text_input("Mês/Ano:", "FEVEREIRO 2026")
            data_inicio = st.date_input("Data início:", value=datetime.date(2026, 2, 1))
            data_fim = st.date_input("Data fim:", value=datetime.date(2026, 2, 28))
            
            delta_dias = (data_fim - data_inicio).days + 1
            
            culturas_prioritarias = st.text_area(
                "Culturas (separadas por vírgula, use 'e' para múltiplas):",
                "Soja, Milho, Cana-de-açúcar, Algodão, Soja e Milho, Soja e Cana"
            )
            culturas_lista = [c.strip() for c in culturas_prioritarias.split(",") if c.strip()]
        
        with col2:
            dias_com_1_pauta = st.number_input("Dias com 1 pauta:", 0, delta_dias, 5)
            dias_com_2_pautas = st.number_input("Dias com 2 pautas:", 0, delta_dias, 15)
            dias_com_3_pautas = st.number_input("Dias com 3 pautas:", 0, delta_dias, 3)
            dias_sem_pautas = delta_dias - (dias_com_1_pauta + dias_com_2_pautas + dias_com_3_pautas)
            
            if dias_sem_pautas < 0:
                st.error("Total excede dias disponíveis")
        
        st.subheader("Produtos e Direcionais")
        st.write("Formato: Produto(s) - Cultura(s) - Tema")
        st.write("Ex: Elestal Neo e Fortenza - Soja e Milho - Controle de pragas")
        
        produtos_direcionais = st.text_area(
            "Produtos com culturas e temas:",
            """Verdavis, Megafol e Victrato - Soja e Milho - Tecnologia para feira
Elestal Neo - Soja - Controle de mosca-branca
Fortenza - Milho - Seedcare para cigarrinha
YieldOn - Soja - Bioativador para pegamento
Miravis - Soja - Fungicida para ferrugem
Victrato - Cana - Nematicida para cana-soca
Victrato pelo Brasil - Soja e Cana - Ação nacional""",
            height=150
        )
        
        produtos_com_direcionais = []
        if produtos_direcionais:
            for linha in produtos_direcionais.split('\n'):
                linha = linha.strip()
                if linha and ' - ' in linha:
                    partes = linha.split(' - ')
                    if len(partes) >= 3:
                        produtos = [p.strip() for p in partes[0].split(' e ') if p.strip()]
                        culturas = [c.strip() for c in partes[1].split(' e ') if c.strip()]
                        tema = ' - '.join(partes[2:]).strip()
                        produtos_com_direcionais.append({
                            'produtos': produtos,
                            'culturas': culturas,
                            'tema': tema
                        })
        
        col_feira, col_recorrente = st.columns(2)
        
        with col_feira:
            st.write("Semana com evento (1 post/dia):")
            semana_feira_inicio = st.date_input("Início:", value=datetime.date(2026, 2, 9))
            semana_feira_fim = st.date_input("Fim:", value=datetime.date(2026, 2, 13))
            produtos_prioritarios_feira = st.text_input("Produtos prioritários:", "Verdavis, Megafol, Victrato")
        
        with col_recorrente:
            pauta_recorrente_texto = st.text_input("Pauta fixa:", "Victrato pelo Brasil")
            pauta_recorrente_dias = st.multiselect(
                "Dias da semana:",
                ["Terça", "Quinta"],
                default=["Terça", "Quinta"]
            )
        
        contexto_mensal = st.text_area(
            "Contexto do mês:",
            """FEVEREIRO 2026:
- Soja: colheita no centro-sul
- Milho: plantio da safrinha
- Cana: crescimento vegetativo
- Evento: Feira Nacional do Agronegócio (09-13/02)
- Foco: Verdavis, Megafol, Victrato na feira
- Pauta fixa: Victrato pelo Brasil (terças e quintas)""",
            height=120
        )
        
        evitar_consecutivos_sem_pautas = st.checkbox("Evitar dias consecutivos sem pautas", True)
        max_repeticoes_tema = st.slider("Máx repetições por tema:", 1, 5, 2)
        
        if st.button("Gerar Calendário", type="primary"):
            if data_inicio >= data_fim:
                st.error("Data início deve ser anterior")
            elif not culturas_lista:
                st.error("Digite culturas")
            elif (dias_com_1_pauta + dias_com_2_pautas + dias_com_3_pautas) > delta_dias:
                st.error("Total excede período")
            else:
                with st.spinner("Gerando calendário..."):
                    try:
                        contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                        
                        info_especifica = f"""
                        CONFIGURAÇÕES:
                        1. SEMANA COM EVENTO ({semana_feira_inicio.strftime('%d/%m')} a {semana_feira_fim.strftime('%d/%m')}):
                           - Apenas 1 pauta por dia
                           - Priorizar: {produtos_prioritarios_feira}
                        
                        2. PAUTA FIXA: "{pauta_recorrente_texto}"
                           - Dias: {', '.join(pauta_recorrente_dias)}
                        
                        3. FREQUÊNCIA:
                           - Dias com 1 pauta: {dias_com_1_pauta}
                           - Dias com 2 pautas: {dias_com_2_pautas} 
                           - Dias com 3 pautas: {dias_com_3_pautas}
                           - Dias sem pautas: {max(0, dias_sem_pautas)}
                           - Evitar consecutivos sem pautas: {evitar_consecutivos_sem_pautas}
                        
                        4. CONTROLE REPETIÇÃO:
                           - Máximo repetições por tema: {max_repeticoes_tema}
                           - Células podem ter múltiplas culturas/produtos
                        """
                        

                        info_algodao = """
                        Tocantins: Plantio de novembro (2ª quinzena) até fevereiro (2ª quinzena), com pico intenso em janeiro. Colheita de abril (2ª quinzena) até agosto (1ª quinzena), com pico intenso em junho e julho.
                        Maranhão: Plantio de dezembro (1ª quinzena) até março (2ª quinzena), com pico intenso em janeiro. Colheita de maio (2ª quinzena) até agosto (2ª quinzena), com pico intenso em junho e julho.
                        Piauí: Plantio de dezembro (2ª quinzena) até março (2ª quinzena), com pico intenso em janeiro. Colheita de maio (2ª quinzena) até agosto (1ª quinzena), com pico intenso em junho e julho.
                        Ceará: Plantio de janeiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em fevereiro e março. Colheita de junho (1ª quinzena) até outubro (2ª quinzena), com pico intenso em junho, julho e agosto.
                        Rio Grande do Norte: Plantio de janeiro (1ª quinzena) até abril (2ª quinzena), com pico intenso em fevereiro e março. Colheita de julho (1ª quinzena) até novembro (2ª quinzena), com pico intenso em agosto e setembro.
                        Paraíba: Plantio de fevereiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março. Colheita de agosto (1ª quinzena) até novembro (2ª quinzena), com pico intenso em agosto e setembro.
                        Pernambuco: Plantio de janeiro (1ª quinzena) até junho (2ª quinzena), com pico intenso em março. Colheita de agosto (1ª quinzena) até dezembro (1ª quinzena), com pico intenso em agosto e setembro.
                        Alagoas: Plantio de maio (2ª quinzena) até agosto (2ª quinzena), com pico intenso em junho. Colheita de outubro (2ª quinzena) até janeiro (2ª quinzena), com pico intenso em novembro e dezembro.
                        Bahia: Plantio de novembro (2ª quinzena) até fevereiro (1ª quinzena), com pico intenso em dezembro. Colheita de abril (2ª quinzena) até setembro (1ª quinzena), com pico intenso em maio e junho.
                        Mato Grosso: Plantio de dezembro (1ª quinzena) até fevereiro (2ª quinzena), com pico intenso em janeiro. Colheita de abril (2ª quinzena) até agosto (2ª quinzena), com pico intenso em junho.
                        Mato Grosso do Sul: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de março (2ª quinzena) até junho (1ª quinzena), com pico intenso em abril.
                        Goiás: Plantio de outubro (2ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de março (2ª quinzena) até junho (2ª quinzena), com pico intenso em maio.
                        Distrito Federal: Plantio de outubro (2ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de abril (1ª quinzena) até junho (2ª quinzena), com pico intenso em maio.
                        Minas Gerais: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de março (2ª quinzena) até junho (1ª quinzena), com pico intenso em abril e maio.
                        São Paulo: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de março (1ª quinzena) até junho (1ª quinzena), com pico intenso em abril e maio.
                        Paraná: Plantio de setembro (2ª quinzena) até dezembro (1ª quinzena), com pico intenso em outubro e novembro. Colheita de março (1ª quinzena) até maio (2ª quinzena), com pico intenso em abril.
                        """
                        
                        info_arroz = """
                        Roraima: Plantio de maio (1ª quinzena) até agosto (2ª quinzena), com pico intenso em maio. Colheita de julho (2ª quinzena) até novembro (2ª quinzena), com pico intenso em setembro.
                        Rondônia: Plantio de setembro (2ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de janeiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em fevereiro e março.
                        Acre: Plantio de setembro (2ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de janeiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em fevereiro e março.
                        Amazonas: Plantio de setembro (2ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de janeiro (2ª quinzena) até maio (1ª quinzena), com pico intenso em março.
                        Amapá: Plantio de janeiro (1ª quinzena) até abril (1ª quinzena), com pico intenso em fevereiro. Colheita de maio (2ª quinzena) até agosto (2ª quinzena), com pico intenso em junho.
                        Pará: Plantio de dezembro (1ª quinzena) até abril (2ª quinzena), com pico intenso em janeiro. Colheita de abril (1ª quinzena) até agosto (2ª quinzena), com pico intenso em abril e maio.
                        Tocantins: Plantio de outubro (2ª quinzena) até janeiro (2ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março e abril.
                        Maranhão: Plantio de novembro (1ª quinzena) até março (1ª quinzena), com pico intenso em janeiro. Colheita de abril (1ª quinzena) até julho (1ª quinzena), com pico intenso em abril e maio.
                        Piauí: Plantio de novembro (1ª quinzena) até março (1ª quinzena), com pico intenso em janeiro. Colheita de abril (1ª quinzena) até julho (1ª quinzena), com pico intenso em abril e maio.
                        Ceará: Plantio de janeiro (1ª quinzena) até abril (2ª quinzena), com pico intenso em janeiro e fevereiro. Colheita de maio (1ª quinzena) até julho (1ª quinzena), com pico intenso em maio e junho.
                        Rio Grande do Norte: Plantio de janeiro (2ª quinzena) até maio (1ª quinzena), com pico intenso em março. Colheita de junho (1ª quinzena) até outubro (1ª quinzena), com pico intenso em agosto.
                        Paraíba: Plantio de janeiro (1ª quinzena) até abril (1ª quinzena), com pico intenso em janeiro e fevereiro. Colheita de maio (2ª quinzena) até agosto (1ª quinzena), com pico intenso em junho.
                        Pernambuco: Plantio de janeiro (1ª quinzena) até abril (1ª quinzena), com pico intenso em fevereiro. Colheita de maio (1ª quinzena) até agosto (1ª quinzena), com pico intenso em junho.
                        Alagoas: Plantio de setembro (2ª quinzena) até dezembro (1ª quinzena), com pico intenso em outubro e novembro. Colheita de janeiro (1ª quinzena) até março (2ª quinzena), com pico intenso em fevereiro e março.
                        Sergipe: Plantio de setembro (2ª quinzena) até novembro (2ª quinzena), com pico intenso em outubro. Colheita de janeiro (1ª quinzena) até março (2ª quinzena), com pico intenso em fevereiro.
                        Bahia: Plantio de setembro (2ª quinzena) até dezembro (1ª quinzena), com pico intenso em outubro e novembro. Colheita de janeiro (1ª quinzena) até abril (1ª quinzena), com pico intenso em fevereiro e março.
                        Mato Grosso: Plantio de setembro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de janeiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março.
                        Mato Grosso do Sul: Plantio de setembro (1ª quinzena) até dezembro (2ª quinzena), com pico intenso em outubro e novembro. Colheita de janeiro (1ª quinzena) até abril (2ª quinzena), com pico intenso em fevereiro.
                        Goiás: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (1ª quinzena) até abril (2ª quinzena), com pico intenso em março.
                        Distrito Federal: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (1ª quinzena) até abril (2ª quinzena), com pico intenso em março.
                        Minas Gerais: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março e abril.
                        Espírito Santo: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março e abril.
                        Rio de Janeiro: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de janeiro (1ª quinzena) até maio (2ª quinzena), com pico intenso em março e abril.
                        São Paulo: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março e abril.
                        Paraná: Plantio de setembro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de janeiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março e abril.
                        Santa Catarina: Plantio de agosto (2ª quinzena) até dezembro (1ª quinzena), com pico intenso em outubro e novembro. Colheita de janeiro (1ª quinzena) até abril (2ª quinzena), com pico intenso em fevereiro e março.
                        Rio Grande do Sul: Plantio de setembro (1ª quinzena) até dezembro (1ª quinzena), com pico intenso em outubro e novembro. Colheita de fevereiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março e abril.
"""
                        info_soja = """
                        Roraima: Plantio de abril (2ª quinzena) até junho (2ª quinzena), com pico intenso em maio. Colheita de julho (2ª quinzena) até novembro (2ª quinzena), com pico intenso em setembro.
                        Rondônia: Plantio de setembro (2ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro. Colheita de janeiro (1ª quinzena) até abril (1ª quinzena), com pico intenso em março.
                        Amazonas: Plantio de setembro (1ª quinzena) até dezembro (1ª quinzena), com pico intenso em setembro e outubro. Colheita de dezembro (2ª quinzena) até março (2ª quinzena), com pico intenso em fevereiro.
                        Pará: Plantio de outubro (1ª quinzena) até janeiro (2ª quinzena), com pico intenso em março. Colheita de fevereiro (2ª quinzena) até agosto (2ª quinzena), com pico intenso em março e julho.
                        Tocantins: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março e abril.
                        Maranhão: Plantio de outubro (1ª quinzena) até janeiro (2ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (2ª quinzena) até maio (2ª quinzena), com pico intenso em março e abril.
                        Piauí: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de março (1ª quinzena) até maio (2ª quinzena), com pico intenso em abril.
                        Bahia: Plantio de outubro (1ª quinzena) até janeiro (2ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (2ª quinzena) até maio (1ª quinzena), com pico intenso em março e abril.
                        Mato Grosso: Plantio de setembro (1ª quinzena) até dezembro (2ª quinzena), com pico intenso em outubro e novembro. Colheita de janeiro (1ª quinzena) até abril (1ª quinzena), com pico intenso em fevereiro e março.
                        Mato Grosso do Sul: Plantio de setembro (2ª quinzena) até dezembro (2ª quinzena), com pico intenso em outubro e novembro. Colheita de janeiro (2ª quinzena) até abril (1ª quinzena), com pico intenso em março.
                        Goiás: Plantio de setembro (2ª quinzena) até janeiro (1ª quinzena), com pico intenso em outubro e novembro. Colheita de janeiro (2ª quinzena) até abril (1ª quinzena), com pico intenso em março.
                        Distrito Federal: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (1ª quinzena) até maio (2ª quinzena), com pico intenso em março e abril.
                        Minas Gerais: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março e abril.
                        São Paulo: Plantio de setembro (2ª quinzena) até janeiro (1ª quinzena), com pico intenso em outubro e novembro. Colheita de janeiro (1ª quinzena) até abril (2ª quinzena), com pico intenso em março.
                        Paraná: Plantio de setembro (2ª quinzena) até dezembro (2ª quinzena), com pico intenso em outubro e novembro. Colheita de janeiro (2ª quinzena) até abril (2ª quinzena), com pico intenso em março.
                        Santa Catarina: Plantio de outubro (1ª quinzena) até dezembro (2ª quinzena), com pico intenso em novembro e dezembro. Colheita de janeiro (1ª quinzena) até maio (2ª quinzena), com pico intenso em março e abril.
                        Rio Grande do Sul: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de fevereiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março e abril.
"""
                        info_milho = """
                        Calendário de Safra: Milho 1ª Safra (Ciclo 120-180 dias)
                        Rondônia: Plantio de agosto (2ª quinzena) até novembro (1ª quinzena), com pico intenso em setembro. Colheita de janeiro (2ª quinzena) até abril (2ª quinzena), com pico intenso em fevereiro.
                        Acre: Plantio de setembro (2ª quinzena) até dezembro (1ª quinzena), com pico intenso em outubro. Colheita de fevereiro (1ª quinzena) até maio (2ª quinzena), com pico intenso em março.
                        Amazonas: Plantio de outubro (1ª quinzena) até dezembro (2ª quinzena), com pico intenso em novembro. Colheita de março (2ª quinzena) até junho (2ª quinzena), com pico intenso em abril.
                        Pará: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro. Colheita de março (2ª quinzena) até junho (2ª quinzena), com pico intenso em maio.
                        Tocantins: Plantio de outubro (1ª quinzena) até janeiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de março (2ª quinzena) até junho (2ª quinzena), com pico intenso em maio.
                        Maranhão: Plantio de outubro (1ª quinzena) até janeiro (2ª quinzena), com pico intenso em novembro. Colheita de março (2ª quinzena) até junho (2ª quinzena), com pico intenso em abril.
                        Piauí: Plantio de outubro (1ª quinzena) até janeiro (2ª quinzena), com pico intenso em novembro e dezembro. Colheita de abril (1ª quinzena) até junho (2ª quinzena), com pico intenso em maio.
                        Pernambuco: Plantio de outubro (2ª quinzena) até janeiro (2ª quinzena), com pico intenso em dezembro. Colheita de abril (2ª quinzena) até junho (2ª quinzena), com pico intenso em maio.
                        Bahia: Plantio de outubro (1ª quinzena) até fevereiro (1ª quinzena), com pico intenso em novembro e dezembro. Colheita de março (2ª quinzena) até julho (1ª quinzena), com pico intenso em abril e maio.
                        Mato Grosso: Plantio de setembro (2ª quinzena) até dezembro (1ª quinzena), com pico intenso em outubro e novembro. Colheita de fevereiro (1ª quinzena) até maio (1ª quinzena), com pico intenso em março e abril.
                        Mato Grosso do Sul: Plantio de agosto (2ª quinzena) até novembro (2ª quinzena), com pico intenso em setembro e outubro. Colheita de janeiro (2ª quinzena) até abril (1ª quinzena), com pico intenso em março.
                        Goiás: Plantio de setembro (2ª quinzena) até janeiro (1ª quinzena), com pico intenso em outubro e novembro. Colheita de janeiro (2ª quinzena) até junho (1ª quinzena), com pico intenso em março e abril.
                        Distrito Federal: Plantio de setembro (2ª quinzena) até dezembro (2ª quinzena), com pico intenso em outubro e novembro. Colheita de fevereiro (2ª quinzena) até junho (1ª quinzena), com pico intenso em abril.
                        Minas Gerais: Plantio de setembro (2ª quinzena) até janeiro (1ª quinzena), com pico intenso em outubro, novembro e dezembro. Colheita de fevereiro (2ª quinzena) até junho (1ª quinzena), com pico intenso em maio.
                        Espírito Santo: Plantio de agosto (2ª quinzena) até dezembro (2ª quinzena), com pico intenso em setembro e outubro. Colheita de janeiro (2ª quinzena) até maio (2ª quinzena), com pico intenso em março.
                        Rio de Janeiro: Plantio de setembro (1ª quinzena) até dezembro (2ª quinzena), com pico intenso em outubro e novembro. Colheita de fevereiro (1ª quinzena) até junho (1ª quinzena), com pico intenso em março e abril.
                        São Paulo: Plantio de setembro (1ª quinzena) até dezembro (2ª quinzena), com pico intenso em outubro e novembro. Colheita de janeiro (1ª quinzena) até julho (1ª quinzena), com pico intenso em março e abril.
                        Paraná: Plantio de agosto (2ª quinzena) até dezembro (1ª quinzena), com pico intenso em setembro e outubro. Colheita de janeiro (2ª quinzena) até junho (2ª quinzena), com pico intenso em março.
                        Santa Catarina: Plantio de agosto (1ª quinzena) até dezembro (1ª quinzena), com pico intenso em setembro e outubro. Colheita de janeiro (1ª quinzena) até maio (2ª quinzena), com pico intenso em março e abril.
                        Rio Grande do Sul: Plantio de agosto (1ª quinzena) até novembro (2ª quinzena), com pico intenso em setembro e outubro. Colheita de dezembro (2ª quinzena) até maio (2ª quinzena), com pico intenso em fevereiro e março.
                        Calendário de Safra: Milho 2ª Safra (Ciclo 120-180 dias)
                        Roraima: Plantio de maio (1ª quinzena) até junho (2ª quinzena), com pico intenso em maio. Colheita de setembro (2ª quinzena) até novembro (2ª quinzena), com pico intenso em outubro.
                        Rondônia: Plantio de janeiro (2ª quinzena) até março (1ª quinzena), com pico intenso em fevereiro. Colheita de maio (2ª quinzena) até agosto (2ª quinzena), com pico intenso em julho e agosto.
                        Amapá: Plantio de fevereiro (1ª quinzena) até março (2ª quinzena), com pico intenso em fevereiro. Colheita de maio (2ª quinzena) até julho (1ª quinzena), com pico intenso em maio e junho.
                        Pará: Plantio de janeiro (1ª quinzena) até março (1ª quinzena), com pico intenso em janeiro e fevereiro. Colheita de abril (2ª quinzena) até novembro (2ª quinzena), com pico intenso em maio.
                        Tocantins: Plantio de janeiro (1ª quinzena) até março (2ª quinzena), com pico intenso de janeiro a março. Colheita de maio (2ª quinzena) até agosto (1ª quinzena), com pico intenso em julho.
                        Maranhão: Plantio de janeiro (1ª quinzena) até março (2ª quinzena), com pico intenso em janeiro e fevereiro. Colheita de maio (1ª quinzena) até agosto (2ª quinzena), com pico intenso em junho e julho.
                        Piauí: Plantio de janeiro (1ª quinzena) até março (1ª quinzena), com pico intenso em janeiro e fevereiro. Colheita de maio (1ª quinzena) até agosto (1ª quinzena), com pico intenso em junho e julho.
                        Ceará: Plantio de janeiro (1ª quinzena) até março (2ª quinzena), com pico intenso em janeiro e fevereiro. Colheita de maio (2ª quinzena) até agosto (1ª quinzena), com pico intenso em julho.
                        Rio Grande do Norte: Plantio de fevereiro (1ª quinzena) até março (2ª quinzena), com pico intenso em fevereiro. Colheita de julho (1ª quinzena) até setembro (2ª quinzena), com pico intenso em agosto.
                        Paraíba: Plantio de março (1ª quinzena) até abril (2ª quinzena), com pico intenso em março e abril. Colheita de julho (1ª quinzena) até setembro (2ª quinzena), com pico intenso em agosto.
                        Pernambuco: Plantio de março (1ª quinzena) até abril (2ª quinzena), com pico intenso em março e abril. Colheita de julho (1ª quinzena) até outubro (2ª quinzena), com pico intenso em agosto e setembro.
                        Alagoas: Plantio de abril (2ª quinzena) até junho (2ª quinzena), com pico intenso em maio. Colheita de setembro (1ª quinzena) até dezembro (1ª quinzena), com pico intenso em outubro.
                        Sergipe: Plantio de abril (2ª quinzena) até junho (1ª quinzena), com pico intenso em maio. Colheita de setembro (1ª quinzena) até dezembro (1ª quinzena), com pico intenso em outubro e novembro.
                        Bahia: Plantio de abril (2ª quinzena) até junho (1ª quinzena), com pico intenso em maio. Colheita de agosto (2ª quinzena) até novembro (2ª quinzena), com pico intenso em outubro.
                        Mato Grosso: Plantio de janeiro (2ª quinzena) até março (1ª quinzena), com pico intenso em fevereiro. Colheita de maio (2ª quinzena) até agosto (1ª quinzena), com pico intenso em junho e julho.
                        Mato Grosso do Sul: Plantio de janeiro (1ª quinzena) até março (2ª quinzena), com pico intenso em fevereiro. Colheita de maio (2ª quinzena) até setembro (1ª quinzena), com pico intenso em julho.
                        Goiás: Plantio de janeiro (1ª quinzena) até março (1ª quinzena), com pico intenso em fevereiro. Colheita de maio (1ª quinzena) até setembro (1ª quinzena), com pico intenso em junho e julho.
                        Distrito Federal: Plantio de janeiro (1ª quinzena) até fevereiro (2ª quinzena), com pico intenso em janeiro e fevereiro. Colheita de maio (2ª quinzena) até agosto (2ª quinzena), com pico intenso em junho e julho.
                        Minas Gerais: Plantio de janeiro (1ª quinzena) até março (2ª quinzena), com pico intenso em fevereiro. Colheita de maio (2ª quinzena) até setembro (1ª quinzena), com pico intenso em julho.
                        Espírito Santo: Plantio de fevereiro (1ª quinzena) até março (1ª quinzena), com pico intenso em fevereiro. Colheita de junho (1ª quinzena) até agosto (2ª quinzena), com pico intenso em julho.
                        Rio de Janeiro: Plantio de fevereiro (1ª quinzena) até março (1ª quinzena), com pico intenso em fevereiro. Colheita de junho (1ª quinzena) até agosto (1ª quinzena), com pico intenso em julho.
                        São Paulo: Plantio de janeiro (2ª quinzena) até março (2ª quinzena), com pico intenso em fevereiro e março. Colheita de junho (1ª quinzena) até setembro (2ª quinzena), com pico intenso em julho e agosto.
                        Paraná: Plantio de janeiro (2ª quinzena) até abril (1ª quinzena), com pico intenso em março. Colheita de junho (1ª quinzena) até outubro (1ª quinzena), com pico intenso em agosto e setembro.
                        Santa Catarina: Plantio de janeiro (1ª quinzena) até fevereiro (1ª quinzena), com pico intenso em janeiro. Colheita de maio (1ª quinzena) até junho (2ª quinzena), com pico intenso em maio e junho.

"""
                        info_trigo_cana = """
                        Calendário de Safra: Trigo (Ciclo 120-135 dias)
                        Mato Grosso do Sul: Plantio de março (2ª quinzena) até maio (2ª quinzena), com pico intenso em abril. Colheita de agosto (1ª quinzena) até setembro (2ª quinzena), com pico intenso em agosto.
                        Goiás: Plantio de abril (1ª quinzena) até maio (2ª quinzena), com pico intenso em maio. Colheita de agosto (1ª quinzena) até outubro (1ª quinzena), com pico intenso em setembro.
                        Distrito Federal: Plantio de abril (1ª quinzena) até maio (2ª quinzena), com pico intenso em maio. Colheita de agosto (1ª quinzena) até outubro (1ª quinzena), com pico intenso em setembro.
                        Minas Gerais: Plantio de fevereiro (2ª quinzena) até maio (2ª quinzena), com pico intenso em março e abril. Colheita de julho (1ª quinzena) até setembro (1ª quinzena), com pico intenso em julho e agosto.
                        São Paulo: Plantio de março (2ª quinzena) até junho (1ª quinzena), com pico intenso em abril e maio. Colheita de julho (2ª quinzena) até outubro (2ª quinzena), com pico intenso em agosto e setembro.
                        Paraná: Plantio de abril (1ª quinzena) até julho (1ª quinzena), com pico intenso em maio e junho. Colheita de agosto (2ª quinzena) até novembro (2ª quinzena), com pico intenso em setembro e outubro.
                        Santa Catarina: Plantio de maio (2ª quinzena) até agosto (2ª quinzena), com pico intenso em junho e julho. Colheita de outubro (2ª quinzena) até dezembro (2ª quinzena), com pico intenso em novembro e dezembro.
                        Rio Grande do Sul: Plantio de maio (1ª quinzena) até agosto (1ª quinzena), com pico intenso em junho e julho. Colheita de outubro (1ª quinzena) até dezembro (2ª quinzena), com pico intenso em novembro e dezembro.
                        Calendário de Safra: Cana-de-Açúcar
                        (Diferente dos grãos, a cana possui ciclos de colheita e plantio mais extensos e contínuos em várias regiões)
                        Centro-Oeste: Plantio de janeiro (1ª quinzena) até julho (1ª quinzena) e de outubro (1ª quinzena) até dezembro (2ª quinzena). Colheita de abril (1ª quinzena) até novembro (2ª quinzena).
                        Nordeste: Plantio de janeiro (1ª quinzena) até abril (2ª quinzena) e de setembro (1ª quinzena) até dezembro (2ª quinzena). Colheita de janeiro (2ª quinzena) até maio (1ª quinzena) e de agosto (2ª quinzena) até outubro (2ª quinzena).
                        Norte: Plantio de outubro (1ª quinzena) até dezembro (2ª quinzena). Colheita de maio (1ª quinzena) até outubro (2ª quinzena).
                        Sudeste: Plantio de janeiro (1ª quinzena) até julho (1ª quinzena) e de outubro (1ª quinzena) até dezembro (2ª quinzena). Colheita de abril (1ª quinzena) até novembro (2ª quinzena).
                        Sul: Plantio de janeiro (1ª quinzena) até julho (1ª quinzena) e de outubro (1ª quinzena) até dezembro (2ª quinzena). Colheita de abril (1ª quinzena) até novembro (2ª quinzena).
"""

                        conhecimento_safras = f"""
                        ### BEGIN DADOS_SAFRA ###
                        {info_algodao}
                        {info_arroz}
                        {info_soja}
                        {info_milho}
                        {info_trigo_cana}
                        ### END DADOS_SAFRA ###
                        """

                        prompt_calendario = f'''
                        {contexto_agente}

                        {conhecimento_safras}

                        GERAR CALENDÁRIO COM ESTAS REGRAS:

                        PERÍODO: {data_inicio.strftime('%d/%m/%Y')} a {data_fim.strftime('%d/%m/%Y')}
                        MÊS: {mes_ano}
                        
                        {info_especifica}
                        
                        CONTEXTO: {contexto_mensal}
                        
                        PRODUTOS E TEMAS:
                        {chr(10).join([f"- {', '.join(p['produtos'])} - {', '.join(p['culturas'])} - {p['tema']}" for p in produtos_com_direcionais])}
                        
                        REGRAS CRÍTICAS:
                        1. Semana {semana_feira_inicio.strftime('%d/%m')} a {semana_feira_fim.strftime('%d/%m')}: APENAS 1 PAUTA POR DIA
                        2. Priorizar produtos: {produtos_prioritarios_feira} na semana da feira
                        3. Inserir "{pauta_recorrente_texto}" em TODAS as {', '.join(pauta_recorrente_dias)}
                        4. NÃO repetir temas (máximo {max_repeticoes_tema} repetições)
                        5. Células podem ter múltiplas culturas: "Soja e Milho", "Verdavis e Megafol"
                        6. Praticamente todos os dias com conteúdo
                        7. NUNCA 3 dias consecutivos sem pautas
                        8. Baseie pautas no contexto do mês
                        9. As pautas devem respeitar COM RIGIDEZ as fases reais de cada cultura por estados descritos no bloco 'DADOS_SAFRA'
                        
                        FORMATO:
                        - Célula: "[EMOJI] Produto(s) - Cultura(s) - Tema - Breve descrição"
                        - Ex: "🔵 Verdavis e Megafol - Soja e Milho - Tecnologia feira - Soluções apresentadas na feira"
                        - Ex: "🟢 Victrato pelo Brasil - Soja e Cana - Ação nacional - Resultados em diferentes regiões"
                        
                        Retorne CSV pronto para Excel.
                        '''
                        
                        resposta = modelo_texto.generate_content(prompt_calendario)
                        calendario_csv = resposta.text
                        
                        calendario_limpo = calendario_csv.strip()
                        if '```csv' in calendario_limpo:
                            calendario_limpo = calendario_limpo.replace('```csv', '').replace('```', '')
                        if '```' in calendario_limpo:
                            calendario_limpo = calendario_limpo.replace('```', '')
                        
                        st.session_state.calendario_gerado = calendario_limpo
                        st.session_state.mes_ano_calendario = mes_ano
                        
                        st.success("Calendário gerado")
                        
                    except Exception as e:
                        st.error(f"Erro: {str(e)}")
        
        if 'calendario_gerado' in st.session_state:
            st.subheader(f"Calendário - {st.session_state.mes_ano_calendario}")
            
            tab_csv, tab_xlsx = st.tabs(["CSV", "XLSX"])
            
            with tab_csv:
                st.text_area("CSV:", st.session_state.calendario_gerado, height=400)
                
                st.download_button(
                    "Baixar CSV",
                    data=st.session_state.calendario_gerado,
                    file_name=f"calendario_{mes_ano.replace(' ', '_').lower()}.csv",
                    mime="text/csv"
                )
            
            with tab_xlsx:
                try:
                    import openpyxl
                    from openpyxl.styles import Font, Alignment, Border, Side
                    from io import BytesIO
                    
                    def gerar_xlsx():
                        wb = openpyxl.Workbook()
                        ws = wb.active
                        ws.title = f"Calendário {mes_ano}"
                        
                        ws.merge_cells('A1:G1')
                        ws['A1'] = f"CALENDÁRIO - {mes_ano}"
                        ws['A1'].font = Font(bold=True, size=14)
                        ws['A1'].alignment = Alignment(horizontal='center')
                        
                        dias_semana = ["DOMINGO", "SEGUNDA", "TERÇA", "QUARTA", "QUINTA", "SEXTA", "SÁBADO"]
                        for col, dia in enumerate(dias_semana, 1):
                            cell = ws.cell(row=3, column=col)
                            cell.value = dia
                            cell.font = Font(bold=True)
                            cell.alignment = Alignment(horizontal='center')
                        
                        linhas = st.session_state.calendario_gerado.split('\n')
                        linha_atual = 4
                        
                        for linha in linhas:
                            if linha.strip() and not linha.startswith(',,'):
                                celulas = linha.split(',')
                                for col, conteudo in enumerate(celulas, 1):
                                    if conteudo.strip():
                                        cell = ws.cell(row=linha_atual, column=col)
                                        cell.value = conteudo.strip()
                                        cell.alignment = Alignment(wrap_text=True, vertical='top')
                                        cell.border = Border(
                                            left=Side(style='thin'),
                                            right=Side(style='thin'),
                                            top=Side(style='thin'),
                                            bottom=Side(style='thin')
                                        )
                                linha_atual += 1
                        
                        for col in range(1, 8):
                            ws.column_dimensions[openpyxl.utils.get_column_letter(col)].width = 30
                            for row in range(4, linha_atual):
                                ws.row_dimensions[row].height = 60
                        
                        buffer = BytesIO()
                        wb.save(buffer)
                        buffer.seek(0)
                        return buffer
                    
                    if st.button("Gerar XLSX"):
                        buffer_xlsx = gerar_xlsx()
                        
                        st.download_button(
                            "Baixar XLSX",
                            data=buffer_xlsx.getvalue(),
                            file_name=f"calendario_{mes_ano.replace(' ', '_').lower()}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        )
                    
                except ImportError:
                    st.write("Para XLSX: pip install openpyxl")
                    st.code("pip install openpyxl")
                except Exception as e:
                    st.error(f"Erro XLSX: {str(e)}")

# ========== ABA: GERADOR DE BRIEFINGS ==========
with tab_briefings:
    st.header("📋 Gerador de Briefings")
    
    # Verificar se há agente selecionado
    if not st.session_state.agente_selecionado:
        st.warning("⚠️ Selecione um agente na parte superior do app para usar esta funcionalidade.")
    else:
        agente = st.session_state.agente_selecionado
        st.success(f"🎯 Gerando briefings com base no agente: **{agente['nome']}**")
        
        # Inicializar session_state para briefings
        if 'briefings_gerados' not in st.session_state:
            st.session_state.briefings_gerados = []
        if 'briefing_atual_selecionado' not in st.session_state:
            st.session_state.briefing_atual_selecionado = None
        if 'briefing_em_edicao' not in st.session_state:
            st.session_state.briefing_em_edicao = None
        
        # ABAS PRINCIPAIS: Upload Calendário vs Texto Único
        modo_entrada = st.radio(
            "Escolha o modo de entrada:",
            ["📅 Upload de Calendário (múltiplos briefings)", "📝 Texto Único (briefing individual)"],
            horizontal=True
        )
        
        # --- MODO 1: UPLOAD DE CALENDÁRIO ---
        if modo_entrada == "📅 Upload de Calendário (múltiplos briefings)":
            st.subheader("📅 Gerar Múltiplos Briefings a partir do Calendário")
            
            col_upload1, col_upload2 = st.columns([2, 1])
            
            with col_upload1:
                usar_calendario_existente = st.checkbox("Usar calendário gerado anteriormente", 
                                                      value='calendario_gerado' in st.session_state)
                
                if not usar_calendario_existente or 'calendario_gerado' not in st.session_state:
                    arquivo_calendario = st.file_uploader("📅 Upload do calendário CSV:", type=['csv'])
                else:
                    st.info("✅ Usando calendário gerado anteriormente")
                    arquivo_calendario = None
            
            with col_upload2:
                mes_referencia = st.text_input("Mês de referência:", "JANEIRO 2026")
                ano_referencia = st.text_input("Ano de referência:", "2026")
            
            # Contexto adicional para os briefings
            contexto_briefings = st.text_area(
                "Informações contextuais para orientar a criação dos briefings:",
                placeholder="Exemplo: Foco em campanha de posicionamento de produtos, linguagem técnica mas acessível...",
                height=80
            )
            
            # Botão para processar calendário
            if st.button("🔄 Processar Calendário e Gerar Briefings", type="primary", use_container_width=True):
                # Obter o conteúdo do CSV
                conteudo_csv = ""
                
                if usar_calendario_existente and 'calendario_gerado' in st.session_state:
                    conteudo_csv = st.session_state.calendario_gerado
                    st.success("✅ Usando calendário da sessão")
                elif arquivo_calendario is not None:
                    try:
                        file_bytes = arquivo_calendario.getvalue()
                        
                        # Tentar diferentes encodings
                        try:
                            conteudo_csv = file_bytes.decode('utf-8')
                        except UnicodeDecodeError:
                            try:
                                conteudo_csv = file_bytes.decode('latin-1')
                            except UnicodeDecodeError:
                                conteudo_csv = file_bytes.decode('utf-8', errors='ignore')
                        
                        st.success("✅ Arquivo CSV carregado")
                    except Exception as e:
                        st.error(f"❌ Erro ao ler arquivo: {str(e)}")
                        st.stop()
                else:
                    st.error("❌ Nenhum calendário disponível para processar")
                    st.stop()
                
                # Processar o CSV para extrair TODAS as células de conteúdo
                with st.spinner("📋 Processando calendário e extraindo pautas..."):
                    try:
                        linhas = conteudo_csv.split('\n')
                        todas_pautas = []
                        
                        for linha_num, linha in enumerate(linhas):
                            linha_limpa = linha.strip().replace('\r', '').replace('﻿', '')
                            if not linha_limpa:
                                continue
                                
                            celulas = linha_limpa.split(',')
                            for celula_num, celula in enumerate(celulas):
                                celula_limpa = celula.strip()
                                
                                if (celula_limpa and 
                                    len(celula_limpa) > 15 and 
                                    not celula_limpa.replace('.', '').isdigit() and
                                    not any(header in celula_limpa for header in ['DOMINGO', 'SEGUNDA', 'TERÇA', 'QUARTA', 'QUINTA', 'SEXTA', 'SÁBADO', 'CALENDÁRIO']) and
                                    'CX,' not in celula_limpa):
                                    
                                    pautas_na_celula = []
                                    
                                    if '\n' in celula_limpa:
                                        sub_pautas = celula_limpa.split('\n')
                                        for sub_pauta in sub_pautas:
                                            sub_pauta_limpa = sub_pauta.strip()
                                            if sub_pauta_limpa and len(sub_pauta_limpa) > 15:
                                                pautas_na_celula.append(sub_pauta_limpa)
                                    else:
                                        pautas_na_celula.append(celula_limpa)
                                    
                                    for pauta in pautas_na_celula:
                                        pauta_limpa = pauta.strip()
                                        pauta_limpa = ' '.join(pauta_limpa.split())
                                        
                                        todas_pautas.append({
                                            'conteudo': pauta_limpa,
                                            'linha': linha_num,
                                            'coluna': celula_num,
                                            'indice': len(todas_pautas) + 1
                                        })
                        
                        st.success(f"✅ Encontradas {len(todas_pautas)} pautas individuais no calendário")
                        
                        if not todas_pautas:
                            st.error("❌ Nenhuma pauta válida encontrada no CSV")
                            st.stop()
                        
                        # Mostrar preview das pautas encontradas
                        with st.expander("👀 Visualizar Pautas Detectadas", expanded=True):
                            st.write(f"**Total de pautas detectadas:** {len(todas_pautas)}")
                            st.write("**Primeiras 10 pautas:**")
                            for i, pauta in enumerate(todas_pautas[:10]):
                                st.write(f"{i+1}. {pauta['conteudo']}")
                        
                        # Gerar briefings
                        st.subheader("📄 Gerando Briefings para Cada Pauta")
                        
                        contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                        
                        briefings_gerados = []
                        
                        progress_bar = st.progress(0)
                        status_text = st.empty()
                        
                        for idx, pauta in enumerate(todas_pautas):
                            status_text.text(f"Fazendo briefing da pauta {idx+1}/{len(todas_pautas)}: {pauta['conteudo'][:50]}...")
                            progress_bar.progress((idx + 1) / len(todas_pautas))
                            
                            try:
                                prompt_briefing = f"""
                                {contexto_agente}

                                ## TAREFA: GERAR BRIEFING COMPLETO PARA ESTA PAUTA ESPECÍFICA

                                **PAUTA ESPECÍFICA:**
                                {pauta['conteudo']}

                                **MÊS DE REFERÊNCIA:** {mes_referencia}

                                **CONTEXTO ADICIONAL:**
                                {contexto_briefings if contexto_briefings else "Nenhum contexto adicional fornecido."}

                                Gere um briefing completo baseado APENAS nesta pauta específica.
                                Use a base de conhecimento fornecida para identificar produtos, culturas e informações técnicas.
                                Formato completo com contexto, objetivos e formatos. Traga informações chave dos produtos exatamente como são, sem alterar o texto. Mas 
                                posicione, crie um tema, discorra sobre o produto, agregue o tema, após trazer as informações brutas dos produtos que não deve ser alterada,
                                o posicione em termos de benefícios, como que ele deve ser discorrido.

                                # [TÍTULO DO BRIEFING]

                                ## 1. OBJETIVO DO CONTEÚDO
                                [Descreva claramente o objetivo principal deste conteúdo]

                                ## 2. PÚBLICO-ALVO
                                [Descreva a persona, nível técnico, perfil do produtor/leitor]

                                ## 3. TEMA PRINCIPAL E ABORDAGEM
                                [Detalhe o tema central e a abordagem sugerida]

                                ## 4. PRODUTOS ENVOLVIDOS
                                [Liste os produtos e seus papéis no conteúdo]

                                ## 5. CULTURAS ALVO
                                [Especifique as culturas agrícolas relevantes]

                                ## 6. PONTOS-CHAVE OBRIGATÓRIOS
                                - [Ponto 1]
                                - [Ponto 2]
                                - [Ponto 3]
                                [Continue conforme necessário]

                                ## 7. TOM DE VOZ E ESTILO
                                [Especifique o tom: técnico, educativo, comercial, etc.]

                                ## 8. FORMATOS SUGERIDOS
                                - [Formato 1: ex: Post para Instagram]
                                - [Formato 2: ex: Artigo para blog]
                                - [Formato 3: ex: Roteiro para vídeo]

                                ## 9. PALAVRAS-CHAVE (SEO)
                                - Palavra-chave principal:
                                - Palavras-chave secundárias:

                                ## 10. CALL TO ACTION (CTA) SUGERIDO
                                [Texto sugerido para o CTA]

                                ## 11. INFORMAÇÕES TÉCNICAS RELEVANTES
                                [Dados técnicos, estatísticas, informações de manejo que devem ser incluídas]

                                ## 12. RESTRIÇÕES E CUIDADOS
                                [O que evitar, termos proibidos, cuidados especiais]

                                ## 13. REFERÊNCIAS SUGERIDAS
                                [Fontes, materiais de apoio, links úteis]

                                Seja detalhado e específico. O briefing deve servir como um guia completo para a criação do conteúdo.

                                Quando trouxer informações de produtos, os traga exatamente como são sem reescrita.
                                """

                                resposta = modelo_texto.generate_content(prompt_briefing)
                                briefing_gerado = resposta.text
                                
                                briefing_limpo = briefing_gerado.strip()
                                if '```' in briefing_limpo:
                                    briefing_limpo = briefing_limpo.replace('```', '')
                                
                                briefings_gerados.append({
                                    'indice': idx + 1,
                                    'conteudo_original': pauta['conteudo'],
                                    'briefing': briefing_limpo,
                                    'mes_referencia': mes_referencia
                                })
                                
                            except Exception as e:
                                st.error(f"❌ Erro ao gerar briefing para pauta {idx+1}: {str(e)}")
                                briefings_gerados.append({
                                    'indice': idx + 1,
                                    'conteudo_original': pauta['conteudo'],
                                    'briefing': f"ERRO: Não foi possível gerar o briefing.\n{str(e)}",
                                    'mes_referencia': mes_referencia
                                })
                        
                        progress_bar.empty()
                        status_text.empty()
                        
                        st.session_state.briefings_gerados = briefings_gerados
                        st.success(f"✅ {len(briefings_gerados)} briefings gerados com sucesso!")
                        
                    except Exception as e:
                        st.error(f"❌ Erro ao processar calendário: {str(e)}")
        
        # --- MODO 2: TEXTO ÚNICO PARA BRIEFING INDIVIDUAL ---
        else:  # modo_entrada == "📝 Texto Único (briefing individual)"
            st.subheader("📝 Gerar Briefing Individual a partir de Texto")
            
            # Campos para briefing individual
            col_texto1, col_texto2 = st.columns([2, 1])
            
            with col_texto1:
                titulo_briefing = st.text_input(
                    "Título do briefing:",
                    placeholder="Ex: Lançamento do produto X na cultura Y",
                    key="titulo_briefing_individual"
                )
            
            with col_texto2:
                mes_referencia_individual = st.text_input(
                    "Mês de referência:", 
                    "JANEIRO 2026",
                    key="mes_ref_individual"
                )
            
            # Texto base para o briefing
            texto_base_briefing = st.text_area(
                "Texto base para gerar o briefing:",
                height=150,
                placeholder="Cole aqui o texto que servirá de base para o briefing. Pode ser uma pauta, um resumo, instruções do cliente, etc.",
                key="texto_base_individual"
            )
            
            # Contexto adicional
            contexto_individual = st.text_area(
                "Contexto adicional (opcional):",
                height=80,
                placeholder="Informações complementares para orientar a criação do briefing...",
                key="contexto_individual"
            )
            
            # Botão para gerar briefing individual
            col_btn_ind1, col_btn_ind2, col_btn_ind3 = st.columns([1, 2, 1])
            with col_btn_ind2:
                if st.button("📄 GERAR BRIEFING INDIVIDUAL", type="primary", use_container_width=True):
                    if not texto_base_briefing:
                        st.error("❌ O texto base é obrigatório!")
                    elif not titulo_briefing:
                        st.error("❌ O título do briefing é obrigatório!")
                    else:
                        with st.spinner("🔄 Gerando briefing individual..."):
                            try:
                                contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                                
                                prompt_briefing_individual = f"""
                                {contexto_agente}

                                ## TAREFA: GERAR BRIEFING COMPLETO E ESTRUTURADO

                                **TÍTULO DO BRIEFING:** {titulo_briefing}
                                **MÊS DE REFERÊNCIA:** {mes_referencia_individual}

                                **TEXTO BASE:**
                                {texto_base_briefing}

                                **CONTEXTO ADICIONAL:**
                                {contexto_individual if contexto_individual else "Nenhum contexto adicional fornecido."}

                                ## INSTRUÇÕES PARA O FORMATO DO BRIEFING:
                                Traga informações chave dos produtos exatamente como são, sem alterar o texto. Mas 
                                posicione, crie um tema, discorra sobre o produto, agregue o tema, após trazer as informações brutas dos produtos que não deve ser alterada,
                                o posicione em termos de benefícios, como que ele deve ser discorrido.

                                Gere um briefing completo seguindo EXATAMENTE esta estrutura:

                                # [TÍTULO DO BRIEFING]

                                ## 1. OBJETIVO DO CONTEÚDO
                                [Descreva claramente o objetivo principal deste conteúdo]

                                ## 2. PÚBLICO-ALVO
                                [Descreva a persona, nível técnico, perfil do produtor/leitor]

                                ## 3. TEMA PRINCIPAL E ABORDAGEM
                                [Detalhe o tema central e a abordagem sugerida]

                                ## 4. PRODUTOS ENVOLVIDOS
                                [Liste os produtos e seus papéis no conteúdo]

                                ## 5. CULTURAS ALVO
                                [Especifique as culturas agrícolas relevantes]

                                ## 6. PONTOS-CHAVE OBRIGATÓRIOS
                                - [Ponto 1]
                                - [Ponto 2]
                                - [Ponto 3]
                                [Continue conforme necessário]

                                ## 7. TOM DE VOZ E ESTILO
                                [Especifique o tom: técnico, educativo, comercial, etc.]

                                ## 8. FORMATOS SUGERIDOS
                                - [Formato 1: ex: Post para Instagram]
                                - [Formato 2: ex: Artigo para blog]
                                - [Formato 3: ex: Roteiro para vídeo]

                                ## 9. PALAVRAS-CHAVE (SEO)
                                - Palavra-chave principal:
                                - Palavras-chave secundárias:

                                ## 10. CALL TO ACTION (CTA) SUGERIDO
                                [Texto sugerido para o CTA]

                                ## 11. INFORMAÇÕES TÉCNICAS RELEVANTES
                                [Dados técnicos, estatísticas, informações de manejo que devem ser incluídas]

                                ## 12. RESTRIÇÕES E CUIDADOS
                                [O que evitar, termos proibidos, cuidados especiais]

                                ## 13. REFERÊNCIAS SUGERIDAS
                                [Fontes, materiais de apoio, links úteis]

                                Seja detalhado e específico. O briefing deve servir como um guia completo para a criação do conteúdo.

                                Quando trouxer informações de produtos, os traga exatamente como são sem reescrita.
                                """

                                resposta = modelo_texto.generate_content(prompt_briefing_individual)
                                briefing_gerado = resposta.text
                                
                                briefing_limpo = briefing_gerado.strip()
                                if '```' in briefing_limpo:
                                    briefing_limpo = briefing_limpo.replace('```', '')
                                
                                # Adicionar aos briefings gerados
                                novo_briefing = {
                                    'indice': len(st.session_state.briefings_gerados) + 1,
                                    'titulo': titulo_briefing,
                                    'conteudo_original': texto_base_briefing,
                                    'briefing': briefing_limpo,
                                    'mes_referencia': mes_referencia_individual,
                                    'tipo': 'individual'
                                }
                                
                                st.session_state.briefings_gerados.append(novo_briefing)
                                st.session_state.briefing_atual_selecionado = novo_briefing
                                
                                st.success(f"✅ Briefing '{titulo_briefing}' gerado com sucesso!")
                                
                            except Exception as e:
                                st.error(f"❌ Erro ao gerar briefing: {str(e)}")
        
        # --- SEÇÃO DE VISUALIZAÇÃO E AJUSTE DOS BRIEFINGS (comum aos dois modos) ---
        if st.session_state.briefings_gerados:
            st.markdown("---")
            st.header("📋 Briefings Gerados")
            
            briefings = st.session_state.briefings_gerados
            
            # Seletor de briefing para visualizar/editar
            briefing_options = {}
            for b in briefings:
                if 'titulo' in b:
                    # Briefing individual
                    label = f"{b['indice']}. {b['titulo']} ({b.get('mes_referencia', 'N/A')})"
                else:
                    # Briefing de calendário
                    label = f"{b['indice']}. {b['conteudo_original'][:60]}... ({b.get('mes_referencia', 'N/A')})"
                briefing_options[label] = b
            
            if briefing_options:
                col_sel1, col_sel2 = st.columns([3, 1])
                
                with col_sel1:
                    briefing_selecionado_label = st.selectbox(
                        "Selecione um briefing para visualizar/editar:",
                        list(briefing_options.keys()),
                        key="seletor_briefing_edicao"
                    )
                
                with col_sel2:
                    if st.button("🔄 Carregar Briefing", key="carregar_briefing"):
                        st.session_state.briefing_atual_selecionado = briefing_options[briefing_selecionado_label]
                        st.session_state.briefing_em_edicao = briefing_options[briefing_selecionado_label]['briefing']
                        st.rerun()
            
            # Briefing atual selecionado
            if st.session_state.briefing_atual_selecionado:
                briefing_atual = st.session_state.briefing_atual_selecionado
                
                st.markdown("---")
                st.subheader(f"📄 Briefing {briefing_atual['indice']}")
                
                # Mostrar informações do briefing original
                if 'titulo' in briefing_atual:
                    st.info(f"**Título:** {briefing_atual['titulo']}")
                else:
                    st.info(f"**Pauta original:** {briefing_atual['conteudo_original']}")
                
                st.write(f"**Mês referência:** {briefing_atual.get('mes_referencia', 'N/A')}")
                
                # ============================================
                # SEÇÃO DE AJUSTE PONTUAL DO BRIEFING
                # ============================================
                st.markdown("---")
                st.subheader("✏️ Ajuste Pontual do Briefing")
                st.markdown("**Mantenha a estrutura - altere apenas o solicitado**")
                
                col_ajuste1, col_ajuste2 = st.columns([3, 1])
                
                with col_ajuste1:
                    solicitacao_ajuste_briefing = st.text_area(
                        "Descreva o ajuste desejado:",
                        placeholder="Exemplos:\n- Adicione mais detalhes sobre o público-alvo\n- Inclua informações sobre o produto X na seção de produtos\n- Reforce a necessidade de dados técnicos\n- Simplifique a linguagem na seção de tom de voz\n- Adicione um formato de conteúdo a mais",
                        height=100,
                        key="ajuste_briefing"
                    )
                
                with col_ajuste2:
                    st.markdown("#####")  # Espaçamento
                    if st.button("✅ APLICAR AJUSTE", key="aplicar_ajuste_briefing", use_container_width=True):
                        if solicitacao_ajuste_briefing.strip():
                            with st.spinner("🔄 Aplicando ajuste pontual ao briefing..."):
                                try:
                                    contexto_agente = construir_contexto(agente, st.session_state.segmentos_selecionados)
                                    
                                    prompt_ajuste_briefing = f"""
                                    {contexto_agente}

                                    ## INSTRUÇÕES: AJUSTE PONTUAL DO BRIEFING
                                    ## MANTENHA A ESTRUTURA ORIGINAL - ALTERE APENAS O SOLICITADO

                                    --------------------------------------------------------------------

                                    ### BRIEFING ORIGINAL COMPLETO:
                                    
                                    {briefing_atual['briefing']}

                                    --------------------------------------------------------------------

                                    ### SOLICITAÇÃO ESPECÍFICA DE AJUSTE:
                                    "{solicitacao_ajuste_briefing}"

                                    --------------------------------------------------------------------

                                    ## INFORMAÇÕES DE CONTEXTO:
                                    
                                    **Título/Pauta original:** {briefing_atual.get('titulo', briefing_atual.get('conteudo_original', 'N/A'))}
                                    **Mês de referência:** {briefing_atual.get('mes_referencia', 'N/A')}

                                    --------------------------------------------------------------------

                                    ## REGRAS ABSOLUTAS:

                                    1. **MANTENHA A ESTRUTURA ORIGINAL COMPLETA**
                                       - NÃO remova seções
                                       - NÃO adicione novas seções
                                       - NÃO renomeie títulos das seções
                                       - NÃO altere a ordem do conteúdo

                                    2. **ALTERE APENAS O ESTRITAMENTE SOLICITADO**
                                       - Se o usuário pediu para "adicionar X na seção Y", adicione APENAS isso
                                       - Se o usuário pediu para "corrigir Z", corrija APENAS Z
                                       - TODO o resto do briefing deve permanecer IDÊNTICO

                                    3. **PRESERVE FORMATAÇÃO E ESTILO**
                                       - Mantenha todos os negritos, itálicos, markdown exatamente iguais
                                       - Mantenha a numeração das seções
                                       - Mantenha os bullets points exatamente como estão

                                    --------------------------------------------------------------------

                                    ## SUA TAREFA:

                                    1. IDENTIFIQUE exatamente o que o usuário quer modificar
                                    2. LOCALIZE esse trecho no briefing original
                                    3. APLIQUE a modificação solicitada com PRECISÃO CIRÚRGICA
                                    4. RETORNE O BRIEFING COMPLETO com a alteração feita

                                    **IMPORTANTE:** O briefing retornado deve ser IDÊNTICO ao original, 
                                    EXCETO pela modificação pontual solicitada.

                                    RETORNE APENAS O BRIEFING AJUSTADO, SEM COMENTÁRIOS ADICIONAIS.
                                    """
                                    
                                    resposta_ajuste = modelo_texto.generate_content(prompt_ajuste_briefing)
                                    briefing_ajustado = resposta_ajuste.text
                                    
                                    # Limpar possíveis markdown
                                    if '```' in briefing_ajustado:
                                        briefing_ajustado = briefing_ajustado.replace('```', '')
                                    
                                    # Atualizar o briefing
                                    briefing_atual['briefing'] = briefing_ajustado
                                    briefing_atual['historico_ajustes'] = briefing_atual.get('historico_ajustes', [])
                                    briefing_atual['historico_ajustes'].append({
                                        'data': datetime.datetime.now(),
                                        'solicitacao': solicitacao_ajuste_briefing
                                    })
                                    
                                    st.session_state.briefing_em_edicao = briefing_ajustado
                                    
                                    st.success("✅ Ajuste aplicado com sucesso! Estrutura original preservada.")
                                    st.rerun()
                                    
                                except Exception as e:
                                    st.error(f"❌ Erro ao aplicar ajuste: {str(e)}")
                        else:
                            st.warning("⚠️ Por favor, descreva o ajuste desejado.")
                
                # ============================================
                # VISUALIZAÇÃO DO BRIEFING ATUAL
                # ============================================
                
                # Usar o briefing em edição se existir, senão usar o original
                briefing_para_mostrar = st.session_state.briefing_em_edicao if st.session_state.briefing_em_edicao else briefing_atual['briefing']
                
                # Editor de texto para visualização/edição direta
                briefing_editado = st.text_area(
                    "📝 Conteúdo do Briefing (você pode editar diretamente):",
                    value=briefing_para_mostrar,
                    height=400,
                    key="editor_briefing_direto"
                )
                
                # Botão para salvar edições diretas
                col_save1, col_save2, col_save3 = st.columns([1, 1, 2])
                
                with col_save1:
                    if st.button("💾 Salvar Edições Diretas", type="primary", use_container_width=True):
                        if briefing_editado != briefing_atual['briefing']:
                            briefing_atual['briefing'] = briefing_editado
                            briefing_atual['historico_ajustes'] = briefing_atual.get('historico_ajustes', [])
                            briefing_atual['historico_ajustes'].append({
                                'data': datetime.datetime.now(),
                                'solicitacao': 'Edição direta no editor'
                            })
                            st.session_state.briefing_em_edicao = briefing_editado
                            st.success("✅ Briefing atualizado com sucesso!")
                            st.rerun()
                
                with col_save2:
                    if st.button("🔄 Restaurar Original", use_container_width=True):
                        briefing_atual['briefing'] = briefing_atual.get('briefing_original', briefing_atual['briefing'])
                        st.session_state.briefing_em_edicao = None
                        st.success("✅ Briefing original restaurado!")
                        st.rerun()
                
                # ============================================
                # HISTÓRICO DE AJUSTES
                # ============================================
                if briefing_atual.get('historico_ajustes'):
                    with st.expander("📋 Histórico de Ajustes Realizados"):
                        for i, ajuste in enumerate(briefing_atual['historico_ajustes']):
                            data_ajuste = ajuste.get('data', '')
                            if isinstance(data_ajuste, datetime.datetime):
                                data_str = data_ajuste.strftime('%d/%m/%Y %H:%M:%S')
                            else:
                                data_str = 'Data desconhecida'
                            
                            st.write(f"**{i+1}. {data_str}**")
                            st.write(f"*Solicitação:* {ajuste['solicitacao']}")
                            st.divider()
                
                # ============================================
                # BOTÕES DE DOWNLOAD
                # ============================================
                st.markdown("---")
                col_dl1, col_dl2, col_dl3 = st.columns(3)
                
                with col_dl1:
                    # Download individual
                    nome_arquivo = f"briefing_{briefing_atual['indice']}.txt"
                    if 'titulo' in briefing_atual:
                        nome_arquivo = f"briefing_{briefing_atual['titulo'].replace(' ', '_')}.txt"
                    
                    st.download_button(
                        "💾 Baixar Este Briefing",
                        data=briefing_atual['briefing'],
                        file_name=nome_arquivo,
                        mime="text/plain",
                        use_container_width=True
                    )
                
                with col_dl2:
                    # Download com histórico
                    if briefing_atual.get('historico_ajustes'):
                        briefing_com_historico = f"""# BRIEFING {briefing_atual['indice']}
                        
## INFORMAÇÕES ORIGINAIS
- Título/Pauta: {briefing_atual.get('titulo', briefing_atual.get('conteudo_original', 'N/A'))}
- Mês referência: {briefing_atual.get('mes_referencia', 'N/A')}

## BRIEFING ATUAL
{briefing_atual['briefing']}

## HISTÓRICO DE AJUSTES
"""
                        for i, ajuste in enumerate(briefing_atual['historico_ajustes'], 1):
                            data_ajuste = ajuste.get('data', '')
                            if isinstance(data_ajuste, datetime.datetime):
                                data_str = data_ajuste.strftime('%d/%m/%Y %H:%M:%S')
                            else:
                                data_str = 'Data desconhecida'
                            
                            briefing_com_historico += f"\n{i}. {data_str}\n"
                            briefing_com_historico += f"   Solicitação: {ajuste['solicitacao']}\n"
                        
                        st.download_button(
                            "📋 Baixar com Histórico",
                            data=briefing_com_historico,
                            file_name=f"briefing_{briefing_atual['indice']}_com_historico.txt",
                            mime="text/plain",
                            use_container_width=True
                        )
                
                with col_dl3:
                    # Download em lote (todos os briefings)
                    if len(briefings) > 1:
                        import zipfile
                        import io
                        
                        zip_buffer = io.BytesIO()
                        with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                            for b in briefings:
                                nome_b = f"briefing_{b['indice']}.txt"
                                if 'titulo' in b:
                                    nome_b = f"briefing_{b['titulo'].replace(' ', '_')}.txt"
                                zip_file.writestr(nome_b, b['briefing'])
                            
                            # Arquivo consolidado
                            consolidado = f"TODOS OS BRIEFINGS\n"
                            consolidado += f"Total: {len(briefings)}\n"
                            consolidado += "="*60 + "\n\n"
                            
                            for b in briefings:
                                consolidado += f"BRIEFING {b['indice']}\n"
                                if 'titulo' in b:
                                    consolidado += f"Título: {b['titulo']}\n"
                                else:
                                    consolidado += f"Pauta: {b['conteudo_original']}\n"
                                consolidado += "-"*40 + "\n"
                                consolidado += f"{b['briefing']}\n"
                                consolidado += "="*60 + "\n\n"
                            
                            zip_file.writestr("briefings_consolidados.txt", consolidado)
                        
                        st.download_button(
                            "📦 Baixar Todos (ZIP)",
                            data=zip_buffer.getvalue(),
                            file_name=f"todos_briefings_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.zip",
                            mime="application/zip",
                            use_container_width=True
                        )
with tab_revisao_tecnica2:
    st.header("🔬 Revisão Técnica Completa")
    st.markdown("**Análise rigorosa com expertise técnica em agronomia**")
    
    # Criar duas colunas para visualização lado a lado
    col_original, col_revisado = st.columns(2)
    
    with col_original:
        st.subheader("📄 Conteúdo Original")
        texto_tecnico = st.text_area(
            "Cole o conteúdo técnico agrícola para revisão:", 
            height=300,
            placeholder="Cole aqui qualquer conteúdo agrícola que precisa ser revisado tecnicamente...",
            key="texto_tecnico_original",
            label_visibility="collapsed"  # Esconde o label para usar o subheader
        )

    with col_revisado:
        st.subheader("✨ Conteúdo Revisado")
        # Placeholder para o conteúdo revisado
        revisao_placeholder = st.empty()
        revisao_placeholder.info("📝 Aguardando revisão... O conteúdo revisado aparecerá aqui.")

    # Botão para realizar revisão técnica completa - agora centralizado
    st.markdown("---")
    col_btn1, col_btn2, col_btn3 = st.columns([1, 2, 1])
    
    with col_btn2:
        if st.button("🔬 Realizar Revisão Técnica Completa", type="primary", key="revisao_inicial", use_container_width=True):
            if texto_tecnico:
                with st.spinner("🔍 Analisando conteúdo com rigor técnico..."):
                    try:
                        # Prompt para revisão técnica no formato específico
                        prompt_revisao = f"""
                        VOCÊ É: Um engenheiro agrônomo com ampla experiência técnica.

                        SUA TAREFA: Realizar uma revisão técnica completa do conteúdo fornecido seguindo EXATAMENTE o formato abaixo.

                        ANALISE ESTE CONTEÚDO:
                        {texto_tecnico}

                        RETORNE APENAS ESTE FORMATO EXATO:

                        ✅ O QUE ESTÁ CORRETO NO TEXTO (visão geral)
                        Antes das correções, é importante destacar que o texto está bem escrito, com boa estrutura, e a maior parte das informações está correta:
                        [Liste aqui os pontos que estão corretos em bullet points]
                        Ou seja: o conteúdo é bom, faltando apenas alguns ajustes e correções pontuais.

                        ❗ PONTOS INCORRETOS, IMPRECISOS OU QUE PRECISAM SER AJUSTADOS
                        Abaixo, estão todos os erros e imprecisões técnicas do texto, com explicação e sugestão.

                        ❌ 1. [Título do primeiro erro]
                        No trecho:
                        "[Citação exata do trecho problemático]"
                        Correção técnica:
                        [Explicação detalhada do erro]
                        ➡ Portanto, [conclusão técnica]
                        Como corrigir:
                        "[Sugestão de texto corrigido]"

                        ❌ 2. [Título do segundo erro]
                        No trecho:
                        "[Citação exata do trecho problemático]"
                        Correção técnica:
                        [Explicação detalhada do erro]
                        ➡ Portanto, [conclusão técnica]
                        Como corrigir:
                        "[Sugestão de texto corrigido]"

                        [Continue numerando para cada erro encontrado...]

                        🧪 CONCLUSÃO TÉCNICA
                        O texto está bem escrito e majoritariamente correto, mas contém:
                        ✔ [X] erro(s) crítico(s)
                        [Descrição dos erros críticos]
                        ✔ [Y] afirmações que precisam correção ou moderação
                        [Descrição das correções necessárias]
                        ✔ [Z] pontos que não estão errados, mas precisam maior precisão
                        [Descrição dos pontos que precisam de precisão]
                        ✔ [W] pontos incompletos (não são erros, mas faltam informações-chave)
                        [Descrição dos pontos incompletos]

                        🔧 Se quiser, posso agora:
                        - Reescrever o texto totalmente revisado e técnico, já corrigido
                        - Criar uma versão mais curta para redes sociais
                        - Criar uma versão para material comercial
                        - Montar um quadro comparativo entre técnicas/culturas
                        - Fazer uma versão para cultura específica

                        Seja direto e técnico. Mantenha o formato exato.
                        """

                        resposta = modelo_texto2.generate_content(prompt_revisao)
                        revisao_completa = resposta.text
                        
                        # Salvar no session state para uso posterior
                        st.session_state.ultima_revisao = revisao_completa
                        st.session_state.texto_original_revisao = texto_tecnico
                        
                        # Atualizar a coluna direita com o conteúdo revisado
                        with col_revisado:
                            revisao_placeholder.empty()
                            st.success("✅ Revisão concluída!")
                            
                            # Criar abas para organizar o conteúdo revisado
                            tab_relatorio, tab_texto_corrigido = st.tabs(["📋 Relatório Completo", "📝 Texto Corrigido"])
                            
                            with tab_relatorio:
                                st.markdown(revisao_completa)
                            
                            with tab_texto_corrigido:
                                # Extrair e mostrar apenas as sugestões de texto corrigido
                                st.info("📝 **Texto revisado com correções aplicadas:**")
                                
                                # Extrair todas as sugestões de correção do relatório
                                linhas = revisao_completa.split('\n')
                                texto_corrigido_final = texto_tecnico
                                
                                # Procurar por sugestões de correção no formato "Como corrigir:"
                                for i, linha in enumerate(linhas):
                                    if "Como corrigir:" in linha and i + 1 < len(linhas):
                                        sugestao = linhas[i + 1].strip().strip('"')
                                        if sugestao:
                                            # Encontrar o trecho original que está sendo corrigido
                                            for j in range(i-3, i):
                                                if j >= 0 and "No trecho:" in linhas[j] and j + 1 < len(linhas):
                                                    trecho_original = linhas[j + 1].strip().strip('"')
                                                    if trecho_original:
                                                        # Substituir no texto corrigido
                                                        texto_corrigido_final = texto_corrigido_final.replace(
                                                            trecho_original, sugestao
                                                        )
                                
                                # Se nenhuma substituição foi feita, mostrar o original
                                if texto_corrigido_final == texto_tecnico:
                                    st.warning("⚠️ Não foi possível extrair automaticamente o texto corrigido. Mostrando o relatório completo.")
                                    st.markdown(revisao_completa)
                                else:
                                    st.text_area(
                                        "Texto com correções aplicadas:",
                                        texto_corrigido_final,
                                        height=300,
                                        label_visibility="collapsed"
                                    )
                        
                        # Botões de download na parte inferior
                        st.markdown("---")
                        col_dl1, col_dl2 = st.columns(2)
                        
                        with col_dl1:
                            st.download_button(
                                "💾 Baixar Relatório Completo",
                                data=revisao_completa,
                                file_name=f"revisao_tecnica_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                                mime="text/plain",
                                use_container_width=True
                            )
                        
                        with col_dl2:
                            # Tentar extrair o texto corrigido para download
                            texto_para_download = texto_corrigido_final if 'texto_corrigido_final' in locals() else texto_tecnico
                            st.download_button(
                                "💾 Baixar Texto Corrigido",
                                data=texto_para_download,
                                file_name=f"texto_corrigido_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                                mime="text/plain",
                                use_container_width=True
                            )
                    
                    except Exception as e:
                        st.error(f"❌ Erro na revisão técnica: {str(e)}")
                        with col_revisado:
                            revisao_placeholder.error(f"❌ Erro: {str(e)}")
            else:
                st.warning("Por favor, cole um conteúdo técnico para revisão.")

    # Seção para ajustes incrementais (só aparece após a primeira revisão)
    if 'ultima_revisao' in st.session_state:
        st.markdown("---")
        st.subheader("🔄 Ajustes Incrementais")
        
        st.info("Use o campo abaixo para solicitar ajustes específicos na última revisão realizada.")
        
        # Caixa de texto para comandos de ajuste
        comando_ajuste = st.text_area(
            "Comandos para ajustar a última revisão:",
            height=150,
            placeholder="Exemplos:\n- Foque mais na adubação nitrogenada\n- Adicione informações sobre irrigação\n- Corrija os termos técnicos sobre pragas\n- Simplifique a linguagem para produtores\n- Inclua recomendações para clima tropical",
            key="comando_ajuste"
        )
        
        # Botão para revisar novamente com base nos ajustes
        if st.button("🔄 Revisar Novamente com Ajustes", type="secondary", use_container_width=True):
            if comando_ajuste:
                with st.spinner("🔄 Aplicando ajustes solicitados..."):
                    try:
                        # Prompt para revisão com ajustes
                        prompt_ajuste = f"""
                        VOCÊ É: Um engenheiro agrônomo com ampla experiência técnica.

                        SUA TAREFA: Revisar e ajustar o relatório técnico anterior com base nas solicitações específicas do usuário.

                        RELATÓRIO TÉCNICO ANTERIOR:
                        {st.session_state.ultima_revisao}

                        TEXTO ORIGINAL ANALISADO:
                        {st.session_state.texto_original_revisao}

                        SOLICITAÇÕES DE AJUSTE DO USUÁRIO:
                        {comando_ajuste}

                        INSTRUÇÕES:
                        1. Mantenha o MESMO FORMATO EXATO do relatório anterior
                        2. Aplique TODOS os ajustes solicitados pelo usuário
                        3. Mantenha a qualidade técnica e rigor científico
                        4. Se o ajuste solicitar foco em algum aspecto específico, dê mais ênfase a esse tópico
                        5. Se o ajuste pedir adição de informações, inclua-as de forma coerente
                        6. Se o ajuste for sobre estilo ou linguagem, adapte conforme solicitado

                        RETORNE APENAS O RELATÓRIO REVISADO NO MESMO FORMATO, SEM COMENTÁRIOS ADICIONAIS.
                        """

                        resposta_ajuste = modelo_texto2.generate_content(prompt_ajuste)
                        revisao_ajustada = resposta_ajuste.text
                        
                        # Atualizar o session state com a nova versão
                        st.session_state.ultima_revisao = revisao_ajustada
                        
                        # Atualizar a visualização da coluna direita
                        with col_revisado:
                            revisao_placeholder.empty()
                            st.success("✅ Revisão ajustada concluída!")
                            
                            # Criar abas para organizar o conteúdo revisado
                            tab_relatorio, tab_texto_corrigido = st.tabs(["📋 Relatório Ajustado", "📝 Texto Corrigido"])
                            
                            with tab_relatorio:
                                st.markdown(revisao_ajustada)
                            
                            with tab_texto_corrigido:
                                # Extrair e mostrar apenas as sugestões de texto corrigido
                                st.info("📝 **Texto revisado com correções aplicadas:**")
                                
                                # Extrair todas as sugestões de correção do relatório
                                linhas = revisao_ajustada.split('\n')
                                texto_corrigido_final = st.session_state.texto_original_revisao
                                
                                # Procurar por sugestões de correção no formato "Como corrigir:"
                                for i, linha in enumerate(linhas):
                                    if "Como corrigir:" in linha and i + 1 < len(linhas):
                                        sugestao = linhas[i + 1].strip().strip('"')
                                        if sugestao:
                                            # Encontrar o trecho original que está sendo corrigido
                                            for j in range(i-3, i):
                                                if j >= 0 and "No trecho:" in linhas[j] and j + 1 < len(linhas):
                                                    trecho_original = linhas[j + 1].strip().strip('"')
                                                    if trecho_original:
                                                        # Substituir no texto corrigido
                                                        texto_corrigido_final = texto_corrigido_final.replace(
                                                            trecho_original, sugestao
                                                        )
                                
                                # Se nenhuma substituição foi feita, mostrar o original
                                if texto_corrigido_final == st.session_state.texto_original_revisao:
                                    st.warning("⚠️ Não foi possível extrair automaticamente o texto corrigido. Mostrando o relatório completo.")
                                    st.markdown(revisao_ajustada)
                                else:
                                    st.text_area(
                                        "Texto com correções aplicadas:",
                                        texto_corrigido_final,
                                        height=300,
                                        label_visibility="collapsed"
                                    )
                        
                        # Botões de download atualizados
                        st.markdown("---")
                        col_dl1, col_dl2 = st.columns(2)
                        
                        with col_dl1:
                            st.download_button(
                                "💾 Baixar Relatório Ajustado",
                                data=revisao_ajustada,
                                file_name=f"revisao_ajustada_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                                mime="text/plain",
                                key="download_ajustado",
                                use_container_width=True
                            )
                        
                        with col_dl2:
                            # Tentar extrair o texto corrigido para download
                            texto_para_download = texto_corrigido_final if 'texto_corrigido_final' in locals() else st.session_state.texto_original_revisao
                            st.download_button(
                                "💾 Baixar Texto Corrigido",
                                data=texto_para_download,
                                file_name=f"texto_corrigido_{datetime.datetime.now().strftime('%Y%m%d_%H%M')}.txt",
                                mime="text/plain",
                                use_container_width=True
                            )
                    
                    except Exception as e:
                        st.error(f"❌ Erro ao aplicar ajustes: {str(e)}")
                        with col_revisado:
                            revisao_placeholder.error(f"❌ Erro: {str(e)}")
            else:
                st.warning("Por favor, digite os comandos de ajuste desejados.")

            
# --- Estilização ---
st.markdown("""
<style>
    .stChatMessage {
        padding: 1rem;
        border-radius: 0.5rem;
        margin-bottom: 1rem;
    }
    [data-testid="stChatMessageContent"] {
        font-size: 1rem;
    }
    div[data-testid="stTabs"] {
        margin-top: -30px;
    }
    .segment-indicator {
        background-color: #f0f2f6;
        padding: 0.5rem;
        border-radius: 0.5rem;
        margin: 0.5rem 0;
        border-left: 4px solid #4CAF50;
    }
    /* Estilo para o pipeline */
    .pipeline-step {
        background-color: #f8f9fa;
        border-radius: 10px;
        padding: 20px;
        margin: 10px 0;
        border-left: 5px solid #4CAF50;
    }
    .pipeline-complete {
        border-left-color: #4CAF50;
    }
    .pipeline-current {
        border-left-color: #2196F3;
    }
    .pipeline-pending {
        border-left-color: #ff9800;
    }
</style>
""", unsafe_allow_html=True)
