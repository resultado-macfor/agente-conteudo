"""
Serviço de extração de texto de arquivos.
Suporta PDF, TXT, PPTX, DOCX.
"""
import io


def extrair_texto_arquivo(arquivo) -> str:
    """
    Extrai texto de diferentes formatos de arquivo.

    Args:
        arquivo: Arquivo com atributos .name e método .read()

    Returns:
        Texto extraído ou mensagem de erro
    """
    try:
        extensao = arquivo.name.split('.')[-1].lower()

        if extensao == 'pdf':
            return extrair_texto_pdf(arquivo)
        elif extensao == 'txt':
            return extrair_texto_txt(arquivo)
        elif extensao in ['pptx', 'ppt']:
            return extrair_texto_pptx(arquivo)
        elif extensao in ['docx', 'doc']:
            return extrair_texto_docx(arquivo)
        else:
            return f"Formato {extensao} não suportado para extração de texto."

    except Exception as e:
        return f"Erro ao extrair texto do arquivo {arquivo.name}: {str(e)}"


def extrair_texto_pdf(arquivo) -> str:
    """
    Extrai texto de arquivos PDF.

    Args:
        arquivo: Arquivo PDF

    Returns:
        Texto extraído
    """
    try:
        import PyPDF2
        pdf_reader = PyPDF2.PdfReader(arquivo)
        texto = ""
        for pagina in pdf_reader.pages:
            texto += pagina.extract_text() + "\n"
        return texto
    except Exception as e:
        return f"Erro na leitura do PDF: {str(e)}"


def extrair_texto_txt(arquivo) -> str:
    """
    Extrai texto de arquivos TXT.

    Args:
        arquivo: Arquivo TXT

    Returns:
        Texto extraído
    """
    try:
        return arquivo.read().decode('utf-8')
    except:
        try:
            arquivo.seek(0)
            return arquivo.read().decode('latin-1')
        except Exception as e:
            return f"Erro na leitura do TXT: {str(e)}"


def extrair_texto_pptx(arquivo) -> str:
    """
    Extrai texto de arquivos PowerPoint.

    Args:
        arquivo: Arquivo PPTX

    Returns:
        Texto extraído
    """
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(arquivo.read()))
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
        return texto
    except Exception as e:
        return f"Erro na leitura do PowerPoint: {str(e)}"


def extrair_texto_docx(arquivo) -> str:
    """
    Extrai texto de arquivos Word.

    Args:
        arquivo: Arquivo DOCX

    Returns:
        Texto extraído
    """
    try:
        import docx
        doc = docx.Document(io.BytesIO(arquivo.read()))
        texto = ""
        for para in doc.paragraphs:
            texto += para.text + "\n"
        return texto
    except Exception as e:
        return f"Erro na leitura do Word: {str(e)}"
