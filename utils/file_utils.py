def extrair_texto_arquivo(arquivo):
    try:
        extensao = arquivo.name.split('.')[-1].lower()

        if extensao == 'pdf':
            return _extrair_pdf(arquivo)
        elif extensao == 'txt':
            return _extrair_txt(arquivo)
        elif extensao in ['pptx', 'ppt']:
            return _extrair_pptx(arquivo)
        elif extensao in ['docx', 'doc']:
            return _extrair_docx(arquivo)
        else:
            return f"Formato {extensao} não suportado para extração de texto."

    except Exception as e:
        return f"Erro ao extrair texto do arquivo {arquivo.name}: {str(e)}"


def _extrair_pdf(arquivo):
    try:
        import PyPDF2
        pdf_reader = PyPDF2.PdfReader(arquivo)
        texto = ""
        for pagina in pdf_reader.pages:
            texto += pagina.extract_text() + "\n"
        return texto
    except Exception as e:
        return f"Erro na leitura do PDF: {str(e)}"


def _extrair_txt(arquivo):
    try:
        return arquivo.read().decode('utf-8')
    except Exception:
        try:
            return arquivo.read().decode('latin-1')
        except Exception as e:
            return f"Erro na leitura do TXT: {str(e)}"


def _extrair_pptx(arquivo):
    try:
        from pptx import Presentation
        import io
        prs = Presentation(io.BytesIO(arquivo.read()))
        texto = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto += shape.text + "\n"
        return texto
    except Exception as e:
        return f"Erro na leitura do PowerPoint: {str(e)}"


def _extrair_docx(arquivo):
    try:
        import docx
        import io
        doc = docx.Document(io.BytesIO(arquivo.read()))
        texto = ""
        for para in doc.paragraphs:
            texto += para.text + "\n"
        return texto
    except Exception as e:
        return f"Erro na leitura do Word: {str(e)}"
